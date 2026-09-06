"""
PPT 统一渲染管线测试

测试 generate_pptx_file_enhanced 使用 PptGenerator 引擎的完整流程。
不需要 LLM，使用 mock 大纲数据。
"""
import pytest
import sys
import os
import json
import tempfile
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch
from datetime import datetime
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from app.api.v1.aiGeneratorPptx import (
    _content_slides_for_total,
    _fit_editorial_text,
    _normalize_approved_outline,
    generate_pptx_file_enhanced,
    PPT_TEMPLATES,
)


def test_editorial_text_fits_long_copy_into_fixed_box():
    text, size = _fit_editorial_text("这是一段需要在有限空间内自动适配的长文本。" * 8, 2.8, 0.8, 16)

    assert size < 16
    assert text.endswith("…") or len(text) < 64


class TestTemplateMapping:
    """测试 API 模板配置"""

    def test_all_api_templates_exist(self):
        """所有 API 模板都存在"""
        expected_templates = ["modern", "business", "creative", "minimal", "academic", "tech", "education", "medical", "elegant"]
        for tpl in expected_templates:
            assert tpl in PPT_TEMPLATES, f"Missing template: {tpl}"

    def test_template_has_required_fields(self):
        """每个模板都有必需的字段"""
        required_fields = ["name", "primary_color", "secondary_color", "font_family", "background"]
        for tpl_name, tpl_config in PPT_TEMPLATES.items():
            for field in required_fields:
                assert field in tpl_config, f"Template '{tpl_name}' missing field: {field}"

    def test_template_colors_are_valid_hex(self):
        """模板颜色是有效的十六进制格式"""
        for tpl_name, tpl_config in PPT_TEMPLATES.items():
            for color_field in ["primary_color", "secondary_color", "background"]:
                color = tpl_config[color_field]
                assert color.startswith("#"), f"Template '{tpl_name}' {color_field} should start with #"
                assert len(color) == 7, f"Template '{tpl_name}' {color_field} should be 7 chars"


class TestUnifiedGeneration:
    """测试统一渲染管线的完整生成流程"""

    @pytest.fixture
    def mock_request(self):
        """创建 mock PPTGenerationRequest"""
        req = MagicMock()
        req.template = "modern"
        req.language = "zh-CN"
        req.style = "professional"
        req.api_key_token = None
        return req

    @pytest.fixture
    def sample_outline(self):
        """示例 AI 大纲"""
        return {
            "title": "AI 技术发展趋势",
            "slides": [
                {
                    "slide_type": "cover",
                    "title": "AI 技术发展趋势",
                    "subtitle": "2026 年展望",
                    "content": "",
                },
                {
                    "slide_type": "toc",
                    "title": "目录",
                    "content": "第一章：大模型\n第二章：多模态\n第三章：Agent",
                },
                {
                    "slide_type": "content",
                    "title": "大模型发展",
                    "content": "参数规模持续增长\n推理能力显著提升\n多模态融合加深",
                    "notes": "这是演讲者备注",
                },
                {
                    "slide_type": "bullet",
                    "title": "关键技术突破",
                    "content": "Transformer 架构优化\n长上下文窗口\n思维链推理\n工具使用能力",
                },
                {
                    "slide_type": "image",
                    "title": "应用场景",
                    "content": "智能客服\n代码生成\n内容创作",
                    "image_keywords": ["AI", "technology", "robot"],
                },
                {
                    "slide_type": "end",
                    "title": "感谢观看",
                    "subtitle": "Q&A",
                    "content": "",
                },
            ],
        }

    @pytest.mark.asyncio
    async def test_generate_pptx_file_enhanced(self, mock_request, sample_outline):
        """测试增强版 PPTX 生成"""
        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "test_output.pptx"
            await generate_pptx_file_enhanced(filepath, sample_outline, mock_request)

            # 验证文件生成
            assert filepath.exists(), "PPTX file should be created"
            assert filepath.stat().st_size > 0, "PPTX file should not be empty"

            # 验证 JSON 快照
            json_path = filepath.parent / f"{filepath.stem}_slides.json"
            assert json_path.exists(), "JSON snapshot should be created"
            with open(json_path, "r", encoding="utf-8") as f:
                slides_data = json.load(f)
            assert len(slides_data) == 5, f"Expected 5 content slides, got {len(slides_data)}"
            assert len(Presentation(filepath).slides) == 6

    def test_total_slide_budget_removes_input_cover(self, sample_outline):
        slides = _content_slides_for_total(sample_outline["slides"], 5)

        assert len(slides) == 4
        assert slides[0]["slide_type"] == "toc"

    @pytest.mark.asyncio
    async def test_requested_total_is_final_pptx_slide_count(self):
        outline = {
            "title": "五页演示",
            "slides": [
                {
                    "slide_type": "key_points",
                    "title": f"内容页 {index}",
                    "content": ["论点", "证据", "行动"],
                }
                for index in range(1, 6)
            ],
        }
        req = MagicMock(template="modern", api_key_token=None, slide_count=5)

        with tempfile.TemporaryDirectory() as tmpdir, patch(
            "app.api.v1.aiGeneratorPptx.visual_analyzer.analyze_ppt_content",
            new=AsyncMock(return_value=None),
        ):
            filepath = Path(tmpdir) / "five-slides.pptx"
            await generate_pptx_file_enhanced(filepath, outline, req)
            presentation = Presentation(filepath)

        assert len(presentation.slides) == 5

    @pytest.mark.asyncio
    async def test_design_tokens_are_applied_to_generated_pptx(self):
        outline = {
            "title": "令牌渲染",
            "slides": [{
                "slide_type": "key_points",
                "title": "内容页",
                "content": ["论点", "证据", "行动"],
            }],
        }
        req = MagicMock(template="business_report", api_key_token=None, slide_count=2)

        with tempfile.TemporaryDirectory() as tmpdir, patch(
            "app.api.v1.aiGeneratorPptx.visual_analyzer.analyze_ppt_content",
            new=AsyncMock(return_value=None),
        ):
            filepath = Path(tmpdir) / "token-style.pptx"
            await generate_pptx_file_enhanced(filepath, outline, req)
            presentation = Presentation(filepath)

        solid_colors = {
            str(shape.fill.fore_color.rgb)
            for slide in presentation.slides
            for shape in slide.shapes
            if shape.fill.type == MSO_FILL_TYPE.SOLID
        }
        assert "1F4E79" in solid_colors

    @pytest.mark.asyncio
    async def test_generate_with_different_templates(self, sample_outline):
        """测试不同模板的生成"""
        for template_name in PPT_TEMPLATES.keys():
            req = MagicMock()
            req.template = template_name
            req.language = "zh-CN"

            with tempfile.TemporaryDirectory() as tmpdir:
                filepath = Path(tmpdir) / f"test_{template_name}.pptx"
                await generate_pptx_file_enhanced(filepath, sample_outline, req)
                assert filepath.exists(), f"PPTX should be created for template '{template_name}'"
                assert filepath.stat().st_size > 0, f"PPTX should not be empty for template '{template_name}'"

    @pytest.mark.asyncio
    async def test_generate_with_content_list(self, mock_request):
        """测试 content 为列表格式的生成"""
        outline = {
            "title": "列表内容测试",
            "slides": [
                {
                    "slide_type": "content",
                    "title": "列表页",
                    "content": ["第一项", "第二项", "第三项"],
                },
            ],
        }
        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "list_content.pptx"
            await generate_pptx_file_enhanced(filepath, outline, mock_request)
            assert filepath.exists()

    @pytest.mark.asyncio
    async def test_generate_with_empty_content(self, mock_request):
        """测试空内容的生成"""
        outline = {
            "title": "空内容测试",
            "slides": [
                {
                    "slide_type": "content",
                    "title": "空页",
                    "content": "",
                },
            ],
        }
        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "empty_content.pptx"
            await generate_pptx_file_enhanced(filepath, outline, mock_request)
            assert filepath.exists()

    @pytest.mark.asyncio
    async def test_generate_with_unknown_slide_type(self, mock_request):
        """测试未知幻灯片类型"""
        outline = {
            "title": "未知类型测试",
            "slides": [
                {
                    "slide_type": "unknown_type",
                    "title": "未知类型页",
                    "content": "一些内容",
                },
            ],
        }
        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "unknown_type.pptx"
            await generate_pptx_file_enhanced(filepath, outline, mock_request)
            assert filepath.exists()

    @pytest.mark.asyncio
    async def test_generate_with_missing_fields(self, mock_request):
        """测试缺少字段的幻灯片"""
        outline = {
            "slides": [
                {"slide_type": "content"},  # 缺少 title 和 content
            ],
        }
        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "missing_fields.pptx"
            await generate_pptx_file_enhanced(filepath, outline, mock_request)
            assert filepath.exists()

    @pytest.mark.asyncio
    async def test_generate_with_progress_callback(self, mock_request, sample_outline):
        """测试进度回调"""
        progress_calls = []

        async def mock_progress(progress=None, message=None):
            progress_calls.append({"progress": progress, "message": message})

        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "with_progress.pptx"
            await generate_pptx_file_enhanced(filepath, sample_outline, mock_request, update_progress=mock_progress)
            assert filepath.exists()
            assert len(progress_calls) >= 1, "Progress callback should be called at least once"

    @pytest.mark.asyncio
    async def test_json_snapshot_matches_outline(self, mock_request, sample_outline):
        """测试 JSON 快照与输入大纲一致"""
        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "snapshot_test.pptx"
            await generate_pptx_file_enhanced(filepath, sample_outline, mock_request)

            json_path = filepath.parent / f"{filepath.stem}_slides.json"
            with open(json_path, "r", encoding="utf-8") as f:
                saved_slides = json.load(f)

            # 系统封面由渲染器统一生成，快照只保存内容页。
            original_slides = sample_outline["slides"][1:]
            assert len(saved_slides) == len(original_slides)
            for orig, saved in zip(original_slides, saved_slides):
                assert orig["title"] == saved["title"]
                assert orig["slide_type"] == saved["slide_type"]

    @pytest.mark.asyncio
    async def test_visual_plan_preserves_narrative_role_layouts(self, mock_request):
        """视觉分析成功时仍使用五类叙事角色专属版式。"""
        mock_request.template = "business"
        roles = [
            ("opportunity_map", "机会 01"),
            ("evidence_story", "EVIDENCE"),
            ("strategic_choice", "方案 A"),
            ("execution_roadmap", "进入下一阶段的门槛"),
            ("decision_close", "决策 01"),
        ]
        outline = {
            "title": "叙事版式回归测试",
            "slides": [
                {
                    "slide_type": "content",
                    "title": f"页面 {index}",
                    "content": ["要点一", "要点二", "要点三", "验证指标"],
                    "narrative_role": role,
                }
                for index, (role, _) in enumerate(roles, 1)
            ],
        }
        visual_plan = MagicMock()
        visual_plan.slides = [MagicMock() for _ in roles]

        with tempfile.TemporaryDirectory() as tmpdir, patch(
            "app.api.v1.aiGeneratorPptx.visual_analyzer.analyze_ppt_content",
            new=AsyncMock(return_value=visual_plan),
        ):
            filepath = Path(tmpdir) / "narrative_roles.pptx"
            await generate_pptx_file_enhanced(filepath, outline, mock_request)
            presentation = Presentation(filepath)

        for slide, (_, expected_label) in zip(list(presentation.slides)[1:], roles):
            slide_text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
            assert expected_label in slide_text

    @pytest.mark.asyncio
    async def test_creative_theme_uses_distinct_geometry(self):
        """Creative 主题拥有独立构图，避免退化为颜色换肤。"""
        outline = {
            "title": "主题构图测试",
            "slides": [{
                "slide_type": "comparison",
                "title": "战略取舍",
                "content": ["方案一", "方案二", "推荐方案二"],
                "narrative_role": "strategic_choice",
            }],
        }

        with tempfile.TemporaryDirectory() as tmpdir, patch(
            "app.api.v1.aiGeneratorPptx.visual_analyzer.analyze_ppt_content",
            new=AsyncMock(return_value=None),
        ):
            presentations = []
            for template in ("business", "creative"):
                req = MagicMock(template=template, api_key_token=None)
                filepath = Path(tmpdir) / f"{template}.pptx"
                await generate_pptx_file_enhanced(filepath, outline, req)
                presentations.append(Presentation(filepath))

        geometries = []
        for presentation in presentations:
            geometries.append([
                (shape.shape_type, shape.left, shape.top, shape.width, shape.height)
                for shape in presentation.slides[1].shapes
            ])
        assert geometries[0] != geometries[1]

    @pytest.mark.asyncio
    async def test_priority_themes_use_pairwise_distinct_geometry(self):
        """高频主题使用成对不同的页面骨架。"""
        outline = {
            "title": "高频主题构图测试",
            "slides": [{
                "slide_type": "comparison",
                "title": "优先路径选择",
                "key_message": "四周内验证高价值路径",
                "content": ["路径一", "路径二", "优先路径二", "复盘后扩大投入"],
                "narrative_role": "strategic_choice",
            }],
        }

        with tempfile.TemporaryDirectory() as tmpdir, patch(
            "app.api.v1.aiGeneratorPptx.visual_analyzer.analyze_ppt_content",
            new=AsyncMock(return_value=None),
        ):
            geometries = {}
            for template in ("modern", "minimal", "academic", "education", "medical", "elegant", "tech", "business", "creative"):
                req = MagicMock(template=template, api_key_token=None)
                filepath = Path(tmpdir) / f"{template}.pptx"
                await generate_pptx_file_enhanced(filepath, outline, req)
                presentation = Presentation(filepath)
                geometries[template] = [
                    (shape.shape_type, shape.left, shape.top, shape.width, shape.height)
                    for shape in presentation.slides[1].shapes
                ]

        signatures = {tuple(geometry) for geometry in geometries.values()}
        assert len(signatures) == len(geometries)

    @pytest.mark.asyncio
    @pytest.mark.parametrize(
        ("template", "expected_label"),
        [
            ("modern", "DECISION"),
            ("minimal", "DECISION /"),
            ("academic", "RESEARCH CONCLUSION"),
            ("education", "LEARNING CHECK"),
            ("medical", "CLINICAL RATIONALE"),
            ("elegant", "BOARD RECOMMENDATION"),
            ("tech", "LOCK / RECOMMENDATION"),
        ],
    )
    async def test_priority_theme_renders_semantic_choice(self, template, expected_label):
        """高频主题保留策略选择语义和结论标签。"""
        outline = {
            "title": "语义主题测试",
            "slides": [{
                "slide_type": "comparison",
                "title": "战略选择",
                "key_message": "先验证闭环，再扩大投入",
                "content": ["完整平台", "场景闭环", "推荐场景闭环"],
                "narrative_role": "strategic_choice",
            }],
        }

        with tempfile.TemporaryDirectory() as tmpdir, patch(
            "app.api.v1.aiGeneratorPptx.visual_analyzer.analyze_ppt_content",
            new=AsyncMock(return_value=None),
        ):
            req = MagicMock(template=template, api_key_token=None)
            filepath = Path(tmpdir) / f"{template}.pptx"
            await generate_pptx_file_enhanced(filepath, outline, req)
            presentation = Presentation(filepath)

        slide_text = "\n".join(
            shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text")
        )
        assert expected_label in slide_text
        assert "推荐场景闭环" in slide_text

    def test_approved_outline_preserves_structured_commercial_blocks(self):
        outline = {
            "title": "商业字段",
            "slides": [{
                "title": "机会判断",
                "key_message": "两周内验证",
                "slide_type": "key_points",
                "narrative_role": "opportunity_map",
                "content_blocks": [{
                    "type": "metric",
                    "content": "验证投入产出比",
                    "metadata": {"roi": "≥3.0", "validation_period": "2 周"},
                }],
            }],
        }

        normalized = _normalize_approved_outline(outline)

        assert normalized["slides"][0]["content_blocks"][0]["metadata"]["roi"] == "≥3.0"
        assert normalized["slides"][0]["key_message"] == "两周内验证"

    def test_normalization_repairs_repeated_commercial_roles(self):
        outline = {
            "title": "重复角色修复",
            "slides": [
                {
                    "title": f"页面 {index}",
                    "slide_type": "key_points",
                    "narrative_role": "opportunity_map",
                    "content": ["内容"],
                }
                for index in range(1, 5)
            ],
        }

        normalized = _normalize_approved_outline(outline)

        assert [slide["narrative_role"] for slide in normalized["slides"]] == [
            "opportunity_map",
            "evidence_story",
            "strategic_choice",
            "execution_roadmap",
        ]
        assert [slide["slide_type"] for slide in normalized["slides"]] == [
            "key_points",
            "data",
            "comparison",
            "timeline",
        ]

    @pytest.mark.asyncio
    async def test_renderer_displays_commercial_metadata(self):
        outline = {
            "title": "商业字段渲染",
            "slides": [{
                "slide_type": "comparison",
                "title": "战略选择",
                "narrative_role": "strategic_choice",
                "content_blocks": [
                    {"type": "option", "content": "完整平台", "metadata": {"cost": "高", "timeframe": "12 周", "risk": "扩散"}},
                    {"type": "option", "content": "场景闭环", "metadata": {"cost": "中", "timeframe": "4 周", "risk": "选型"}},
                    {"type": "recommendation", "content": "推荐场景闭环", "metadata": {"rationale": "反馈更快"}},
                ],
            }],
        }
        req = MagicMock(template="business", api_key_token=None)

        with tempfile.TemporaryDirectory() as tmpdir, patch(
            "app.api.v1.aiGeneratorPptx.visual_analyzer.analyze_ppt_content",
            new=AsyncMock(return_value=None),
        ):
            filepath = Path(tmpdir) / "commercial_metadata.pptx"
            await generate_pptx_file_enhanced(filepath, outline, req)
            presentation = Presentation(filepath)

        slide_text = "\n".join(
            shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text")
        )
        assert "成本：高" in slide_text
        assert "周期：12 周" in slide_text
        assert "风险：扩散" in slide_text
        assert "依据：反馈更快" in slide_text


class TestEndToEndGeneration:
    """端到端生成测试（模拟完整 API 调用流程）"""

    @pytest.mark.asyncio
    async def test_full_generation_flow(self):
        """模拟完整生成流程：大纲构建 -> 渲染 -> 文件验证"""
        req = MagicMock()
        req.template = "business"
        req.language = "zh-CN"

        outline = {
            "title": "季度业务报告",
            "subtitle": "2026 Q1",
            "slides": [
                {"slide_type": "cover", "title": "季度业务报告", "subtitle": "2026 Q1"},
                {"slide_type": "toc", "title": "议程", "content": "业绩回顾\n市场分析\n下季度计划"},
                {"slide_type": "content", "title": "业绩回顾", "content": "营收增长 25%\n用户增长 30%\n新客户 15 家"},
                {"slide_type": "chart", "title": "收入趋势", "content": "Q1: 100 万\nQ2: 120 万\nQ3: 150 万"},
                {"slide_type": "bullet", "title": "关键成就", "content": "产品上线\n团队扩张\n市场拓展"},
                {"slide_type": "end", "title": "谢谢", "subtitle": "Questions?"},
            ],
        }

        with tempfile.TemporaryDirectory() as tmpdir:
            filepath = Path(tmpdir) / "quarterly_report.pptx"
            await generate_pptx_file_enhanced(filepath, outline, req)

            # 验证文件
            assert filepath.exists()
            file_size = filepath.stat().st_size
            assert file_size > 10000, f"PPTX file too small: {file_size} bytes"

            # 验证是有效的 PPTX (ZIP 格式)
            import zipfile
            assert zipfile.is_zipfile(filepath), "PPTX should be a valid ZIP file"

            # 验证包含必需的 PPTX 组件
            with zipfile.ZipFile(filepath) as zf:
                names = zf.namelist()
                assert "[Content_Types].xml" in names, "PPTX should contain [Content_Types].xml"
                assert any("ppt/slides/slide" in n for n in names), "PPTX should contain slide files"

            # 验证 JSON 快照
            json_path = filepath.parent / f"{filepath.stem}_slides.json"
            assert json_path.exists()
            with open(json_path, "r", encoding="utf-8") as f:
                saved = json.load(f)
            assert len(saved) == 5


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
