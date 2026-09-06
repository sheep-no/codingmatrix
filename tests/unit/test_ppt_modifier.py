"""
PPT 修改器单元测试

测试 PPTX 文件修改功能
"""
import pytest
import sys
import os
import tempfile
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from app.utils.pptx.ppt_modifier import PPTModifier, modify_ppt, COLOR_MAP
from app.utils.pptx.modify_intent_parser import ModifyIntent, ModifyTarget


def create_test_pptx(path: str, num_slides: int = 3):
    """创建测试用 PPTX 文件"""
    prs = Presentation()

    for i in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 带标题和内容的布局

        # 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = f"标题 {i + 1}"

        # 设置内容
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text_frame.paragraphs[0].text = f"内容 {i + 1}-1"
                p = shape.text_frame.add_paragraph()
                p.text = f"内容 {i + 1}-2"

    prs.save(path)
    return path


class TestPPTModifier:
    """测试 PPT 修改器"""

    def setup_method(self):
        self.tmp_dir = tempfile.mkdtemp()
        self.input_path = os.path.join(self.tmp_dir, "test_input.pptx")
        self.output_path = os.path.join(self.tmp_dir, "test_output.pptx")
        create_test_pptx(self.input_path)

    def teardown_method(self):
        import shutil
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def test_init_with_valid_file(self):
        """测试有效文件初始化"""
        modifier = PPTModifier(self.input_path)
        assert modifier.prs is not None

    def test_init_with_invalid_file(self):
        """测试无效文件初始化"""
        with pytest.raises(FileNotFoundError):
            PPTModifier("/nonexistent/file.pptx")

    def test_apply_font_modification(self):
        """测试字体修改"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(slide_number=1, property_name="font", property_value="微软雅黑")
            ]
        )
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True
        assert os.path.exists(self.output_path)

        # 验证修改后的文件可以打开
        prs = Presentation(self.output_path)
        assert len(prs.slides) > 0

    def test_apply_color_modification(self):
        """测试颜色修改"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(slide_number=1, property_name="color", property_value="红色")
            ]
        )
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True
        assert os.path.exists(self.output_path)

    def test_apply_size_modification(self):
        """测试字号修改"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(slide_number=1, property_name="size", property_value="24")
            ]
        )
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True

    def test_modify_all_slides(self):
        """测试修改所有幻灯片"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(slide_number=None, property_name="font", property_value="Arial")
            ]
        )
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True

    def test_modify_nonexistent_slide(self):
        """测试修改不存在的幻灯片"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(slide_number=999, property_name="font", property_value="Arial")
            ]
        )
        # 不应该抛出异常，只是警告
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True

    def test_modify_with_element_type_title(self):
        """测试按元素类型修改 - 标题"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(
                    slide_number=1,
                    element_type="title",
                    property_name="font",
                    property_value="黑体"
                )
            ]
        )
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True

    def test_modify_with_element_type_text(self):
        """测试按元素类型修改 - 正文"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(
                    slide_number=1,
                    element_type="text",
                    property_name="font",
                    property_value="宋体"
                )
            ]
        )
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True

    def test_empty_intent(self):
        """测试空修改意图"""
        modifier = PPTModifier(self.input_path)
        intent = ModifyIntent(raw_text="test", targets=[])
        success = modifier.apply_modifications(intent, self.output_path)
        assert success is True


class TestModifyPPTFunction:
    """测试 modify_ppt 函数"""

    def setup_method(self):
        self.tmp_dir = tempfile.mkdtemp()
        self.input_path = os.path.join(self.tmp_dir, "test_input.pptx")
        self.output_path = os.path.join(self.tmp_dir, "test_output.pptx")
        create_test_pptx(self.input_path)

    def teardown_method(self):
        import shutil
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def test_basic_modify(self):
        """测试基本修改"""
        intent = ModifyIntent(
            raw_text="test",
            targets=[
                ModifyTarget(slide_number=1, property_name="font", property_value="微软雅黑")
            ]
        )
        success = modify_ppt(self.input_path, intent, self.output_path)
        assert success is True
        assert os.path.exists(self.output_path)


class TestColorMap:
    """测试颜色映射"""

    def test_chinese_colors(self):
        """测试中文颜色"""
        assert "红色" in COLOR_MAP
        assert "蓝色" in COLOR_MAP
        assert "绿色" in COLOR_MAP

    def test_english_colors(self):
        """测试英文颜色"""
        assert "red" in COLOR_MAP
        assert "blue" in COLOR_MAP
        assert "green" in COLOR_MAP

    def test_color_values(self):
        """测试颜色值"""
        assert COLOR_MAP["红色"] == RGBColor(0xFF, 0x00, 0x00)
        assert COLOR_MAP["蓝色"] == RGBColor(0x00, 0x00, 0xFF)
        assert COLOR_MAP["黑色"] == RGBColor(0x00, 0x00, 0x00)


class TestSlideRenderer:
    """测试幻灯片渲染器"""

    def setup_method(self):
        self.tmp_dir = tempfile.mkdtemp()
        self.pptx_path = os.path.join(self.tmp_dir, "test.pptx")
        create_test_pptx(self.pptx_path)

    def teardown_method(self):
        import shutil
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def test_render_slide(self):
        """测试渲染单页幻灯片"""
        from app.utils.pptx.slide_renderer import SlideRenderer

        renderer = SlideRenderer(self.pptx_path)
        img_bytes = renderer.render_slide(1)
        assert img_bytes is not None
        assert len(img_bytes) > 0

    def test_render_all_slides(self):
        """测试渲染所有幻灯片"""
        from app.utils.pptx.slide_renderer import SlideRenderer

        renderer = SlideRenderer(self.pptx_path)
        images = renderer.render_all_slides()
        assert len(images) == 3

    def test_render_nonexistent_slide(self):
        """测试渲染不存在的幻灯片"""
        from app.utils.pptx.slide_renderer import SlideRenderer

        renderer = SlideRenderer(self.pptx_path)
        img_bytes = renderer.render_slide(999)
        assert img_bytes is None

    def test_get_slide_metadata(self):
        """测试获取幻灯片元数据"""
        from app.utils.pptx.slide_renderer import SlideRenderer

        renderer = SlideRenderer(self.pptx_path)
        metadata = renderer.get_slide_metadata(1)
        assert metadata is not None
        assert "fonts" in metadata
        assert "colors" in metadata
        assert "elements" in metadata

    def test_widescreen_coordinates_use_presentation_dimensions(self):
        from app.utils.pptx.slide_renderer import PREVIEW_WIDTH, SlideRenderer

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(12), Inches(1), Inches(1), Inches(1)).text = "右侧文本"
        prs.save(self.pptx_path)

        renderer = SlideRenderer(self.pptx_path)
        preview = renderer._parse_slide(renderer.prs.slides[0], 1)
        text_element = next(element for element in preview.elements if element.text == "右侧文本")

        assert text_element.x + text_element.width <= PREVIEW_WIDTH
        assert text_element.x >= 1100

    def test_text_shape_keeps_background_and_chinese_text(self):
        from app.utils.pptx.slide_renderer import SlideRenderer, _load_preview_font

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(1), Inches(4), Inches(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        shape.text = "策略验证"
        prs.save(self.pptx_path)

        renderer = SlideRenderer(self.pptx_path)
        elements = renderer._parse_shape(renderer.prs.slides[0].shapes[0])

        assert [element.type for element in elements] == ["shape", "text"]
        assert elements[0].fill_color == (0x1F, 0x4E, 0x79)
        assert elements[1].text == "策略验证"
        assert elements[1].is_bullet is False
        assert _load_preview_font(20).getmask("策略验证").getbbox() is not None
