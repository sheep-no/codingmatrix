"""
PPT API 端点单元测试

覆盖 10 个此前无测试的端点：
- POST /api/v1/pptx/generate (同步生成)
- GET  /api/v1/pptx/download/{ppt_id} (下载)
- POST /api/v1/pptx/{task_id}/update (增量更新)
- POST /api/v1/pptx/{task_id}/modify (视觉修改)
- GET  /api/v1/pptx/{task_id}/analyze (分析)
- GET  /api/v1/pptx/templates (模板列表)
- GET  /api/v1/pptx/history (历史记录)
- DELETE /api/v1/pptx/history/{task_id} (删除历史)
- GET  /api/v1/pptx/history/stats (统计)
- POST /api/v1/generate-text (仅大纲)
"""
import json
import os
import sys
import uuid

import pytest
from unittest.mock import patch, AsyncMock, MagicMock
from fastapi.testclient import TestClient

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from app.main import app
from app.utils.security import verify_token


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True)
def _override_auth():
    """覆盖认证，跳过真实鉴权"""
    app.dependency_overrides[verify_token] = lambda: {"sub": "test_user", "role": "user"}
    yield
    app.dependency_overrides.clear()


@pytest.fixture
def client():
    return TestClient(app)


@pytest.fixture
def tmp_output(tmp_path):
    """临时 PPT 输出目录"""
    with patch("app.api.v1.aiGeneratorPptx.PPT_OUTPUT_DIR", tmp_path):
        yield tmp_path


@pytest.fixture
def sample_outline():
    return {
        "title": "AI 发展趋势",
        "slides": [
            {"type": "cover", "title": "AI 发展趋势", "bullets": [], "image_keywords": "AI technology"},
            {"type": "content", "title": "概述", "bullets": ["要点1", "要点2"], "image_keywords": ""},
            {"type": "end", "title": "谢谢", "bullets": [], "image_keywords": ""},
        ],
    }


@pytest.fixture
def mock_outline_obj():
    """模拟 PPTAgent.generate_outline 返回的 PresentationOutline 对象"""
    from app.agent.ppt_agent import PresentationOutline, SlideOutline
    return PresentationOutline(
        title="AI 发展趋势",
        slides=[
            SlideOutline(type="cover", title="AI 发展趋势", bullets=[], image_keywords=["AI", "technology"], notes=""),
            SlideOutline(type="content", title="概述", bullets=["要点1"], image_keywords=[], notes=""),
            SlideOutline(type="end", title="谢谢", bullets=[], image_keywords=[], notes=""),
        ],
    )


def _create_pptx_file(path):
    """创建一个最小可用的 PPTX 文件"""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test"
    prs.save(str(path))


def _write_snapshot(path, slides, user_id="test_user", title="测试"):
    """写入新格式 JSON 快照"""
    with open(path, "w", encoding="utf-8") as f:
        json.dump({"user_id": user_id, "title": title, "slides": slides}, f, ensure_ascii=False)


# ===========================================================================
# 1. POST /api/v1/pptx/generate (同步生成)
# ===========================================================================

class TestGeneratePptSync:
    """同步生成端点测试"""

    @patch("app.api.v1.aiGeneratorPptx.generate_pptx_file_enhanced", new_callable=AsyncMock)
    @patch("app.api.v1.aiGeneratorPptx.generate_ppt_outline", new_callable=AsyncMock)
    def test_generate_pptx_success(self, mock_outline, mock_enhanced, client, tmp_output, sample_outline):
        mock_outline.return_value = sample_outline
        mock_enhanced.return_value = None
        resp = client.post("/api/v1/pptx/generate", json={
            "topic": "AI 发展趋势",
            "slide_count": 5,
            "output_format": "pptx",
        })
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "completed"
        assert "download_url" in data
        assert "preview_url" in data

    @patch("app.api.v1.aiGeneratorPptx.generate_html_ppt", new_callable=AsyncMock)
    @patch("app.api.v1.aiGeneratorPptx.generate_ppt_outline", new_callable=AsyncMock)
    def test_generate_html_success(self, mock_outline, mock_html, client, tmp_output, sample_outline):
        mock_outline.return_value = sample_outline
        mock_html.return_value = None
        resp = client.post("/api/v1/pptx/generate", json={
            "topic": "测试",
            "slide_count": 5,
            "output_format": "html",
        })
        assert resp.status_code == 200
        assert resp.json()["output_format"] == "html"

    @patch("app.api.v1.aiGeneratorPptx.generate_ppt_outline", new_callable=AsyncMock)
    def test_generate_outline_failure(self, mock_outline, client, tmp_output):
        mock_outline.side_effect = RuntimeError("LLM 调用失败")
        resp = client.post("/api/v1/pptx/generate", json={
            "topic": "测试",
            "slide_count": 5,
        })
        assert resp.status_code == 500

    def test_generate_validation_error(self, client):
        """缺少必填字段 topic"""
        resp = client.post("/api/v1/pptx/generate", json={"slide_count": 5})
        assert resp.status_code == 422


# ===========================================================================
# 2. GET /api/v1/pptx/download/{ppt_id}
# ===========================================================================

class TestDownloadPpt:
    """下载端点测试"""

    def test_download_pptx_success(self, client, tmp_output):
        ppt_id = str(uuid.uuid4())
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")
        resp = client.get(f"/api/v1/pptx/download/{ppt_id}")
        assert resp.status_code == 200
        assert resp.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def test_download_not_found(self, client, tmp_output):
        resp = client.get(f"/api/v1/pptx/download/{uuid.uuid4()}")
        assert resp.status_code == 404

    def test_download_html_format(self, client, tmp_output):
        ppt_id = str(uuid.uuid4())
        (tmp_output / f"{ppt_id}.html").write_text("<html>test</html>")
        resp = client.get(f"/api/v1/pptx/download/{ppt_id}?format=html")
        assert resp.status_code == 200
        assert "text/html" in resp.headers["content-type"]

    def test_download_path_traversal_rejected(self, client, tmp_output):
        """路径穿越应被拒绝（ppt_id 包含 ../）"""
        resp = client.get("/api/v1/pptx/download/../../../etc/passwd")
        assert resp.status_code in (404, 422)


# ===========================================================================
# 3. POST /api/v1/pptx/{task_id}/update
# ===========================================================================

class TestUpdatePpt:
    """增量更新端点测试"""

    @patch("app.api.v1.aiGeneratorPptx.task_manager.create_task", new_callable=AsyncMock)
    @patch("app.api.v1.aiGeneratorPptx.generate_ppt_outline", new_callable=AsyncMock)
    def test_update_success(self, mock_outline, mock_create, client, tmp_output, sample_outline):
        task_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{task_id}_slides.json", sample_outline["slides"])

        mock_outline.return_value = sample_outline
        mock_create.return_value = str(uuid.uuid4())

        resp = client.post(f"/api/v1/pptx/{task_id}/update", json={
            "topic": "新增市场分析章节",
            "slide_count": 5,
        })
        assert resp.status_code == 200
        data = resp.json()
        assert "task_id" in data

    def test_update_no_intermediate_state(self, client, tmp_output):
        """没有中间状态文件时返回 404"""
        resp = client.post(f"/api/v1/pptx/{uuid.uuid4()}/update", json={
            "topic": "测试",
            "slide_count": 5,
        })
        assert resp.status_code == 404

    def test_update_validation_error(self, client, tmp_output, sample_outline):
        task_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{task_id}_slides.json", sample_outline["slides"])
        resp = client.post(f"/api/v1/pptx/{task_id}/update", json={"slide_count": 5})
        assert resp.status_code == 422


# ===========================================================================
# 4. POST /api/v1/pptx/{task_id}/modify
# ===========================================================================

class TestModifyPpt:
    """视觉修改端点测试"""

    @patch("app.utils.pptx.visual_modifier.modify_ppt_visual", new_callable=AsyncMock)
    def test_modify_success(self, mock_modify, client, tmp_output):
        ppt_id = str(uuid.uuid4())
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")

        mock_modify.return_value = {
            "success": True,
            "message": "修改完成",
            "intent": None,
            "analysis": None,
            "preview_count": 0,
        }

        resp = client.post(f"/api/v1/pptx/{ppt_id}/modify", json={
            "user_input": "把标题改成红色",
        })
        assert resp.status_code == 200
        assert resp.json()["success"] is True
        assert "task_id" in resp.json()

    def test_modify_pptx_not_found(self, client, tmp_output):
        resp = client.post(f"/api/v1/pptx/{uuid.uuid4()}/modify", json={
            "user_input": "修改颜色",
        })
        assert resp.status_code == 404

    @patch("app.utils.pptx.visual_modifier.modify_ppt_visual", new_callable=AsyncMock)
    def test_modify_failure(self, mock_modify, client, tmp_output):
        ppt_id = str(uuid.uuid4())
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")

        mock_modify.return_value = {"success": False, "message": "修改失败"}

        resp = client.post(f"/api/v1/pptx/{ppt_id}/modify", json={
            "user_input": "修改",
        })
        assert resp.status_code == 200
        assert resp.json()["success"] is False

    @patch("app.utils.pptx.visual_modifier.modify_ppt_visual", new_callable=AsyncMock)
    def test_modify_exception(self, mock_modify, client, tmp_output):
        ppt_id = str(uuid.uuid4())
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")

        mock_modify.side_effect = RuntimeError("内部错误")

        resp = client.post(f"/api/v1/pptx/{ppt_id}/modify", json={
            "user_input": "修改",
        })
        assert resp.status_code == 500


# ===========================================================================
# 5. GET /api/v1/pptx/{task_id}/analyze
# ===========================================================================

class TestAnalyzePpt:
    """分析端点测试"""

    @patch("app.utils.pptx.visual_modifier.analyze_ppt_for_modification", new_callable=AsyncMock)
    def test_analyze_success(self, mock_analyze, client, tmp_output):
        ppt_id = str(uuid.uuid4())
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")

        mock_analyze.return_value = {
            "success": True,
            "slides_analyzed": 3,
            "analysis": {"layout": "standard"},
        }

        resp = client.get(f"/api/v1/pptx/{ppt_id}/analyze")
        assert resp.status_code == 200
        assert resp.json()["success"] is True

    def test_analyze_not_found(self, client, tmp_output):
        resp = client.get(f"/api/v1/pptx/{uuid.uuid4()}/analyze")
        assert resp.status_code == 404

    @patch("app.utils.pptx.visual_modifier.analyze_ppt_for_modification", new_callable=AsyncMock)
    def test_analyze_with_slide_number(self, mock_analyze, client, tmp_output):
        ppt_id = str(uuid.uuid4())
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")

        mock_analyze.return_value = {"success": True}

        resp = client.get(f"/api/v1/pptx/{ppt_id}/analyze?slide_number=2")
        assert resp.status_code == 200
        mock_analyze.assert_called_once()
        call_kwargs = mock_analyze.call_args
        assert call_kwargs.kwargs.get("slide_number") == 2 or call_kwargs[1].get("slide_number") == 2


# ===========================================================================
# 6. GET /api/v1/pptx/templates
# ===========================================================================

class TestListTemplates:
    """模板列表端点测试"""

    def test_list_templates_success(self, client):
        resp = client.get("/api/v1/pptx/templates")
        assert resp.status_code == 200
        data = resp.json()
        assert "templates" in data
        assert len(data["templates"]) > 0
        for tpl in data["templates"]:
            assert "id" in tpl
            assert "name" in tpl

    def test_list_templates_with_category(self, client):
        resp = client.get("/api/v1/pptx/templates?category=business")
        assert resp.status_code == 200
        for tpl in resp.json()["templates"]:
            assert tpl["id"].startswith("business")


# ===========================================================================
# 7. GET /api/v1/pptx/history
# ===========================================================================

class TestListHistory:
    """历史记录端点测试"""

    def test_history_empty(self, client, tmp_output):
        resp = client.get("/api/v1/pptx/history")
        assert resp.status_code == 200
        data = resp.json()
        assert data["records"] == []
        assert data["total"] == 0

    def test_history_with_records(self, client, tmp_output, sample_outline):
        ppt_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{ppt_id}_slides.json", sample_outline["slides"])

        resp = client.get("/api/v1/pptx/history")
        assert resp.status_code == 200
        data = resp.json()
        assert data["total"] == 1
        assert data["records"][0]["task_id"] == ppt_id

    def test_history_pagination(self, client, tmp_output, sample_outline):
        for _ in range(3):
            ppt_id = str(uuid.uuid4())
            _write_snapshot(tmp_output / f"{ppt_id}_slides.json", sample_outline["slides"])

        resp = client.get("/api/v1/pptx/history?page=1&page_size=2")
        assert resp.status_code == 200
        data = resp.json()
        assert data["total"] == 3
        assert len(data["records"]) == 2

    def test_history_filters_by_user_id(self, client, tmp_output, sample_outline):
        """其他用户的记录不应出现"""
        # 当前用户的记录
        my_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{my_id}_slides.json", sample_outline["slides"], user_id="test_user")
        # 其他用户的记录
        other_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{other_id}_slides.json", sample_outline["slides"], user_id="other_user")

        resp = client.get("/api/v1/pptx/history")
        assert resp.status_code == 200
        data = resp.json()
        assert data["total"] == 1
        assert data["records"][0]["task_id"] == my_id

    def test_history_skips_old_format(self, client, tmp_output, sample_outline):
        """旧格式（纯数组）记录应被跳过"""
        old_id = str(uuid.uuid4())
        with open(tmp_output / f"{old_id}_slides.json", "w") as f:
            json.dump(sample_outline["slides"], f)

        resp = client.get("/api/v1/pptx/history")
        assert resp.status_code == 200
        assert resp.json()["total"] == 0


# ===========================================================================
# 8. DELETE /api/v1/pptx/history/{task_id}
# ===========================================================================

class TestDeleteHistory:
    """删除历史端点测试"""

    def test_delete_success(self, client, tmp_output, sample_outline):
        ppt_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{ppt_id}_slides.json", sample_outline["slides"], user_id="test_user")
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")

        resp = client.delete(f"/api/v1/pptx/history/{ppt_id}")
        assert resp.status_code == 200
        assert resp.json()["success"] is True
        assert not (tmp_output / f"{ppt_id}_slides.json").exists()
        assert not (tmp_output / f"{ppt_id}.pptx").exists()

    def test_delete_not_found(self, client, tmp_output):
        resp = client.delete(f"/api/v1/pptx/history/{uuid.uuid4()}")
        assert resp.status_code == 404

    def test_delete_cleans_all_formats(self, client, tmp_output, sample_outline):
        ppt_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{ppt_id}_slides.json", sample_outline["slides"], user_id="test_user")
        (tmp_output / f"{ppt_id}.html").write_text("<html>")
        (tmp_output / f"{ppt_id}.md").write_text("# md")

        resp = client.delete(f"/api/v1/pptx/history/{ppt_id}")
        assert resp.status_code == 200
        assert not (tmp_output / f"{ppt_id}.html").exists()
        assert not (tmp_output / f"{ppt_id}.md").exists()

    def test_delete_rejects_other_user(self, client, tmp_output, sample_outline):
        """其他用户的记录不能删除"""
        ppt_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{ppt_id}_slides.json", sample_outline["slides"], user_id="other_user")

        resp = client.delete(f"/api/v1/pptx/history/{ppt_id}")
        assert resp.status_code == 403

    def test_delete_rejects_old_format(self, client, tmp_output, sample_outline):
        """旧格式（无 user_id）记录不能删除"""
        ppt_id = str(uuid.uuid4())
        with open(tmp_output / f"{ppt_id}_slides.json", "w") as f:
            json.dump(sample_outline["slides"], f)

        resp = client.delete(f"/api/v1/pptx/history/{ppt_id}")
        assert resp.status_code == 403


# ===========================================================================
# 9. GET /api/v1/pptx/history/stats
# ===========================================================================

class TestPptStats:
    """统计端点测试"""

    def test_stats_empty(self, client, tmp_output):
        resp = client.get("/api/v1/pptx/history/stats")
        assert resp.status_code == 200
        data = resp.json()
        assert data["total"] == 0
        assert data["completed"] == 0

    def test_stats_with_files(self, client, tmp_output, sample_outline):
        ppt_id = str(uuid.uuid4())
        _write_snapshot(tmp_output / f"{ppt_id}_slides.json", sample_outline["slides"])
        _create_pptx_file(tmp_output / f"{ppt_id}.pptx")

        resp = client.get("/api/v1/pptx/history/stats")
        assert resp.status_code == 200
        data = resp.json()
        assert data["total"] == 1
        assert data["completed"] == 1


# ===========================================================================
# 10. POST /api/v1/generate-text (仅大纲)
# ===========================================================================

class TestGenerateText:
    """大纲生成端点测试"""

    @patch("app.agent.ppt_agent.PPTAgent")
    def test_generate_text_success(self, MockAgent, client, mock_outline_obj):
        mock_instance = MockAgent.return_value
        mock_instance.generate_outline = AsyncMock(return_value=mock_outline_obj)

        resp = client.post("/api/v1/generate-text", json={
            "topic": "AI 发展趋势",
            "num_slides": 5,
        })
        assert resp.status_code == 200
        data = resp.json()
        assert "title" in data
        assert "slides" in data
        assert data["total_slides"] > 0

    @patch("app.agent.ppt_agent.PPTAgent")
    def test_generate_text_agent_failure(self, MockAgent, client):
        mock_instance = MockAgent.return_value
        mock_instance.generate_outline = AsyncMock(side_effect=RuntimeError("LLM 失败"))

        resp = client.post("/api/v1/generate-text", json={
            "topic": "测试",
            "num_slides": 5,
        })
        assert resp.status_code == 500

    def test_generate_text_validation_error(self, client):
        """缺少必填字段 topic"""
        resp = client.post("/api/v1/generate-text", json={"num_slides": 5})
        assert resp.status_code == 422
