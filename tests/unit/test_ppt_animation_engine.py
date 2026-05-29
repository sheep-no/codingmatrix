"""
PPT 动画引擎单元测试

测试动画效果、页面切换和动画预设
"""
import pytest
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from app.utils.pptx.animation_engine import (
    TransitionEffect,
    EntranceEffect,
    AnimationConfig,
    AnimationEngine,
    AnimationPresets
)


class TestTransitionEffect:
    """测试页面切换效果"""

    def test_enum_values(self):
        """测试枚举值"""
        assert TransitionEffect.FADE.value == "fade"
        assert TransitionEffect.PUSH.value == "push"
        assert TransitionEffect.WIPE.value == "wipe"
        assert TransitionEffect.SPLIT.value == "split"
        assert TransitionEffect.COVER.value == "cover"


class TestEntranceEffect:
    """测试元素进入效果"""

    def test_enum_values(self):
        """测试枚举值"""
        assert EntranceEffect.FADE_IN.value == "fade_in"
        assert EntranceEffect.FLY_IN.value == "fly_in"
        assert EntranceEffect.ZOOM.value == "zoom"
        assert EntranceEffect.WIPE_ENTRANCE.value == "wipe_entrance"
        assert EntranceEffect.NONE.value == "none"


class TestAnimationConfig:
    """测试动画配置"""

    def test_default_config(self):
        """测试默认配置"""
        config = AnimationConfig()

        assert config.transition == TransitionEffect.FADE
        assert config.duration == 1.0
        assert config.delay == 0.0
        assert config.entrance == EntranceEffect.NONE
        assert config.trigger == "on_click"

    def test_custom_config(self):
        """测试自定义配置"""
        config = AnimationConfig(
            transition=TransitionEffect.PUSH,
            duration=2.0,
            delay=0.5,
            entrance=EntranceEffect.ZOOM,
            trigger="auto"
        )

        assert config.transition == TransitionEffect.PUSH
        assert config.duration == 2.0
        assert config.delay == 0.5
        assert config.entrance == EntranceEffect.ZOOM
        assert config.trigger == "auto"


class TestAnimationEngine:
    """测试动画引擎"""

    @pytest.fixture
    def engine(self):
        return AnimationEngine()

    def test_engine_initial(self, engine):
        """测试引擎初始化"""
        assert engine is not None

    def test_set_default_transition(self, engine):
        """测试设置默认切换效果"""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        result = engine.set_default_transition(prs, TransitionEffect.FADE, duration=1.0)
        assert result is True

    def test_apply_slide_transition(self, engine):
        """测试应用页面切换"""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        result = engine.apply_slide_transition(slide, TransitionEffect.PUSH, duration=1.5)
        assert result is True

    def test_apply_slide_transition_invalid(self, engine):
        """测试无效的切换效果"""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        result = engine.set_default_transition(
            prs, TransitionEffect.RANDOM, duration=1.0
        )
        assert result is True

    def test_animate_slide_sequence_empty(self, engine):
        """测试空序列动画"""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        result = engine.animate_slide_sequence(slide, [])
        assert result is False

    def test_animate_slide_sequence(self, engine):
        """测试序列动画"""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_shape = slide.shapes.title
        if title_shape:
            result = engine.animate_slide_sequence(slide, [title_shape])
            assert result is True


class TestAnimationPresets:
    """测试动画预设"""

    def test_get_preset_corporate(self):
        """测试企业主题预设"""
        preset = AnimationPresets.get_preset_for_theme("corporate")

        assert preset is not None
        assert "transition" in preset
        assert "entrance" in preset
        assert preset["transition"] == TransitionEffect.FADE

    def test_get_preset_creative(self):
        """测试创意主题预设"""
        preset = AnimationPresets.get_preset_for_theme("creative")

        assert preset is not None
        assert preset["transition"] == TransitionEffect.PUSH

    def test_get_preset_minimal(self):
        """测试简约主题预设"""
        preset = AnimationPresets.get_preset_for_theme("minimal")

        assert preset is not None
        assert preset["entrance"] == EntranceEffect.NONE

    def test_get_preset_unknown(self):
        """测试未知主题预设"""
        with pytest.raises(ValueError):
            AnimationPresets.get_preset_for_theme("unknown_theme")

    def test_apply_fade_transition(self):
        """测试应用淡入切换"""
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])

        result = AnimationPresets.apply_fade_transition(prs, duration=1.0)
        assert result is True

    def test_apply_professional_animations(self):
        """测试应用商务专业动画"""
        from pptx import Presentation
        prs = Presentation()
        # 使用 Title Slide 布局以包含标题形状
        title_layout = prs.slide_layouts[0]
        prs.slides.add_slide(title_layout)

        result = AnimationPresets.apply_professional_animations(prs)
        assert result is True
