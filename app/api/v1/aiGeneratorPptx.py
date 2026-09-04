"""
PPT 生成 API - 统一增强版

功能整合：
1. 合并原 aiGeneratorPptx.py 和 ppt_enhanced.py 的功能。
2. 支持会话历史、素材上传绑定、文件权限验证。
3. 支持视觉决策与智能布局 (Visual Decision & Layout)。
4. 支持多种输出格式 (PPTX, PDF, HTML, Markdown)。
5. 支持模板系统、在线预览、代码审查集成。
6. 提供同步/异步生成、下载、预览、幻灯片详情接口。
"""
import asyncio
import html
import json
import logging
import os
import re
import uuid
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import List, Optional, Dict, Any

from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile, File as FastAPIFile, Form, WebSocket, WebSocketDisconnect
from fastapi.responses import StreamingResponse, HTMLResponse, FileResponse
from sqlalchemy.ext.asyncio import AsyncSession
from pydantic import BaseModel, Field

from app.db.database import get_db
from app.db.database import async_session
from app.utils.security import verify_token, verify_token_ws
from app.utils.task_manager import task_manager
from app.schema.task_schema import TaskResponse
from app.services.unified_state_service import (
    StateNotFoundError,
    StateOwnershipError,
    get_owned_task,
    transition_task,
)
from app.schema.ppt_outline import (
    OutlineCreateRequest,
    OutlineDraft,
    OutlineGenerateRequest,
    SlideRegenerateRequest,
    OutlineUpdateRequest,
)
from app.services.ppt_state_service import (
    approve_ppt_outline as persist_approve_ppt_outline,
    create_ppt_outline as persist_create_ppt_outline,
    delete_ppt_outline as persist_delete_ppt_outline,
    get_ppt_quality_report as persist_get_ppt_quality_report,
    get_ppt_outline as persist_get_ppt_outline,
    update_ppt_outline as persist_update_ppt_outline,
)
from app.services.ppt_generation_persistence import (
    build_ppt_trace_context,
    persist_ppt_generation_result,
    save_ppt_stage_checkpoint,
    serialize_quality_report,
)
from app.services.ppt_quality_orchestrator import run_quality_pipeline
from app.utils.pptx.semantic_renderer import build_render_metadata
from app.utils.pptx.semantic_renderer import normalize_slide_type
from app.utils.pptx.templates.manager import TemplateManager
from app.utils.pptx.commercial_content import (
    NARRATIVE_ROLES,
    build_commercial_page_blueprint,
    format_commercial_metadata,
)
from app.celery_app import celery_app

# 视觉决策模块
from app.utils.visual import (
    visual_analyzer,
    image_manager,
    ImageType,
)

from app.models.file import File
from app.models.task import Task
from sqlalchemy import select

# PPT 工具模块
from app.utils.pptx.text_processor import (
    prevent_text_overflow as prevent_text_overflow_v2,
)
from app.utils.pptx.image_search import ImageSearchManager
from app.utils.pptx.ppt_style import PPTStyle, PPT_TEMPLATES
from app.agent.models import DEFAULT_PPT_MODEL

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from PIL import Image

logger = logging.getLogger(__name__)
router = APIRouter(tags=["PPT 生成 (增强版)"])


@router.post("/pptx/outlines", response_model=OutlineDraft, status_code=201)
async def create_ppt_outline(
    req: OutlineCreateRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    """创建用户作用域的 PPT 大纲草稿。"""
    try:
        return await persist_create_ppt_outline(db, str(token.get("sub", "anonymous")), req)
    except StateNotFoundError:
        raise HTTPException(status_code=404, detail="素材文件不存在")
    except ValueError as exc:
        raise HTTPException(status_code=422, detail=str(exc))


@router.get("/pptx/outlines/{outline_id}", response_model=OutlineDraft)
async def get_ppt_outline(
    outline_id: str,
    version: Optional[int] = Query(default=None, ge=1),
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    try:
        return await persist_get_ppt_outline(db, str(token.get("sub", "anonymous")), outline_id, version)
    except (StateNotFoundError, StateOwnershipError):
        raise HTTPException(status_code=404, detail="大纲不存在")


@router.patch("/pptx/outlines/{outline_id}", response_model=OutlineDraft)
async def update_ppt_outline(
    outline_id: str,
    req: OutlineUpdateRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    try:
        return await persist_update_ppt_outline(db, str(token.get("sub", "anonymous")), outline_id, req)
    except (StateNotFoundError, StateOwnershipError):
        raise HTTPException(status_code=404, detail="大纲不存在")


@router.delete("/pptx/outlines/{outline_id}")
async def delete_ppt_outline(
    outline_id: str,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    try:
        await persist_delete_ppt_outline(db, str(token.get("sub", "anonymous")), outline_id)
    except (StateNotFoundError, StateOwnershipError):
        raise HTTPException(status_code=404, detail="大纲不存在")
    return {"deleted": True}


@router.post("/pptx/outlines/{outline_id}/approve", response_model=OutlineDraft)
async def approve_ppt_outline(
    outline_id: str,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    try:
        return await persist_approve_ppt_outline(db, str(token.get("sub", "anonymous")), outline_id)
    except (StateNotFoundError, StateOwnershipError):
        raise HTTPException(status_code=404, detail="大纲不存在")
    except ValueError as exc:
        raise HTTPException(status_code=422, detail=str(exc))


@router.post("/pptx/outlines/{outline_id}/generate", response_model=TaskResponse)
async def generate_ppt_from_approved_outline(
    outline_id: str,
    req: OutlineGenerateRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    """Create a generation task from an approved outline snapshot."""
    try:
        outline = await persist_get_ppt_outline(
            db,
            str(token.get("sub", "anonymous")),
            outline_id,
            req.outline_version,
        )
    except (StateNotFoundError, StateOwnershipError):
        raise HTTPException(status_code=404, detail="大纲不存在")
    if outline.status != "approved":
        raise HTTPException(status_code=409, detail="大纲尚未批准")

    outline_prompt = {
        "id": outline.id,
        "version": outline.version,
        "status": outline.status,
        "title": outline.title,
        "scenario": outline.scenario,
        "template_id": outline.template_id,
        "slides": [slide.model_dump(mode="json") for slide in outline.slides],
    }
    generation_request = PPTGenerationRequest(
        topic=outline.title,
        template=outline.template_id,
        slide_count=len(outline.slides),
        quality="high",
        api_key_token=None,
        options={
            "quality_mode": req.quality_mode,
            "outline_id": outline.id,
            "outline_version": outline.version,
            "approved_outline": outline_prompt,
        },
    )
    return await generate_ppt_task(generation_request, token=token, db=db)


@router.post("/pptx/outlines/{outline_id}/slides/{slide_id}/regenerate", response_model=TaskResponse)
async def regenerate_ppt_slide(
    outline_id: str,
    slide_id: str,
    req: SlideRegenerateRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    """Create a new approved outline version while preserving non-target slides."""
    user_id = str(token.get("sub", "anonymous"))
    try:
        current = await persist_get_ppt_outline(db, user_id, outline_id)
    except StateNotFoundError:
        raise HTTPException(status_code=404, detail="大纲不存在")

    target_index = next((index for index, slide in enumerate(current.slides) if slide.id == slide_id), None)
    if target_index is None:
        raise HTTPException(status_code=404, detail="目标页面不存在")

    slides = list(current.slides)
    if req.slide is not None:
        slides[target_index] = req.slide.model_copy(
            update={"id": slide_id, "position": current.slides[target_index].position}
        )
    updated = await persist_update_ppt_outline(
        db,
        user_id,
        outline_id,
        OutlineUpdateRequest(slides=slides),
    )
    try:
        await persist_approve_ppt_outline(db, user_id, outline_id)
    except ValueError as exc:
        raise HTTPException(status_code=422, detail=str(exc))
    return await generate_ppt_from_approved_outline(
        outline_id,
        OutlineGenerateRequest(quality_mode=req.quality_mode, outline_version=updated.version),
        token,
        db,
    )


@router.get("/pptx/{task_id}/quality-report")
async def get_ppt_quality_report(
    task_id: str,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    try:
        report = await persist_get_ppt_quality_report(db, task_id, int(token.get("sub")))
    except StateNotFoundError:
        raise HTTPException(status_code=404, detail="质量报告不存在")
    return {
        "id": report.id,
        "task_id": report.task_id,
        "outline_id": report.outline_id,
        "outline_version": report.outline_version,
        "quality_mode": report.quality_mode,
        "overall_score": report.overall_score,
        "slide_scores": report.slide_scores_json,
        "issues": report.issues_json,
        "reflow_attempts": report.reflow_attempts_json,
        "degraded_stage": report.degraded_stage,
        "status": report.status,
    }

# =============================================================================
# PPT 全局常量
# =============================================================================

PPT_DEFAULT_MODEL = DEFAULT_PPT_MODEL
PPT_MAX_SLIDES = 50
PPT_OUTPUT_DIR = Path("./pptx_output")
PPT_OWNER_DIR = PPT_OUTPUT_DIR / ".owners"


def _register_ppt_owner(ppt_id: str, user_id: str) -> None:
    """Register the owner of a generated PPT artifact."""
    PPT_OWNER_DIR.mkdir(parents=True, exist_ok=True)
    with open(PPT_OWNER_DIR / f"{ppt_id}.json", "w", encoding="utf-8") as owner_file:
        json.dump({"ppt_id": ppt_id, "user_id": str(user_id)}, owner_file)


def _verify_ppt_owner(ppt_id: str, user_id: str) -> None:
    """Reject access when an artifact has a different registered owner."""
    owner_path = PPT_OWNER_DIR / f"{ppt_id}.json"
    if not owner_path.exists():
        raise HTTPException(status_code=404, detail="PPT 任务不存在或已过期")
    try:
        with open(owner_path, "r", encoding="utf-8") as owner_file:
            owner = json.load(owner_file)
    except (OSError, ValueError, TypeError):
        raise HTTPException(status_code=404, detail="PPT 任务状态不可用")
    if str(owner.get("user_id")) != str(user_id):
        raise HTTPException(status_code=403, detail="无权访问此 PPT")


def _preview_url(ppt_id: str, output_format: "OutputFormat") -> str:
    """Return a preview URL that preserves the generated artifact format."""
    return f"/api/v1/pptx/preview/{ppt_id}?format={output_format.value}"

# 所有合法模板 ID（前端、后端、配置文件共用此列表）
VALID_TEMPLATE_IDS = [
    "modern", "business", "creative", "minimal",
    "academic", "tech", "education", "medical", "elegant",
]

# =============================================================================
# 模型定义
# =============================================================================

class OutputFormat(str, Enum):
    """输出格式枚举"""
    PPTX = "pptx"
    PDF = "pdf"
    HTML = "html"
    MARKDOWN = "markdown"

class PPTGenerationRequest(BaseModel):
    """统一的 PPT 生成请求"""
    # 基础信息
    topic: str = Field(..., description="PPT 主题/提示词", max_length=5000, alias="prompt")
    model: str = Field(default=PPT_DEFAULT_MODEL, description="AI 模型")
    conversation_id: Optional[int] = Field(None, description="会话 ID (用于携带历史上下文)")
    session_id: Optional[str] = Field(None, description="会话 ID (用于素材隔离)")
    material_file_ids: Optional[List[int]] = Field(None, description="已上传素材的文件 ID 列表")
    api_key_token: Optional[str] = Field(None, description="用户 API Key Token（用于从 Redis 获取用户自定义 Key）")
    
    # 增强选项
    template: str = Field(default="modern", description="模板风格")
    slide_count: int = Field(default=10, ge=1, le=PPT_MAX_SLIDES, description="页数")
    output_format: OutputFormat = Field(default=OutputFormat.PPTX, description="输出格式")
    language: str = Field(default="zh-CN", description="语言")
    quality: str = Field(default="high", description="内容质量")
    options: Dict[str, Any] = Field(default_factory=dict, description="高级选项")
    skills: List[str] = Field(default_factory=list, description="Skill 提示词列表")

    class Config:
        populate_by_name = True

class PPTModifyRequest(BaseModel):
    """PPT 修改请求"""
    user_input: str = Field(..., description="用户修改需求（自然语言）", max_length=2000)
    api_key_token: Optional[str] = Field(None, description="用户 API Key Token")
    analyze_before_modify: bool = Field(default=True, description="修改前是否分析当前状态")

    class Config:
        populate_by_name = True


class PPTGenerationResponse(BaseModel):
    """PPT 生成响应"""
    id: str
    status: str
    topic: str
    slide_count: int
    output_format: str
    created_at: str
    download_url: str
    preview_url: str
    slides: Optional[List[Dict[str, Any]]] = None

# =============================================================================
# 模板与样式配置（已提取到 app/utils/pptx/ppt_style.py）
# =============================================================================

# PPT_TEMPLATES 和 PPTStyle 从 app.utils.pptx.ppt_style 导入（见文件顶部）


def add_decorative_header(slide, prs, title_text, style: PPTStyle):
    """添加装饰性页眉"""
    # 顶部装饰条
    header_height = Inches(0.15)
    header_shape = slide.shapes.add_shape(
        1,  # 矩形
        Inches(0), Inches(0),
        prs.slide_width, header_height
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = style.PRIMARY_COLOR
    header_shape.line.fill.background()

    # 标题区域背景
    title_bg = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0.15),
        prs.slide_width, Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = style.PRIMARY_LIGHT
    title_bg.line.fill.background()


def add_page_number(slide, prs, page_num, total_pages, style: PPTStyle):
    """添加页码"""
    # 页脚装饰条
    footer_height = Inches(0.1)
    footer_shape = slide.shapes.add_shape(
        1,
        Inches(0), prs.slide_height - footer_height,
        prs.slide_width, footer_height
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = style.PRIMARY_DARK
    footer_shape.line.fill.background()

    # 页码文本
    page_text = slide.shapes.add_textbox(
        prs.slide_width - Inches(1),
        prs.slide_height - Inches(0.6),
        Inches(0.9), Inches(0.4)
    )
    tf = page_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total_pages}"
    p.font.name = style.FONT_MAIN
    p.font.size = Pt(10)
    p.font.color.rgb = style.PRIMARY_COLOR
    p.alignment = PP_ALIGN.RIGHT


def add_bullet_with_icon(slide, left, top, width, text, level, style: PPTStyle, icon=None):
    """添加带图标的bullet"""
    # bullet 符号
    if level == 0:
        bullet_char = "●"
        bullet_color = style.PRIMARY_COLOR
        font_size = Pt(16)
    else:
        bullet_char = "○"
        bullet_color = style.PRIMARY_LIGHT
        font_size = Pt(14)

    txbox = slide.shapes.add_textbox(left, top, Inches(0.3), Inches(0.3))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = bullet_char
    p.font.name = style.FONT_MAIN
    p.font.size = font_size
    p.font.color.rgb = bullet_color

    # 文本内容
    textbox = slide.shapes.add_textbox(
        left + Inches(0.3), top,
        width - Inches(0.3), Inches(0.5)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = style.FONT_MAIN
    p.font.size = font_size
    p.font.color.rgb = style.TEXT_DARK
    p.space_after = Pt(6)


def style_title_shape(title_shape, text, style: PPTStyle, font_size=Pt(36)):
    """样式化标题形状"""
    title_shape.text = text
    for para in title_shape.text_frame.paragraphs:
        para.font.name = style.FONT_MAIN
        para.font.size = font_size
        para.font.bold = True
        para.font.color.rgb = style.TEXT_WHITE
        para.alignment = PP_ALIGN.LEFT


def add_slide_background(slide, prs, style: PPTStyle, light=True):
    """添加幻灯片背景"""
    bg_color = style.BG_WHITE if light else style.BG_LIGHT_BLUE
    background = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()
    # 移到最底层
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_decorative_corner(slide, prs, style: PPTStyle):
    """添加角落装饰"""
    # 右下角装饰圆
    corner = slide.shapes.add_shape(
        9,  # 椭圆
        prs.slide_width - Inches(2),
        prs.slide_height - Inches(1.5),
        Inches(2), Inches(1.5)
    )
    corner.fill.solid()
    corner.fill.fore_color.rgb = style.ACCENT_LIGHT
    corner.fill.fore_color.brightness = 0.3
    corner.line.fill.background()


def _add_editorial_text(slide, text, left, top, width, height, style, size=18, color=None, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.font.name = style.FONT_MAIN
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color or style.TEXT_DARK
    if align:
        paragraph.alignment = align
    return box


def _add_editorial_card(slide, title, body, left, top, width, height, style, accent=None, filled=False):
    card = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = style.PRIMARY_COLOR if filled else style.BG_WHITE
    card.line.color.rgb = style.PRIMARY_COLOR if filled else style.PRIMARY_LIGHT
    card.line.width = Pt(1)
    marker = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.08), Inches(height))
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent or style.ACCENT_COLOR
    marker.line.fill.background()
    title_color = style.TEXT_WHITE if filled else style.PRIMARY_COLOR
    body_color = style.TEXT_WHITE if filled else style.TEXT_DARK
    if height <= 1.1:
        _add_editorial_text(slide, title, left + 0.28, top + 0.18, 1.75, height - 0.36, style, 12, title_color, True)
        _add_editorial_text(slide, body, left + 2.05, top + 0.14, width - 2.25, height - 0.28, style, 14, body_color)
    else:
        _add_editorial_text(slide, title, left + 0.28, top + 0.2, width - 0.45, 0.35, style, 14, title_color, True)
        _add_editorial_text(slide, body, left + 0.28, top + 0.65, width - 0.45, height - 0.8, style, 16, body_color)


def _add_cropped_picture(slide, image_path, left, top, width, height):
    """Fill a visual box while preserving the source image aspect ratio."""
    with Image.open(image_path) as image:
        source_ratio = image.width / image.height
    box_ratio = width / height
    picture = slide.shapes.add_picture(
        image_path, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if source_ratio > box_ratio:
        picture.crop_left = picture.crop_right = (1 - box_ratio / source_ratio) / 2
    elif source_ratio < box_ratio:
        picture.crop_top = picture.crop_bottom = (1 - source_ratio / box_ratio) / 2
    return picture


def _normalize_slide_items(content):
    """Turn model prose into short presentation-ready statements."""
    if isinstance(content, list):
        raw_items = content
    else:
        raw_items = re.split(r"[\n；;]+", str(content or ""))
    items = []
    for item in raw_items:
        text = re.sub(r"^\s*[-•●\d]+[.、)]?\s*", "", str(item)).strip()
        if text:
            items.append(text)
    return items or ["待补充内容"]


def _commercial_slide_items(slide_data: Dict[str, Any]) -> List[str]:
    """Prefer semantic content blocks and expose their commercial metadata."""
    role = slide_data.get("narrative_role") or "opportunity_map"
    blocks = slide_data.get("content_blocks") or []
    if not isinstance(blocks, list) or not blocks:
        return _normalize_slide_items(slide_data.get("content", []))

    items = []
    for block in blocks:
        if not isinstance(block, dict):
            continue
        content = str(block.get("content", "")).strip()
        if not content:
            continue
        detail = format_commercial_metadata(role, block.get("metadata") or {})
        items.append(f"{content}\n{detail}" if detail else content)
    return items or _normalize_slide_items(slide_data.get("content", []))


def _slide_role(slide_data: Dict[str, Any]) -> str:
    return slide_data.get("narrative_role") or {
        "data": "evidence_story",
        "comparison": "strategic_choice",
        "timeline": "execution_roadmap",
        "process": "execution_roadmap",
        "summary": "decision_close",
        "closing": "decision_close",
    }.get(normalize_slide_type(slide_data.get("slide_type")), "opportunity_map")


def _item_at(items: List[str], index: int) -> str:
    return items[index] if index < len(items) else items[-1]


def _split_commercial_item(item: str) -> tuple[str, str]:
    headline, separator, detail = str(item).partition("\n")
    return headline.strip(), detail.strip() if separator else ""


def _add_academic_note(slide, label, body, left, top, width, height, style, featured=False):
    """Add a paper-like evidence note with a restrained research hierarchy."""
    note = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    note.fill.solid()
    note.fill.fore_color.rgb = style.PRIMARY_COLOR if featured else style.TEXT_WHITE
    note.line.color.rgb = style.PRIMARY_COLOR if featured else style.PRIMARY_LIGHT
    note.line.width = Pt(1)
    label_color = style.TEXT_WHITE if featured else style.PRIMARY_COLOR
    body_color = style.TEXT_WHITE if featured else style.TEXT_DARK
    headline, detail = _split_commercial_item(body)
    if height <= 1.4:
        _add_editorial_text(slide, label, left + 0.25, top + 0.14, 2.05, 0.26, style, 9, label_color, True)
        _add_editorial_text(slide, headline, left + 2.35, top + 0.13, width - 2.6, 0.38, style, 14, body_color, True)
        if detail:
            _add_editorial_text(slide, detail, left + 2.35, top + 0.62, width - 2.6, 0.32, style, 10, body_color)
        return note
    _add_editorial_text(slide, label, left + 0.25, top + 0.18, width - 0.5, 0.28, style, 10, label_color, True)
    _add_editorial_text(slide, headline, left + 0.25, top + 0.62, width - 0.5, 0.62, style, 16, body_color, True)
    if detail:
        _add_editorial_text(slide, detail, left + 0.25, top + 1.38, width - 0.5, max(height - 1.58, 0.25), style, 11, body_color)
    return note


def _render_slide_academic(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render an evidence-led research brief with paper and citation motifs."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    items = _commercial_slide_items(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    key_message = slide_data.get("key_message") or "以证据回答一个明确问题"
    _add_editorial_text(slide, "RESEARCH BRIEF", 0.7, 0.34, 2.7, 0.28, style, 9, style.PRIMARY_COLOR, True)
    _add_editorial_text(slide, f"FIG. {idx:02d}  /  {role.replace('_', ' ').upper()}", 8.8, 0.34, 3.8, 0.28, style, 9, style.PRIMARY_COLOR, True, PP_ALIGN.RIGHT)
    rule = slide.shapes.add_shape(1, Inches(0.7), Inches(0.83), Inches(11.9), Inches(0.025))
    rule.fill.solid(); rule.fill.fore_color.rgb = style.PRIMARY_COLOR; rule.line.fill.background()
    _add_editorial_text(slide, title, 0.7, 1.02, 8.5, 0.68, style, 28, style.PRIMARY_DARK, True)
    _add_editorial_text(slide, key_message, 9.2, 1.05, 3.4, 0.58, style, 12, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)

    if role == "evidence_story":
        _add_academic_note(slide, "PRIMARY FINDING", items[0], 0.7, 1.92, 4.0, 4.72, style, True)
        for number in range(1, 4):
            _add_academic_note(slide, f"EVIDENCE [{number}]", _item_at(items, number), 5.05, 1.92 + (number - 1) * 1.57, 7.55, 1.28, style)
    elif role == "strategic_choice":
        _add_academic_note(slide, "HYPOTHESIS A", items[0], 0.7, 1.92, 5.55, 2.72, style)
        _add_academic_note(slide, "HYPOTHESIS B", _item_at(items, 1), 6.65, 1.92, 5.95, 2.72, style)
        _add_academic_note(slide, "RESEARCH CONCLUSION", "；".join(items[2:]) or key_message, 2.15, 5.02, 9.05, 1.38, style, True)
    elif role == "execution_roadmap":
        protocol = slide.shapes.add_shape(1, Inches(1.28), Inches(2.18), Inches(0.035), Inches(3.72))
        protocol.fill.solid(); protocol.fill.fore_color.rgb = style.PRIMARY_COLOR; protocol.line.fill.background()
        for number in range(3):
            top = 1.9 + number * 1.48
            marker = slide.shapes.add_shape(9, Inches(1.03), Inches(top + 0.3), Inches(0.52), Inches(0.52))
            marker.fill.solid(); marker.fill.fore_color.rgb = style.PRIMARY_COLOR; marker.line.fill.background()
            _add_editorial_text(slide, str(number + 1), 1.03, top + 0.42, 0.52, 0.2, style, 11, style.TEXT_WHITE, True, PP_ALIGN.CENTER)
            _add_academic_note(slide, f"PROTOCOL / {('PILOT', 'VALIDATE', 'SCALE')[number]}", _item_at(items, number), 1.85, top, 10.45, 1.12, style, number == 2)
        _add_editorial_text(slide, _item_at(items, 3), 1.85, 6.42, 10.45, 0.32, style, 10, style.TEXT_GRAY)
    elif role == "decision_close":
        _add_academic_note(slide, "CONCLUSION", items[0], 0.7, 1.92, 11.9, 2.05, style, True)
        for number in range(1, 4):
            _add_academic_note(slide, f"IMPLICATION [{number}]", _item_at(items, number), 0.7 + (number - 1) * 4.05, 4.35, 3.8, 2.05, style)
    else:
        _add_academic_note(slide, "RESEARCH QUESTION", items[0], 0.7, 1.92, 7.15, 4.72, style, True)
        _add_academic_note(slide, "OBSERVATION [1]", _item_at(items, 1), 8.2, 1.92, 4.4, 2.05, style)
        _add_academic_note(slide, "OBSERVATION [2]", _item_at(items, 2), 8.2, 4.35, 4.4, 2.29, style)
        _add_editorial_text(slide, _item_at(items, 3), 1.02, 5.92, 6.5, 0.38, style, 10, style.TEXT_WHITE, True)
    sources = slide_data.get("evidence_sources") or []
    source_label = sources[0].get("title", "已批准大纲") if sources else "待补：基线数据、用户访谈或公开研究"
    _add_editorial_text(slide, f"SOURCE / {source_label[:72]}", 0.7, 6.92, 10.2, 0.24, style, 9, style.PRIMARY_COLOR)
    return slide


def _add_learning_card(slide, label, body, left, top, width, height, style, featured=False):
    """Add a workshop card with a clear label, lesson and supporting detail."""
    card = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = style.PRIMARY_COLOR if featured else style.TEXT_WHITE
    card.line.color.rgb = style.PRIMARY_COLOR if featured else style.PRIMARY_LIGHT
    card.line.width = Pt(1.25)
    label_color = style.TEXT_WHITE if featured else style.PRIMARY_COLOR
    text_color = style.TEXT_WHITE if featured else style.TEXT_DARK
    headline, detail = _split_commercial_item(body)
    _add_editorial_text(slide, label, left + 0.28, top + 0.18, width - 0.56, 0.3, style, 9, label_color, True)
    _add_editorial_text(slide, headline, left + 0.28, top + 0.65, width - 0.56, 0.62, style, 15, text_color, True)
    if detail and height >= 1.75:
        _add_editorial_text(slide, detail, left + 0.28, top + 1.4, width - 0.56, height - 1.62, style, 12, text_color)
    return card


def _render_slide_education(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render an adult learning workshop with objectives, practice and checks."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    items = _commercial_slide_items(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    key_message = slide_data.get("key_message") or "先理解，再练习，最后带走行动"
    chapter = slide.shapes.add_shape(5, Inches(0.68), Inches(0.38), Inches(1.35), Inches(0.42))
    chapter.fill.solid(); chapter.fill.fore_color.rgb = style.PRIMARY_COLOR; chapter.line.fill.background()
    _add_editorial_text(slide, f"LESSON {idx:02d}", 0.68, 0.46, 1.35, 0.22, style, 9, style.TEXT_WHITE, True, PP_ALIGN.CENTER)
    _add_editorial_text(slide, role.replace("_", " ").upper(), 2.25, 0.46, 4.5, 0.25, style, 9, style.PRIMARY_COLOR, True)
    _add_editorial_text(slide, f"{idx + 1} / {total_slides}", 11.68, 0.46, 0.9, 0.25, style, 9, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    _add_editorial_text(slide, title, 0.68, 1.02, 8.65, 0.68, style, 29, style.TEXT_DARK, True)
    _add_editorial_text(slide, key_message, 9.25, 1.05, 3.33, 0.58, style, 12, style.PRIMARY_DARK, True, PP_ALIGN.RIGHT)

    if role == "evidence_story":
        _add_learning_card(slide, "KEY LEARNING", items[0], 0.68, 1.92, 4.2, 4.72, style, True)
        for number in range(1, 4):
            _add_learning_card(slide, f"PROOF POINT {number}", _item_at(items, number), 5.18, 1.92 + (number - 1) * 1.57, 7.4, 1.28, style)
    elif role == "strategic_choice":
        _add_learning_card(slide, "PRACTICE A", items[0], 0.68, 1.92, 5.55, 2.75, style)
        _add_learning_card(slide, "PRACTICE B", _item_at(items, 1), 6.65, 1.92, 5.93, 2.75, style, True)
        _add_learning_card(slide, "LEARNING CHECK", "；".join(items[2:]) or key_message, 2.1, 5.03, 9.15, 1.38, style, True)
    elif role == "execution_roadmap":
        rail = slide.shapes.add_shape(1, Inches(1.25), Inches(3.16), Inches(10.7), Inches(0.06))
        rail.fill.solid(); rail.fill.fore_color.rgb = style.PRIMARY_LIGHT; rail.line.fill.background()
        for number in range(3):
            left = 0.68 + number * 4.08
            badge = slide.shapes.add_shape(9, Inches(left + 1.5), Inches(2.87), Inches(0.62), Inches(0.62))
            badge.fill.solid(); badge.fill.fore_color.rgb = style.PRIMARY_COLOR; badge.line.fill.background()
            _add_editorial_text(slide, str(number + 1), left + 1.5, 3.0, 0.62, 0.22, style, 12, style.TEXT_WHITE, True, PP_ALIGN.CENTER)
            _add_learning_card(slide, f"MODULE / {('TRY', 'PRACTICE', 'APPLY')[number]}", _item_at(items, number), left, 3.55, 3.65, 2.35, style, number == 2)
        _add_editorial_text(slide, _item_at(items, 3), 1.3, 6.3, 10.7, 0.38, style, 11, style.PRIMARY_DARK, True, PP_ALIGN.CENTER)
    elif role == "decision_close":
        _add_learning_card(slide, "TAKEAWAY", items[0], 0.68, 1.92, 11.9, 2.02, style, True)
        for number in range(1, 4):
            _add_learning_card(slide, f"NEXT STEP {number}", _item_at(items, number), 0.68 + (number - 1) * 4.08, 4.35, 3.66, 2.05, style)
    else:
        _add_learning_card(slide, "WHY THIS MATTERS", items[0], 0.68, 1.92, 6.95, 4.72, style, True)
        _add_learning_card(slide, "NOTICE", _item_at(items, 1), 7.98, 1.92, 4.6, 2.05, style)
        _add_learning_card(slide, "TRY NEXT", _item_at(items, 2), 7.98, 4.35, 4.6, 2.29, style)
        _add_editorial_text(slide, _item_at(items, 3), 1.0, 5.9, 6.1, 0.42, style, 10, style.TEXT_WHITE, True)
    return slide


def _add_clinical_card(slide, label, body, left, top, width, height, style, featured=False):
    """Add a calm clinical brief card with evidence-first hierarchy."""
    card = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = style.PRIMARY_COLOR if featured else style.TEXT_WHITE
    card.line.color.rgb = style.PRIMARY_COLOR if featured else style.PRIMARY_LIGHT
    card.line.width = Pt(1)
    label_color = style.TEXT_WHITE if featured else style.PRIMARY_COLOR
    text_color = style.TEXT_WHITE if featured else style.TEXT_DARK
    headline, detail = _split_commercial_item(body)
    _add_editorial_text(slide, label, left + 0.3, top + 0.2, width - 0.6, 0.28, style, 9, label_color, True)
    _add_editorial_text(slide, headline, left + 0.3, top + 0.68, width - 0.6, 0.6, style, 15, text_color, True)
    if detail and height >= 1.7:
        _add_editorial_text(slide, detail, left + 0.3, top + 1.42, width - 0.6, height - 1.62, style, 12, text_color)
    return card


def _render_slide_medical(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render a clinical brief with findings, options, pathway and care actions."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    items = _commercial_slide_items(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    key_message = slide_data.get("key_message") or "以证据明确判断，以路径推进照护"
    _add_editorial_text(slide, "CLINICAL BRIEF", 0.72, 0.42, 2.4, 0.28, style, 10, style.PRIMARY_COLOR, True)
    _add_editorial_text(slide, f"CASE / {idx:02d}", 10.75, 0.42, 1.8, 0.28, style, 9, style.TEXT_GRAY, True, PP_ALIGN.RIGHT)
    _add_editorial_text(slide, title, 0.72, 0.98, 8.4, 0.7, style, 28, style.TEXT_DARK, True)
    _add_editorial_text(slide, key_message, 9.28, 1.02, 3.3, 0.58, style, 12, style.PRIMARY_DARK, True, PP_ALIGN.RIGHT)
    rule = slide.shapes.add_shape(1, Inches(0.72), Inches(1.78), Inches(11.88), Inches(0.025))
    rule.fill.solid(); rule.fill.fore_color.rgb = style.PRIMARY_LIGHT; rule.line.fill.background()

    if role == "evidence_story":
        _add_clinical_card(slide, "PRIMARY FINDING", items[0], 0.72, 2.05, 4.25, 4.45, style, True)
        for number in range(1, 4):
            _add_clinical_card(slide, f"EVIDENCE / {number}", _item_at(items, number), 5.25, 2.05 + (number - 1) * 1.5, 7.35, 1.23, style)
    elif role == "strategic_choice":
        _add_clinical_card(slide, "OPTION A", items[0], 0.72, 2.05, 5.5, 2.65, style)
        _add_clinical_card(slide, "RECOMMENDED OPTION", _item_at(items, 1), 6.65, 2.05, 5.95, 2.65, style, True)
        _add_clinical_card(slide, "CLINICAL RATIONALE", "；".join(items[2:]) or key_message, 2.05, 5.03, 9.2, 1.35, style, True)
    elif role == "execution_roadmap":
        _add_editorial_text(slide, "CARE PATHWAY", 0.72, 2.18, 2.4, 0.3, style, 10, style.PRIMARY_COLOR, True)
        rail = slide.shapes.add_shape(1, Inches(1.0), Inches(3.2), Inches(11.0), Inches(0.04))
        rail.fill.solid(); rail.fill.fore_color.rgb = style.PRIMARY_LIGHT; rail.line.fill.background()
        for number in range(3):
            left = 0.72 + number * 4.08
            marker = slide.shapes.add_shape(9, Inches(left + 1.48), Inches(2.88), Inches(0.64), Inches(0.64))
            marker.fill.solid(); marker.fill.fore_color.rgb = style.PRIMARY_COLOR; marker.line.fill.background()
            _add_editorial_text(slide, str(number + 1), left + 1.48, 3.03, 0.64, 0.2, style, 12, style.TEXT_WHITE, True, PP_ALIGN.CENTER)
            _add_clinical_card(slide, ("ASSESS", "TREAT", "FOLLOW UP")[number], _item_at(items, number), left, 3.58, 3.65, 2.15, style, number == 2)
        _add_editorial_text(slide, _item_at(items, 3), 1.25, 6.25, 10.8, 0.35, style, 11, style.PRIMARY_DARK, True, PP_ALIGN.CENTER)
    elif role == "decision_close":
        _add_clinical_card(slide, "CARE DECISION", items[0], 0.72, 2.05, 11.88, 1.95, style, True)
        for number in range(1, 4):
            _add_clinical_card(slide, f"ACTION / {number}", _item_at(items, number), 0.72 + (number - 1) * 4.08, 4.35, 3.66, 2.15, style)
    else:
        _add_clinical_card(slide, "CLINICAL SIGNAL", items[0], 0.72, 2.05, 6.95, 4.45, style, True)
        _add_clinical_card(slide, "PATIENT IMPACT", _item_at(items, 1), 7.98, 2.05, 4.62, 1.95, style)
        _add_clinical_card(slide, "NEXT REVIEW", _item_at(items, 2), 7.98, 4.35, 4.62, 2.15, style)
    sources = slide_data.get("evidence_sources") or []
    source_label = sources[0].get("title", "已批准大纲") if sources else "待补：临床指南、病例数据或用户访谈"
    _add_editorial_text(slide, f"SOURCE / {source_label[:68]}", 0.72, 6.9, 10.8, 0.28, style, 10, style.PRIMARY_COLOR)
    _add_editorial_text(slide, f"{idx + 1} / {total_slides}", 11.68, 6.92, 0.9, 0.24, style, 9, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    return slide


def _add_elegant_note(slide, label, body, left, top, width, height, style, featured=False):
    """Add a restrained editorial note for executive presentations."""
    note = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    note.fill.solid()
    note.fill.fore_color.rgb = style.PRIMARY_COLOR if featured else style.TEXT_WHITE
    note.line.color.rgb = style.PRIMARY_COLOR if featured else style.PRIMARY_LIGHT
    note.line.width = Pt(0.75)
    label_color = style.TEXT_WHITE if featured else style.PRIMARY_COLOR
    text_color = style.TEXT_WHITE if featured else style.TEXT_DARK
    headline, detail = _split_commercial_item(body)
    _add_editorial_text(slide, label, left + 0.28, top + 0.18, width - 0.56, 0.28, style, 9, label_color, True)
    divider = slide.shapes.add_shape(1, Inches(left + 0.28), Inches(top + 0.58), Inches(min(width - 0.56, 1.05)), Inches(0.025))
    divider.fill.solid(); divider.fill.fore_color.rgb = label_color; divider.line.fill.background()
    if height < 1.75:
        _add_editorial_text(slide, headline, left + 0.28, top + 0.67, width - 0.56, max(height - 0.84, 0.28), style, 12, text_color, featured)
    else:
        _add_editorial_text(slide, headline, left + 0.28, top + 0.72, width - 0.56, 0.62, style, 15, text_color, True)
    if detail and height >= 1.75:
        _add_editorial_text(slide, detail, left + 0.28, top + 1.48, width - 0.56, height - 1.68, style, 12, text_color)
    return note


def _render_slide_elegant(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render a premium board memo with editorial hierarchy and fine rules."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    items = _commercial_slide_items(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    key_message = slide_data.get("key_message") or "聚焦关键判断，形成清晰决议"
    _add_editorial_text(slide, "EXECUTIVE MEMO", 0.72, 0.42, 2.7, 0.28, style, 9, style.PRIMARY_COLOR, True)
    _add_editorial_text(slide, f"{idx + 1:02d}", 11.78, 0.36, 0.8, 0.38, style, 16, style.PRIMARY_COLOR, True, PP_ALIGN.RIGHT)
    _add_editorial_text(slide, title, 0.72, 0.95, 8.4, 0.72, style, 29, style.TEXT_DARK, True)
    _add_editorial_text(slide, key_message, 9.15, 1.0, 3.43, 0.58, style, 11, style.PRIMARY_DARK, False, PP_ALIGN.RIGHT)
    top_rule = slide.shapes.add_shape(1, Inches(0.72), Inches(1.78), Inches(11.86), Inches(0.02))
    top_rule.fill.solid(); top_rule.fill.fore_color.rgb = style.PRIMARY_COLOR; top_rule.line.fill.background()

    if role == "evidence_story":
        _add_elegant_note(slide, "THE EVIDENCE", items[0], 0.72, 2.08, 4.05, 4.35, style, True)
        for number in range(1, 4):
            top = 2.08 + (number - 1) * 1.48
            _add_editorial_text(slide, f"0{number}", 5.25, top + 0.05, 0.55, 0.35, style, 16, style.PRIMARY_COLOR, True)
            _add_elegant_note(slide, "EVIDENCE NOTE", _item_at(items, number), 5.95, top, 6.63, 1.18, style)
    elif role == "strategic_choice":
        _add_elegant_note(slide, "PATH A", items[0], 0.72, 2.08, 4.7, 2.65, style)
        _add_elegant_note(slide, "PATH B", _item_at(items, 1), 5.78, 2.08, 6.8, 2.65, style)
        _add_elegant_note(slide, "BOARD RECOMMENDATION", "；".join(items[2:]) or key_message, 2.08, 5.08, 9.18, 1.3, style, True)
    elif role == "execution_roadmap":
        spine = slide.shapes.add_shape(1, Inches(2.15), Inches(2.18), Inches(0.025), Inches(4.05))
        spine.fill.solid(); spine.fill.fore_color.rgb = style.PRIMARY_COLOR; spine.line.fill.background()
        for number in range(3):
            top = 2.02 + number * 1.42
            _add_editorial_text(slide, f"0{number + 1}", 0.72, top + 0.28, 0.7, 0.35, style, 16, style.PRIMARY_COLOR, True)
            marker = slide.shapes.add_shape(9, Inches(1.98), Inches(top + 0.31), Inches(0.36), Inches(0.36))
            marker.fill.solid(); marker.fill.fore_color.rgb = style.PRIMARY_COLOR; marker.line.fill.background()
            _add_elegant_note(slide, ("COMMIT", "PROVE", "EXPAND")[number], _item_at(items, number), 2.72, top, 9.86, 1.12, style, number == 2)
        _add_editorial_text(slide, _item_at(items, 3), 2.72, 6.38, 9.86, 0.32, style, 11, style.PRIMARY_DARK, True)
    elif role == "decision_close":
        _add_elegant_note(slide, "RESOLUTION", items[0], 0.72, 2.08, 11.86, 1.86, style, True)
        for number in range(1, 4):
            _add_elegant_note(slide, f"COMMITMENT / 0{number}", _item_at(items, number), 0.72 + (number - 1) * 4.08, 4.38, 3.66, 2.02, style)
    else:
        _add_editorial_text(slide, "01", 0.72, 2.12, 1.4, 0.9, style, 44, style.PRIMARY_COLOR, True)
        _add_elegant_note(slide, "PRIMARY SIGNAL", items[0], 2.05, 2.08, 6.12, 4.35, style, True)
        _add_elegant_note(slide, "IMPLICATION", _item_at(items, 1), 8.55, 2.08, 4.03, 1.88, style)
        _add_elegant_note(slide, "VALIDATION", _item_at(items, 2), 8.55, 4.38, 4.03, 2.05, style)
    _add_editorial_text(slide, "PRIVATE & CONFIDENTIAL", 0.72, 6.92, 3.0, 0.22, style, 8, style.TEXT_GRAY, True)
    return slide


def _add_modern_tile(slide, label, body, left, top, width, height, style, featured=False):
    tile = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    tile.fill.solid()
    tile.fill.fore_color.rgb = style.PRIMARY_COLOR if featured else style.BG_WHITE
    tile.line.color.rgb = style.PRIMARY_COLOR if featured else style.PRIMARY_LIGHT
    tile.line.width = Pt(1.25)
    text_color = style.TEXT_WHITE if featured else style.TEXT_DARK
    label_color = style.TEXT_WHITE if featured else style.PRIMARY_COLOR
    _add_editorial_text(slide, label, left + 0.28, top + 0.2, width - 0.56, 0.3, style, 11, label_color, True)
    headline, detail = _split_commercial_item(body)
    if detail and height >= 2.0:
        _add_editorial_text(slide, headline, left + 0.28, top + 0.65, width - 0.56, 0.68, style, 17, text_color, True)
        _add_editorial_text(slide, detail, left + 0.28, top + 1.5, width - 0.56, height - 1.72, style, 11, text_color)
    else:
        _add_editorial_text(slide, body, left + 0.28, top + 0.62, width - 0.56, height - 0.82, style, 15, text_color, featured)
    return tile


def _render_slide_modern(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render a modular product-story layout with strong metric tiles."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    items = _commercial_slide_items(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    key_message = slide_data.get("key_message") or "用清晰证据推动下一步决策"
    add_page_number(slide, prs, idx + 1, total_slides, style)
    _add_editorial_text(slide, f"{idx:02d} / {role.replace('_', ' ').upper()}", 0.72, 0.38, 4.8, 0.28, style, 10, style.PRIMARY_COLOR, True)
    _add_editorial_text(slide, title, 0.72, 0.72, 8.6, 0.58, style, 28, style.TEXT_DARK, True)
    _add_editorial_text(slide, key_message, 9.05, 0.68, 3.55, 0.55, style, 12, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)

    if role == "evidence_story":
        _add_modern_tile(slide, "PRIMARY SIGNAL", items[0], 0.72, 1.58, 4.0, 4.85, style, True)
        for number in range(1, 4):
            _add_modern_tile(slide, f"EVIDENCE {number:02d}", _item_at(items, number), 5.05, 1.58 + (number - 1) * 1.62, 7.55, 1.32, style)
    elif role == "strategic_choice":
        _add_modern_tile(slide, "OPTION 01", items[0], 0.72, 1.68, 5.7, 2.8, style)
        _add_modern_tile(slide, "OPTION 02 / RECOMMENDED", _item_at(items, 1), 6.9, 1.68, 5.7, 2.8, style, True)
        _add_modern_tile(slide, "DECISION", "；".join(items[2:]) or key_message, 2.65, 4.92, 8.05, 1.35, style, True)
    elif role == "execution_roadmap":
        rail = slide.shapes.add_shape(1, Inches(1.0), Inches(3.05), Inches(11.25), Inches(0.06))
        rail.fill.solid(); rail.fill.fore_color.rgb = style.PRIMARY_LIGHT; rail.line.fill.background()
        for number in range(3):
            left = 0.72 + number * 4.15
            _add_modern_tile(slide, f"0{number + 1} / {('试点', '扩展', '规模化')[number]}", _item_at(items, number), left, 1.68 + (number % 2) * 1.45, 3.72, 2.25, style, number == 2)
        _add_editorial_text(slide, _item_at(items, 3), 2.1, 5.72, 9.1, 0.52, style, 13, style.PRIMARY_COLOR, True, PP_ALIGN.CENTER)
    elif role == "decision_close":
        _add_modern_tile(slide, "PRIORITY", items[0], 0.72, 1.58, 6.2, 4.82, style, True)
        for number in range(1, 4):
            _add_modern_tile(slide, f"ACTION {number:02d}", _item_at(items, number), 7.28, 1.58 + (number - 1) * 1.62, 5.32, 1.32, style)
    else:
        _add_modern_tile(slide, "OPPORTUNITY / 01", items[0], 0.72, 1.58, 7.05, 4.82, style, True)
        _add_modern_tile(slide, "SIGNAL / 02", _item_at(items, 1), 8.1, 1.58, 4.5, 2.1, style)
        _add_modern_tile(slide, "SIGNAL / 03", _item_at(items, 2), 8.1, 4.02, 4.5, 2.38, style)
        _add_editorial_text(slide, _item_at(items, 3), 1.08, 5.65, 6.35, 0.38, style, 11, style.TEXT_WHITE, True)
    return slide


def _add_minimal_rule(slide, left, top, width, style, weight=1):
    rule = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.02 * weight))
    rule.fill.solid(); rule.fill.fore_color.rgb = style.PRIMARY_LIGHT; rule.line.fill.background()
    return rule


def _render_slide_minimal(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render a restrained editorial layout built from type and rules."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    items = _commercial_slide_items(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    key_message = slide_data.get("key_message") or "一个页面，一个明确判断"
    _add_editorial_text(slide, f"{idx:02d}", 0.65, 0.42, 0.7, 0.4, style, 12, style.TEXT_DARK, True)
    _add_editorial_text(slide, role.replace("_", " / ").upper(), 1.5, 0.42, 4.4, 0.4, style, 10, style.TEXT_GRAY, True)
    _add_editorial_text(slide, f"{idx + 1} / {total_slides}", 11.8, 0.42, 0.85, 0.4, style, 10, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    _add_minimal_rule(slide, 0.65, 1.0, 12.0, style, 2)
    _add_editorial_text(slide, title, 0.65, 1.22, 8.65, 0.72, style, 30, style.TEXT_DARK, True)
    _add_editorial_text(slide, key_message, 9.45, 1.28, 3.2, 0.62, style, 12, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)

    if role == "execution_roadmap":
        timeline = slide.shapes.add_shape(1, Inches(2.13), Inches(2.25), Inches(0.03), Inches(3.55))
        timeline.fill.solid(); timeline.fill.fore_color.rgb = style.TEXT_DARK; timeline.line.fill.background()
        for number in range(3):
            top = 2.08 + number * 1.4
            _add_editorial_text(slide, f"0{number + 1}", 0.75, top, 0.7, 0.45, style, 18, style.TEXT_DARK, True)
            _add_editorial_text(slide, ("试点", "扩展", "规模化")[number], 1.45, top, 0.7, 0.45, style, 11, style.TEXT_GRAY, True)
            _add_editorial_text(slide, _item_at(items, number), 2.65, top - 0.05, 9.7, 0.75, style, 15, style.TEXT_DARK)
            _add_minimal_rule(slide, 2.65, top + 0.85, 9.7, style)
    elif role == "decision_close":
        _add_editorial_text(slide, items[0], 0.65, 2.25, 11.3, 1.25, style, 25, style.TEXT_DARK, True)
        _add_minimal_rule(slide, 0.65, 3.72, 12.0, style)
        for number in range(1, 4):
            left = 0.65 + (number - 1) * 4.0
            _add_editorial_text(slide, f"ACTION / 0{number}", left, 4.05, 3.55, 0.3, style, 10, style.TEXT_GRAY, True)
            _add_editorial_text(slide, _item_at(items, number), left, 4.55, 3.55, 1.25, style, 15, style.TEXT_DARK)
    else:
        labels = {
            "strategic_choice": ("A / OPTION", "B / OPTION", "DECISION /"),
            "evidence_story": ("01 / SIGNAL", "02 / SIGNAL", "03 / MEANING"),
            "opportunity_map": ("01 / CHANGE", "02 / WINDOW", "03 / TEST"),
        }.get(role, ("01 / POINT", "02 / POINT", "03 / POINT"))
        for number in range(3):
            top = 2.28 + number * 1.25
            if role == "strategic_choice" and number == 2:
                decision = slide.shapes.add_shape(1, Inches(0.65), Inches(top - 0.18), Inches(12.0), Inches(1.02))
                decision.fill.solid(); decision.fill.fore_color.rgb = style.TEXT_DARK; decision.line.fill.background()
                _add_editorial_text(slide, labels[number], 0.92, top, 1.2, 0.35, style, 10, style.TEXT_WHITE, True)
                _add_editorial_text(slide, _item_at(items, number), 2.2, top - 0.08, 9.95, 0.68, style, 16, style.TEXT_WHITE, True)
                continue
            _add_editorial_text(slide, labels[number], 0.65, top, 1.35, 0.35, style, 10, style.TEXT_GRAY, True)
            _add_editorial_text(slide, _item_at(items, number), 2.2, top - 0.08, 10.1, 0.68, style, 16, style.TEXT_DARK, False)
            _add_minimal_rule(slide, 2.2, top + 0.72, 10.1, style)
        _add_editorial_text(slide, _item_at(items, 3), 2.2, 6.15, 10.1, 0.45, style, 12, style.TEXT_GRAY)
    return slide


def _add_tech_panel(slide, label, body, left, top, width, height, style, featured=False):
    panel = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid(); panel.fill.fore_color.rgb = style.PRIMARY_DARK if featured else style.BG_WHITE
    panel.line.color.rgb = style.ACCENT_COLOR; panel.line.width = Pt(1.5 if featured else 0.75)
    corner = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.42), Inches(0.06))
    corner.fill.solid(); corner.fill.fore_color.rgb = style.ACCENT_COLOR; corner.line.fill.background()
    _add_editorial_text(slide, label, left + 0.25, top + 0.2, width - 0.5, 0.3, style, 10, style.ACCENT_LIGHT, True)
    headline, detail = _split_commercial_item(body)
    if detail and height >= 2.0:
        _add_editorial_text(slide, headline, left + 0.25, top + 0.65, width - 0.5, 0.68, style, 16, style.TEXT_WHITE, True)
        _add_editorial_text(slide, detail, left + 0.25, top + 1.48, width - 0.5, height - 1.7, style, 11, style.ACCENT_LIGHT)
    else:
        _add_editorial_text(slide, body, left + 0.25, top + 0.62, width - 0.5, height - 0.82, style, 14, style.TEXT_WHITE, featured)
    return panel


def _render_slide_tech(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render a dark technical command-center layout."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    items = _commercial_slide_items(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    key_message = slide_data.get("key_message") or "从信号进入可执行决策"
    _add_editorial_text(slide, "SYSTEM / STRATEGY", 0.7, 0.35, 3.0, 0.25, style, 9, style.ACCENT_LIGHT, True)
    _add_editorial_text(slide, f"NODE {idx:02d}", 10.9, 0.35, 1.7, 0.25, style, 9, style.ACCENT_LIGHT, True, PP_ALIGN.RIGHT)
    _add_editorial_text(slide, title, 0.7, 0.72, 8.75, 0.62, style, 28, style.TEXT_WHITE, True)
    _add_editorial_text(slide, key_message, 9.15, 0.74, 3.45, 0.54, style, 12, style.ACCENT_LIGHT, True, PP_ALIGN.RIGHT)
    header_line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.42), Inches(11.9), Inches(0.04))
    header_line.fill.solid(); header_line.fill.fore_color.rgb = style.ACCENT_COLOR; header_line.line.fill.background()

    if role == "evidence_story":
        _add_tech_panel(slide, "CORE / SIGNAL", items[0], 0.7, 1.78, 4.15, 4.9, style, True)
        for number in range(1, 4):
            _add_tech_panel(slide, f"DATA / 0{number}", _item_at(items, number), 5.2, 1.78 + (number - 1) * 1.62, 7.4, 1.28, style)
    elif role == "strategic_choice":
        _add_tech_panel(slide, "NODE A / OPTION", items[0], 0.7, 1.82, 5.45, 2.75, style)
        _add_tech_panel(slide, "NODE B / RECOMMENDED", _item_at(items, 1), 7.15, 1.82, 5.45, 2.75, style, True)
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.15), Inches(3.18), Inches(7.15), Inches(3.18))
        connector.line.color.rgb = style.ACCENT_COLOR; connector.line.width = Pt(2)
        _add_tech_panel(slide, "LOCK / RECOMMENDATION", "；".join(items[2:]) or key_message, 2.35, 5.05, 8.65, 1.35, style, True)
    elif role == "execution_roadmap":
        for number in range(3):
            left = 0.7 + number * 4.18
            top = 2.0 + number * 0.42
            _add_tech_panel(slide, f"PHASE / 0{number + 1}", _item_at(items, number), left, top, 3.65, 3.25, style, number == 2)
            if number < 2:
                connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left + 3.65), Inches(top + 1.62), Inches(left + 4.18), Inches(top + 2.04))
                connector.line.color.rgb = style.ACCENT_COLOR; connector.line.width = Pt(2)
        _add_editorial_text(slide, _item_at(items, 3), 1.2, 6.35, 10.9, 0.35, style, 11, style.ACCENT_LIGHT, True, PP_ALIGN.CENTER)
    elif role == "decision_close":
        _add_tech_panel(slide, "PRIORITY / LOCKED", items[0], 0.7, 1.82, 7.0, 4.82, style, True)
        for number in range(1, 4):
            _add_tech_panel(slide, f"EXEC / 0{number}", _item_at(items, number), 8.05, 1.82 + (number - 1) * 1.62, 4.55, 1.28, style)
    else:
        for number in range(3):
            left = 0.7 + number * 4.18
            _add_tech_panel(slide, f"SIGNAL / 0{number + 1}", _item_at(items, number), left, 2.0 + (number % 2) * 0.7, 3.65, 3.45, style, number == 0)
        _add_editorial_text(slide, _item_at(items, 3), 2.0, 6.2, 9.35, 0.42, style, 11, style.ACCENT_LIGHT, True, PP_ALIGN.CENTER)
    return slide


def _add_creative_block(
    slide, title, body, left, top, width, height, style, fill_color,
    inverted=True, body_font=17,
):
    block = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    block.fill.solid(); block.fill.fore_color.rgb = fill_color; block.line.fill.background()
    title_color = style.TEXT_WHITE if inverted else style.PRIMARY_COLOR
    body_color = style.TEXT_WHITE if inverted else style.TEXT_DARK
    if height <= 1.2:
        _add_editorial_text(slide, title, left + 0.32, top + 0.2, 2.0, height - 0.4, style, 11, title_color, True)
        _add_editorial_text(slide, body, left + 2.45, top + 0.16, width - 2.77, height - 0.32, style, 15, body_color, True)
    else:
        _add_editorial_text(slide, title, left + 0.32, top + 0.25, width - 0.64, 0.35, style, 12, title_color, True)
        _add_editorial_text(
            slide, body, left + 0.32, top + 0.72, width - 0.64,
            height - 0.87, style, body_font, body_color, True,
        )
    return block


def _render_slide_creative(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render creative slides with asymmetric editorial compositions."""
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = slide_data.get("narrative_role") or "opportunity_map"
    title = slide_data.get("title", f"第 {idx} 页")
    items = _commercial_slide_items(slide_data)
    add_page_number(slide, prs, idx + 1, total_slides, style)
    marker = slide.shapes.add_shape(1, Inches(0.72), Inches(0.48), Inches(0.1), Inches(0.72))
    marker.fill.solid(); marker.fill.fore_color.rgb = style.ACCENT_COLOR; marker.line.fill.background()
    _add_editorial_text(slide, title, 1.05, 0.42, 9.5, 0.72, style, 29, style.PRIMARY_COLOR, True)
    _add_editorial_text(slide, f"{idx:02d}", 11.25, 0.22, 1.2, 0.9, style, 38, style.ACCENT_LIGHT, True, PP_ALIGN.RIGHT)

    if role == "evidence_story":
        _add_editorial_text(slide, "事实先行，判断随后", 0.8, 1.42, 8.2, 0.6, style, 24, style.PRIMARY_COLOR, True)
        band_specs = [(0.8, 2.35, 11.8, style.PRIMARY_COLOR), (1.55, 3.65, 11.05, style.ACCENT_COLOR), (2.3, 4.95, 10.3, style.PRIMARY_DARK)]
        for number, (left, top, width, color) in enumerate(band_specs, 1):
            _add_creative_block(slide, f"EVIDENCE / {number:02d}", items[number - 1] if len(items) >= number else items[-1], left, top, width, 1.05, style, color)
    elif role == "strategic_choice":
        _add_creative_block(slide, "PATH / A", items[0], 0.7, 1.6, 5.45, 3.15, style, style.PRIMARY_COLOR)
        _add_creative_block(slide, "PATH / B", items[1] if len(items) > 1 else items[0], 7.18, 1.6, 5.45, 3.15, style, style.ACCENT_COLOR)
        versus = slide.shapes.add_shape(9, Inches(6.0), Inches(2.68), Inches(1.0), Inches(1.0))
        versus.fill.solid(); versus.fill.fore_color.rgb = style.BG_WHITE; versus.line.color.rgb = style.PRIMARY_COLOR
        _add_editorial_text(slide, "VS", 6.0, 2.96, 1.0, 0.3, style, 14, style.PRIMARY_COLOR, True, PP_ALIGN.CENTER)
        _add_creative_block(slide, "CHOICE / 推荐", "；".join(items[2:]) or "优先验证价值更清晰的路径", 2.05, 5.18, 9.25, 1.12, style, style.PRIMARY_DARK)
    elif role == "execution_roadmap":
        stage_specs = [(0.75, 4.35, style.PRIMARY_COLOR), (4.88, 3.12, style.ACCENT_COLOR), (9.0, 1.9, style.PRIMARY_DARK)]
        connector_specs = [(4.3, 5.1, 4.88, 4.15), (8.43, 3.9, 9.0, 2.95)]
        for start_x, start_y, end_x, end_y in connector_specs:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(start_x), Inches(start_y), Inches(end_x), Inches(end_y),
            )
            connector.line.color.rgb = style.ACCENT_COLOR
            connector.line.width = Pt(2.5)
        for number, (left, top, color) in enumerate(stage_specs, 1):
            stage_name = ("试点", "扩展", "规模化")[number - 1]
            _add_creative_block(
                slide, f"0{number} / {stage_name}",
                items[number - 1] if len(items) >= number else items[-1],
                left, top, 3.55, 1.95, style, color, body_font=12,
            )
        _add_editorial_text(slide, "从左下到右上，逐级放大确定性", 0.78, 1.55, 6.8, 0.45, style, 16, style.TEXT_GRAY)
    elif role == "decision_close":
        _add_creative_block(slide, "PRIORITY / 01", items[0], 0.72, 1.55, 6.75, 4.75, style, style.PRIMARY_COLOR)
        _add_creative_block(slide, "ACTION / 02", items[1] if len(items) > 1 else items[0], 7.85, 1.55, 4.78, 2.1, style, style.ACCENT_COLOR)
        _add_creative_block(slide, "ACTION / 03", items[2] if len(items) > 2 else items[-1], 7.85, 4.18, 4.78, 2.12, style, style.PRIMARY_DARK)
        if len(items) > 3:
            _add_editorial_text(slide, items[3], 0.9, 5.7, 6.35, 0.42, style, 11, style.TEXT_WHITE, True)
    else:
        _add_creative_block(slide, "SIGNAL / 01", items[0], 0.72, 1.55, 5.35, 4.75, style, style.PRIMARY_COLOR)
        _add_creative_block(slide, "SIGNAL / 02", items[1] if len(items) > 1 else items[0], 6.42, 1.55, 6.2, 2.05, style, style.ACCENT_COLOR)
        _add_creative_block(slide, "SIGNAL / 03", items[2] if len(items) > 2 else items[-1], 7.15, 4.02, 5.47, 2.28, style, style.PRIMARY_DARK)

    sources = slide_data.get("evidence_sources", [])
    if sources:
        _add_editorial_text(slide, f"SOURCE / {sources[0].get('title', '公开资料')[:64]}", 0.72, 6.86, 10.2, 0.22, style, 8, style.TEXT_GRAY)
    return slide


def _render_slide_default(prs, blank_layout, style, slide_data, idx, total_slides):
    """Render semantic content with varied editorial compositions."""
    theme_renderers = {
        "academic": _render_slide_academic,
        "creative": _render_slide_creative,
        "education": _render_slide_education,
        "elegant": _render_slide_elegant,
        "medical": _render_slide_medical,
        "modern": _render_slide_modern,
        "minimal": _render_slide_minimal,
        "tech": _render_slide_tech,
    }
    theme_renderer = theme_renderers.get(style.template_name)
    if theme_renderer:
        return theme_renderer(prs, blank_layout, style, slide_data, idx, total_slides)
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    role = _slide_role(slide_data)
    title = slide_data.get("title", f"第 {idx} 页")
    items = _commercial_slide_items(slide_data)
    add_page_number(slide, prs, idx + 1, total_slides, style)
    _add_editorial_text(slide, f"0{idx}", 0.6, 0.45, 0.55, 0.35, style, 12, style.ACCENT_COLOR, True)
    _add_editorial_text(slide, title, 1.25, 0.32, 10.8, 0.65, style, 27, style.PRIMARY_COLOR, True)
    sources = slide_data.get("evidence_sources", [])

    if role == "evidence_story":
        panel = slide.shapes.add_shape(5, Inches(0.65), Inches(1.45), Inches(2.75), Inches(4.95))
        panel.fill.solid(); panel.fill.fore_color.rgb = style.PRIMARY_COLOR; panel.line.fill.background()
        _add_editorial_text(slide, "EVIDENCE", 0.95, 1.85, 2.05, 0.3, style, 12, style.TEXT_WHITE, True)
        _add_editorial_text(slide, "01", 0.9, 2.4, 2.1, 0.9, style, 46, style.TEXT_WHITE, True)
        _add_editorial_text(slide, "从事实出发\n把变化转成判断", 0.95, 3.75, 1.95, 0.9, style, 17, style.TEXT_WHITE, True)
        _add_editorial_text(slide, "可验证  ·  可复核  ·  可追溯", 0.95, 5.45, 2.05, 0.3, style, 10, style.TEXT_WHITE, True)
        for number, item in enumerate(items[:3], 1):
            _add_editorial_card(slide, f"证据 {number:02d}", item, 3.8, 1.45 + (number - 1) * 1.55, 8.85, 1.22, style, style.ACCENT_COLOR)
        if len(items) > 3:
            _add_editorial_card(slide, "结论含义", items[3], 3.8, 6.0, 8.85, 0.75, style, style.PRIMARY_COLOR)
    elif role == "strategic_choice":
        _add_editorial_text(slide, "两条路径，一个明确选择", 0.65, 1.25, 8.5, 0.35, style, 15, style.TEXT_GRAY)
        _add_editorial_card(slide, "方案 A", items[0], 0.65, 1.9, 5.8, 2.45, style, style.PRIMARY_LIGHT)
        _add_editorial_card(slide, "方案 B", items[1] if len(items) > 1 else items[0], 6.85, 1.9, 5.8, 2.45, style, style.ACCENT_COLOR)
        _add_editorial_card(slide, "推荐路径", "；".join(items[2:]) or "优先验证价值更清晰、交付更可控的路径。", 0.65, 4.75, 12.0, 1.45, style, style.ACCENT_COLOR, filled=True)
    elif role == "execution_roadmap":
        _add_editorial_text(slide, "三个阶段，逐步验证，形成可复制交付", 0.65, 1.25, 11.7, 0.4, style, 15, style.TEXT_GRAY)
        stages = items[:3]
        for number, item in enumerate(stages, 1):
            left = 0.7 + (number - 1) * 4.15
            circle = slide.shapes.add_shape(9, Inches(left), Inches(2.0), Inches(0.62), Inches(0.62))
            circle.fill.solid(); circle.fill.fore_color.rgb = style.ACCENT_COLOR; circle.line.fill.background()
            _add_editorial_text(slide, str(number), left, 2.12, 0.62, 0.25, style, 14, style.TEXT_WHITE, True, PP_ALIGN.CENTER)
            if number < 3:
                connector = slide.shapes.add_shape(1, Inches(left + 0.62), Inches(2.29), Inches(3.53), Inches(0.04))
                connector.fill.solid(); connector.fill.fore_color.rgb = style.PRIMARY_LIGHT; connector.line.fill.background()
            stem = slide.shapes.add_shape(1, Inches(left + 0.29), Inches(2.62), Inches(0.04), Inches(0.33))
            stem.fill.solid(); stem.fill.fore_color.rgb = style.PRIMARY_LIGHT; stem.line.fill.background()
            _add_editorial_card(slide, ["试点", "扩展", "规模化"][number - 1], item, left - 0.1, 2.95, 3.65, 1.7, style, style.ACCENT_COLOR)
        _add_editorial_card(slide, "进入下一阶段的门槛", items[3] if len(items) > 3 else "核心指标达到目标，并完成复盘。", 2.1, 5.05, 9.15, 1.1, style, style.PRIMARY_COLOR)
    elif role == "decision_close":
        _add_editorial_text(slide, "把决定转成今天可执行的三个动作", 1.1, 1.45, 11.1, 0.65, style, 25, style.PRIMARY_COLOR, True, PP_ALIGN.CENTER)
        for number, item in enumerate(items[:3], 1):
            _add_editorial_card(slide, f"决策 {number:02d}", item, 0.75 + (number - 1) * 4.15, 2.65, 3.7, 2.05, style, style.ACCENT_COLOR)
        close_metric = items[3] if len(items) > 3 else "下一次评审以结果改善幅度决定是否扩大投入"
        _add_editorial_text(slide, close_metric, 2.0, 5.35, 9.4, 0.75, style, 13, style.TEXT_GRAY, False, PP_ALIGN.CENTER)
    else:
        _add_editorial_text(slide, "从市场变化中锁定最值得优先验证的机会", 0.65, 1.25, 11.7, 0.4, style, 15, style.TEXT_GRAY)
        local_images = slide_data.get("local_images", [])
        usable_image = local_images and Path(local_images[0]).exists()
        if usable_image:
            _add_cropped_picture(slide, local_images[0], 0.7, 1.95, 4.15, 3.75)
            for number, item in enumerate(items[:3], 1):
                _add_editorial_card(slide, f"机会 {number:02d}", item, 5.2, 1.95 + (number - 1) * 1.3, 7.45, 1.08, style, style.ACCENT_COLOR if number == 2 else style.PRIMARY_COLOR)
        else:
            for number, item in enumerate(items[:3], 1):
                _add_editorial_card(slide, f"机会 {number:02d}", item, 0.7 + (number - 1) * 4.15, 2.05, 3.65, 2.55, style, style.ACCENT_COLOR if number == 2 else style.PRIMARY_COLOR)
        _add_editorial_card(slide, "验证指标", items[3] if len(items) > 3 else "在明确周期内验证使用率与结果改善幅度。", 2.15, 4.85, 9.05, 1.05, style, style.PRIMARY_COLOR)

    if sources:
        source_title = sources[0].get("title", "网络检索来源")
        _add_editorial_text(slide, f"来源：{source_title[:72]}", 0.65, 6.88, 10.0, 0.22, style, 9, style.TEXT_GRAY)


# 导入 Aicode 的会话压缩函数
from app.api.v1.Aicode import compress_conversation_history
from app.utils import call_llm

# =============================================================================
# 辅助函数：大纲生成、多格式导出、预览
# =============================================================================

async def generate_ppt_outline(req: PPTGenerationRequest, user_id: str = None) -> Dict[str, Any]:
    """使用 AI 生成 PPT 大纲 (支持 Skills 和多参数)"""
    # 构建 prompt
    skill_prompts = []
    for skill in req.skills:
        skill_prompts.append(f"- 遵循 {skill} 最佳实践")
    
    skill_context = "\n".join(skill_prompts) if skill_prompts else ""
    
    # 构建内容上下文
    content_context = req.topic
    if req.material_file_ids:
        content_context += f"\n\n参考素材 ID: {req.material_file_ids}"
    
    role_schema = " | ".join(NARRATIVE_ROLES)
    prompt = f"""
请根据以下要求生成一个专业的 PPT 大纲，返回 JSON 格式：

{{
    "title": "PPT 标题",
    "slides": [
        {{
            "slide_number": 1,
            "title": "幻灯片标题",
            "key_message": "本页唯一核心结论",
            "content": ["兼容正文 1", "兼容正文 2"],
            "content_blocks": [
                {{
                    "type": "signal|evidence|case|implication|option|recommendation|criteria|stage|gate|decision|action|request|success_metric",
                    "content": "简洁、可展示的论点",
                    "metadata": {{"metric": "指标名称", "target": "目标值"}}
                }}
            ],
            "slide_type": "key_points|data|comparison|timeline|summary",
            "narrative_role": "{role_schema}",
            "notes": "演讲者备注（可选）"
        }}
    ]
}}

要求：
- 主题：{content_context}
- 页数：{req.slide_count}页
- 语言：{req.language}
- 质量：{req.quality}
- 模板风格：{req.template}
{skill_context if skill_context else ""}
- 内容页按 opportunity_map、evidence_story、strategic_choice、execution_roadmap、decision_close 形成完整商业叙事；页数超过 5 时循环使用并保持相邻页面角色不同
- opportunity_map 的 metadata 提供 roi、priority、validation_period、metric、target
- evidence_story 的 metadata 提供 metric、target、validation_period，并用公开来源可验证的数据表达
- strategic_choice 的两个 option 提供 cost、timeframe、risk，recommendation 提供 rationale
- execution_roadmap 的三个 stage 提供 deliverable、metric、target、gate
- decision_close 的 decision、action、request 提供 owner、deadline、priority，success_metric 提供 metric、target
- 每页输出 4 个 content_blocks，并让 content 与 content_blocks 的 content 保持一致

请直接返回 JSON 格式，不要有多余解释。
"""
    
    try:
        response = await call_llm(
            model=req.model,
            prompt=prompt,
            api_key_token=req.api_key_token
        )

        # 从 API 响应中提取 content
        if isinstance(response, dict):
            choices = response.get('choices', [])
            if choices:
                content = choices[0].get('message', {}).get('content', '')
            else:
                content = str(response)
        else:
            content = str(response)

        # 提取 JSON（可能用 markdown 代码块包裹）
        json_match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```|(\{[\s\S]*\}|\[[\s\S]*\])', content)
        if json_match:
            json_str = json_match.group(1) or json_match.group(2)
        else:
            json_str = content

        outline = json.loads(json_str)
        return outline
        
    except Exception as e:
        logger.warning(f"AI 大纲生成失败，使用默认大纲：{e}")
        page_blueprint = build_commercial_page_blueprint(req.topic)
        slides = []
        for index in range(req.slide_count):
            page = page_blueprint[index % len(page_blueprint)]
            slides.append({
                "slide_number": index + 1,
                "title": page["title"],
                "key_message": page["key_message"],
                "content": [block["content"] for block in page["blocks"]],
                "content_blocks": page["blocks"],
                "slide_type": page["slide_type"],
                "narrative_role": page["role"],
                "asset_intent": page["asset_intent"],
                "notes": "",
            })
        return {
            "title": req.topic[:50],
            "slides": slides,
        }


def _normalize_approved_outline(outline: Dict[str, Any]) -> Dict[str, Any]:
    """Convert reviewed semantic slides into the legacy renderer contract."""
    normalized_slides = []
    for index, slide in enumerate(outline.get("slides", []), 1):
        content = slide.get("content", "")
        if not content and slide.get("content_blocks"):
            content = [
                block.get("content", "")
                for block in slide["content_blocks"]
                if block.get("content", "")
            ]
        normalized_slides.append({
            "slide_number": slide.get("slide_number", index),
            "title": slide.get("title", f"第 {index} 页"),
            "content": content,
            "content_blocks": slide.get("content_blocks", []),
            "key_message": slide.get("key_message", ""),
            "slide_type": slide.get("slide_type", "content"),
            "notes": slide.get("notes", slide.get("speaker_notes", "")),
            "asset_intent": slide.get("asset_intent"),
            "narrative_role": slide.get("narrative_role", "opportunity_map"),
            "evidence_sources": slide.get("evidence_sources", outline.get("evidence_sources", [])),
        })
    return {
        "title": outline.get("title", "PPT 标题"),
        "slides": normalized_slides,
    }

async def generate_pptx_file_enhanced(filepath: Path, outline: Dict[str, Any], req: PPTGenerationRequest, update_progress=None):
    """生成 PPTX 文件 (增强版：包含视觉决策和模板支持)"""
    outline = _normalize_approved_outline(outline)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    style = PPTStyle(template_name=req.template)
    template_manager = TemplateManager()
    token_template_id = {"modern": "minimal", "business": "business_report", "creative": "pitch_deck"}.get(req.template, req.template)
    try:
        tokens = template_manager.resolve_design_tokens(token_template_id)
    except KeyError:
        tokens = None
        logger.warning("设计令牌模板不存在，沿用 PPTStyle: %s", req.template)
    total_slides = 1 + len(outline.get('slides', []))
    blank_layout = prs.slide_layouts[6]

    if update_progress: await update_progress(message="正在创建封面...")

    # === 1. 标题页（封面）===
    slide = prs.slides.add_slide(blank_layout)
    add_slide_background(slide, prs, style, light=True)
    
    if style.template_name == "education":
        _add_editorial_text(slide, "LEARNING LAB / 2026", 0.78, 0.68, 3.8, 0.32, style, 10, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 0.78, 1.55, 9.7, 1.62, style, 42, style.TEXT_DARK, True)
        _add_editorial_text(slide, "理解关键判断 · 完成一次练习 · 带走下一步行动", 0.82, 3.5, 7.8, 0.42, style, 14, style.PRIMARY_DARK)
        for number, label in enumerate(("理解", "练习", "应用"), 1):
            left = 0.78 + (number - 1) * 4.05
            _add_learning_card(slide, f"STEP 0{number}", label, left, 5.05, 3.55, 1.35, style, number == 3)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 10.75, 0.68, 1.75, 0.28, style, 9, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    elif style.template_name == "elegant":
        _add_editorial_text(slide, "PRIVATE EDITION / 2026", 0.8, 0.7, 3.8, 0.3, style, 9, style.PRIMARY_COLOR, True)
        cover_rule = slide.shapes.add_shape(1, Inches(0.8), Inches(1.28), Inches(11.75), Inches(0.025))
        cover_rule.fill.solid(); cover_rule.fill.fore_color.rgb = style.PRIMARY_COLOR; cover_rule.line.fill.background()
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 0.8, 1.85, 9.5, 1.62, style, 41, style.TEXT_DARK, True)
        _add_editorial_text(slide, "INSIGHT  /  CHOICE  /  COMMITMENT", 0.83, 4.08, 6.5, 0.34, style, 11, style.PRIMARY_DARK, True)
        _add_editorial_text(slide, "I", 10.75, 4.6, 1.55, 1.25, style, 66, style.PRIMARY_COLOR, True, PP_ALIGN.RIGHT)
        bottom_rule = slide.shapes.add_shape(1, Inches(0.8), Inches(6.2), Inches(11.75), Inches(0.025))
        bottom_rule.fill.solid(); bottom_rule.fill.fore_color.rgb = style.PRIMARY_LIGHT; bottom_rule.line.fill.background()
        _add_editorial_text(slide, "EXECUTIVE MEMORANDUM", 0.8, 6.42, 3.7, 0.28, style, 9, style.TEXT_GRAY, True)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 10.75, 6.42, 1.8, 0.28, style, 9, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    elif style.template_name == "medical":
        _add_editorial_text(slide, "CLINICAL BRIEF / 2026", 0.78, 0.68, 3.8, 0.32, style, 10, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 0.78, 1.55, 9.25, 1.5, style, 40, style.TEXT_DARK, True)
        _add_editorial_text(slide, "证据驱动判断 · 路径化推进照护", 0.82, 3.55, 7.8, 0.42, style, 14, style.PRIMARY_DARK)
        cover_rule = slide.shapes.add_shape(1, Inches(0.78), Inches(4.55), Inches(11.78), Inches(0.035))
        cover_rule.fill.solid(); cover_rule.fill.fore_color.rgb = style.PRIMARY_LIGHT; cover_rule.line.fill.background()
        for number, label in enumerate(("SIGNAL", "PATHWAY", "ACTION"), 1):
            left = 0.78 + (number - 1) * 4.05
            _add_clinical_card(slide, f"0{number}", label, left, 5.18, 3.55, 1.28, style, number == 1)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 10.75, 0.68, 1.75, 0.28, style, 9, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    elif style.template_name == "academic":
        margin = slide.shapes.add_shape(1, Inches(0.72), Inches(0.68), Inches(0.09), Inches(6.15))
        margin.fill.solid(); margin.fill.fore_color.rgb = style.PRIMARY_COLOR; margin.line.fill.background()
        paper = slide.shapes.add_shape(1, Inches(1.18), Inches(0.68), Inches(11.42), Inches(6.15))
        paper.fill.solid(); paper.fill.fore_color.rgb = style.TEXT_WHITE; paper.line.color.rgb = style.PRIMARY_LIGHT
        _add_editorial_text(slide, "RESEARCH BRIEF / 2026", 1.65, 1.15, 4.2, 0.32, style, 10, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 1.65, 2.02, 9.75, 1.7, style, 39, style.PRIMARY_DARK, True)
        _add_editorial_text(slide, "证据 · 方法 · 结论", 1.68, 4.2, 5.3, 0.42, style, 14, style.TEXT_GRAY)
        cover_rule = slide.shapes.add_shape(1, Inches(1.65), Inches(5.35), Inches(9.95), Inches(0.025))
        cover_rule.fill.solid(); cover_rule.fill.fore_color.rgb = style.PRIMARY_LIGHT; cover_rule.line.fill.background()
        _add_editorial_text(slide, "VOLUME 01", 1.65, 5.62, 2.0, 0.28, style, 9, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 9.8, 5.62, 1.8, 0.28, style, 9, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    elif style.template_name == "creative":
        spine = slide.shapes.add_shape(1, Inches(0.75), Inches(0.78), Inches(0.12), Inches(5.95))
        spine.fill.solid(); spine.fill.fore_color.rgb = style.ACCENT_COLOR; spine.line.fill.background()
        _add_editorial_text(slide, "STRATEGY / 2026", 1.3, 1.25, 4.0, 0.35, style, 12, style.ACCENT_COLOR, True)
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 1.3, 2.05, 8.2, 1.65, style, 41, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, "MAKE THE NEXT MOVE VISIBLE", 1.32, 4.2, 5.7, 0.38, style, 13, style.TEXT_GRAY, True)
        _add_editorial_text(slide, "IDEA", 9.15, 0.65, 3.25, 1.2, style, 48, style.ACCENT_LIGHT, True, PP_ALIGN.RIGHT)
        footer_panel = slide.shapes.add_shape(1, Inches(8.1), Inches(4.75), Inches(5.233), Inches(2.75))
        footer_panel.fill.solid(); footer_panel.fill.fore_color.rgb = style.PRIMARY_COLOR; footer_panel.line.fill.background()
        _add_editorial_text(slide, "01  洞察\n02  选择\n03  行动", 8.65, 5.12, 3.0, 1.35, style, 17, style.TEXT_WHITE, True)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 11.15, 6.6, 1.45, 0.3, style, 10, style.ACCENT_LIGHT, True, PP_ALIGN.RIGHT)
    elif style.template_name == "modern":
        _add_editorial_text(slide, "STRATEGY / BRIEF", 0.78, 0.72, 3.5, 0.35, style, 11, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 0.78, 1.55, 9.4, 1.55, style, 42, style.TEXT_DARK, True)
        _add_editorial_text(slide, "从洞察到行动，建立可验证的增长路径", 0.82, 3.48, 7.2, 0.45, style, 15, style.TEXT_GRAY)
        for number, label in enumerate(("洞察", "选择", "行动"), 1):
            left = 0.78 + (number - 1) * 4.15
            _add_modern_tile(slide, f"0{number}", label, left, 5.15, 3.65, 1.25, style, number == 1)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 11.0, 0.72, 1.55, 0.3, style, 10, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
    elif style.template_name == "minimal":
        _add_editorial_text(slide, "PRESENTATION", 0.72, 0.68, 2.2, 0.3, style, 10, style.TEXT_GRAY, True)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 10.85, 0.68, 1.75, 0.3, style, 10, style.TEXT_GRAY, False, PP_ALIGN.RIGHT)
        _add_minimal_rule(slide, 0.72, 1.15, 11.88, style, 2)
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 0.72, 2.0, 10.45, 1.8, style, 44, style.TEXT_DARK, True)
        _add_editorial_text(slide, "01", 10.75, 4.62, 1.85, 1.0, style, 52, style.TEXT_DARK, True, PP_ALIGN.RIGHT)
        _add_minimal_rule(slide, 0.72, 6.1, 3.25, style)
        _add_editorial_text(slide, "INSIGHT / CHOICE / ACTION", 0.72, 6.3, 4.5, 0.3, style, 10, style.TEXT_GRAY, True)
    elif style.template_name == "tech":
        _add_editorial_text(slide, "SYSTEM / STRATEGY DECK", 0.78, 0.62, 4.2, 0.32, style, 10, style.ACCENT_LIGHT, True)
        _add_editorial_text(slide, "ONLINE", 10.7, 0.62, 1.8, 0.32, style, 10, style.ACCENT_LIGHT, True, PP_ALIGN.RIGHT)
        top_line = slide.shapes.add_shape(1, Inches(0.78), Inches(1.12), Inches(11.72), Inches(0.04))
        top_line.fill.solid(); top_line.fill.fore_color.rgb = style.ACCENT_COLOR; top_line.line.fill.background()
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 0.78, 1.72, 10.2, 1.55, style, 42, style.TEXT_WHITE, True)
        _add_editorial_text(slide, "将复杂信号转化为可执行决策", 0.82, 3.65, 6.3, 0.42, style, 14, style.ACCENT_LIGHT)
        for number, label in enumerate(("INSIGHT", "CHOICE", "ACTION"), 1):
            _add_tech_panel(slide, f"MODULE / 0{number}", label, 0.78 + (number - 1) * 4.05, 5.05, 3.55, 1.18, style, number == 3)
        _add_editorial_text(slide, datetime.now().strftime('%Y.%m.%d'), 10.75, 6.72, 1.75, 0.28, style, 9, style.ACCENT_LIGHT, False, PP_ALIGN.RIGHT)
    else:
        hero = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(9.15), prs.slide_height)
        hero.fill.solid(); hero.fill.fore_color.rgb = style.PRIMARY_COLOR; hero.line.fill.background()
        accent = slide.shapes.add_shape(1, Inches(0.85), Inches(1.15), Inches(1.0), Inches(0.08))
        accent.fill.solid(); accent.fill.fore_color.rgb = style.ACCENT_COLOR; accent.line.fill.background()
        _add_editorial_text(slide, "STRATEGY BRIEF", 0.85, 1.4, 4.0, 0.35, style, 12, style.TEXT_WHITE, True)
        _add_editorial_text(slide, outline.get('title', 'PPT 标题'), 0.85, 2.15, 7.55, 1.55, style, 40, style.TEXT_WHITE, True)
        _add_editorial_text(slide, f"商业决策简报  ·  {datetime.now().strftime('%Y.%m.%d')}", 0.88, 5.65, 6.5, 0.4, style, 14, style.TEXT_WHITE)
        _add_editorial_text(slide, "01", 10.1, 1.35, 1.8, 0.8, style, 34, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, "洞察", 10.1, 2.1, 1.8, 0.35, style, 14, style.TEXT_GRAY, True)
        _add_editorial_text(slide, "02", 10.1, 3.15, 1.8, 0.8, style, 34, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, "选择", 10.1, 3.9, 1.8, 0.35, style, 14, style.TEXT_GRAY, True)
        _add_editorial_text(slide, "03", 10.1, 4.95, 1.8, 0.8, style, 34, style.PRIMARY_COLOR, True)
        _add_editorial_text(slide, "行动", 10.1, 5.7, 1.8, 0.35, style, 14, style.TEXT_GRAY, True)
    
    # === 2. 视觉决策阶段 ===
    visual_plan = None
    if update_progress: await update_progress(message="正在分析视觉需求...")
    try:
        visual_plan = await visual_analyzer.analyze_ppt_content(
            title=outline.get("title", "PPT"),
            slides_content=outline.get("slides", []),
            theme=req.template,
            api_key_token=req.api_key_token,
        )
    except Exception as e:
        logger.warning(f"视觉分析失败，使用默认布局: {e}")

    # === 3. 内容幻灯片 ===
    for idx, slide_data in enumerate(outline.get('slides', []), 1):
        if update_progress: 
            await update_progress(progress=int(50 + 50 * idx / total_slides), message=f"正在生成第 {idx} 页...")
        
        # 如果视觉决策启用，尝试使用
        image_asset = None
        role = slide_data.get("narrative_role", "opportunity_map")
        needs_editorial_image = role == "opportunity_map"

        if needs_editorial_image and slide_data.get("asset_intent"):
            try:
                intent = slide_data["asset_intent"]
                image_asset = await image_manager.get_image_for_slide(
                    image_type=intent.get("asset_type", "illustration"),
                    description=intent.get("description", slide_data.get("title", "PPT 视觉素材")),
                    keywords=intent.get("keywords", []),
                    slide_index=idx,
                    style=f"{req.template} 风格",
                )
            except Exception as e:
                logger.warning("页面配图获取失败: %s", e)

        if needs_editorial_image and visual_plan and idx <= len(visual_plan.slides) and not image_asset:
            slide_decision = visual_plan.slides[idx - 1]
            main_image = slide_decision.get_main_image() if slide_decision.need_image else None

            if main_image and main_image.image_type != ImageType.NONE:
                try:
                    image_asset = await image_manager.get_image_for_slide(
                        image_type=main_image.image_type.value,
                        description=main_image.description,
                        keywords=main_image.keywords,
                        slide_index=idx,
                        style=f"{req.template} 风格"
                    )
                except Exception: pass

        if image_asset and image_asset.local_path:
            slide_data = {**slide_data, "local_images": [image_asset.local_path]}
        _render_slide_default(prs, blank_layout, style, slide_data, idx, total_slides)
    
    prs.save(str(filepath))
    logger.info(f"PPTX 文件保存成功 | file: {filepath}")
    
    # 保存 JSON 数据（用于预览）
    json_path = filepath.parent / f"{filepath.stem}_slides.json"
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(outline.get('slides', []), f, ensure_ascii=False, indent=2)

async def generate_html_ppt(filepath: Path, outline: Dict[str, Any], req: PPTGenerationRequest):
    """生成 HTML 格式 PPT"""
    template = PPT_TEMPLATES.get(req.template, PPT_TEMPLATES["modern"])
    
    title = html.escape(str(outline.get('title', 'PPT')))
    language = html.escape(req.language, quote=True)
    html_content = f"""<!DOCTYPE html>
<html lang="{language}">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        :root {{
            --primary-color: {template['primary_color']};
            --secondary-color: {template['secondary_color']};
            --font-family: {template['font_family']};
            --background: {template['background']};
        }}
        body {{
            font-family: var(--font-family);
            background: var(--background);
            margin: 0;
            padding: 20px;
        }}
        .slide {{
            background: white;
            border-radius: 8px;
            padding: 40px;
            margin: 20px 0;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            page-break-after: always;
        }}
        .slide h1 {{ color: var(--primary-color); font-size: 32px; }}
        .slide h2 {{ color: var(--secondary-color); font-size: 24px; }}
        .slide-content {{ font-size: 18px; line-height: 1.6; white-space: pre-line; }}
        .slide-number {{ color: var(--secondary-color); font-size: 14px; text-align: right; margin-top: 20px; }}
    </style>
</head>
<body>
    <div class="slide">
        <h1>{title}</h1>
        <p style="color: var(--secondary-color);">生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}</p>
    </div>
"""
    
    for slide_data in outline.get('slides', []):
        if slide_data.get('slide_type') == 'cover':
            continue
            
        slide_title = html.escape(str(slide_data.get('title', '幻灯片标题')))
        slide_content = html.escape(str(slide_data.get('content', ''))).replace('\n', '<br>')
        slide_number = html.escape(str(slide_data.get('slide_number', 0)))
        html_content += f"""
    <div class="slide">
        <h2>{slide_title}</h2>
        <div class="slide-content">{slide_content}</div>
        <div class="slide-number">{slide_number}</div>
    </div>
"""
    
    html_content += """</body></html>"""
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    logger.info(f"HTML PPT 保存成功 | file: {filepath}")

async def generate_markdown_ppt(filepath: Path, outline: Dict[str, Any], req: PPTGenerationRequest):
    """生成 Markdown 格式 PPT"""
    md_content = f"# {outline.get('title', 'PPT 标题')}\n\n"
    md_content += f"*生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}*\n\n---\n\n"
    
    for slide_data in outline.get('slides', []):
        md_content += f"## {slide_data.get('title', '幻灯片标题')}\n\n"
        md_content += f"{slide_data.get('content', '')}\n\n"
        
        if slide_data.get('notes'):
            md_content += f"> **备注**: {slide_data.get('notes')}\n\n"
        
        md_content += "---\n\n"
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(md_content)
    
    logger.info(f"Markdown PPT 保存成功 | file: {filepath}")

def generate_preview_html(ppt_id: str) -> str:
    """生成预览 HTML 页面"""
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PPT 预览 - {html.escape(ppt_id)}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; background: #f1f5f9; padding: 20px; }}
        .container {{ max-width: 1200px; margin: 0 auto; }}
        .header {{ background: white; padding: 20px 30px; border-radius: 8px; margin-bottom: 20px; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
        .header h1 {{ font-size: 24px; color: #1e293b; }}
        .slides-container {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(500px, 1fr)); gap: 20px; }}
        .slide-card {{ background: white; border-radius: 8px; overflow: hidden; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
        .slide-header {{ background: #2563eb; color: white; padding: 12px 20px; display: flex; justify-content: space-between; align-items: center; }}
        .slide-number {{ font-size: 14px; font-weight: 600; }}
        .slide-type {{ font-size: 12px; padding: 4px 8px; background: rgba(255,255,255,0.2); border-radius: 4px; }}
        .slide-body {{ padding: 20px; }}
        .slide-title {{ font-size: 18px; font-weight: 600; color: #1e293b; margin-bottom: 12px; }}
        .slide-content {{ font-size: 14px; color: #64748b; line-height: 1.6; white-space: pre-line; }}
        .download-section {{ background: white; padding: 20px 30px; border-radius: 8px; margin-top: 20px; text-align: center; }}
        .download-btn {{ display: inline-block; padding: 12px 24px; background: #2563eb; color: white; text-decoration: none; border-radius: 6px; font-weight: 600; margin: 0 10px; }}
        .download-btn:hover {{ background: #1d4ed8; }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>PPT 预览</h1>
            <p style="color: #64748b; margin-top: 8px;">PPT ID: {html.escape(ppt_id)}</p>
        </div>
        
        <div class="download-section">
            <h2 style="margin-bottom: 16px;">下载 PPT</h2>
            <a href="/api/v1/pptx/download/{ppt_id}?format=pptx" class="download-btn">下载 PowerPoint</a>
        </div>
        
        <div id="slides-container" class="slides-container">
            <!-- Slides will be loaded here -->
        </div>
    </div>
    
    <script>
        fetch('/api/v1/pptx/{html.escape(ppt_id)}' + '/slides')
            .then(r => r.json())
            .then(data => {{
                const container = document.getElementById('slides-container');
                if (data.slides) {{
                    data.slides.forEach((slide, idx) => {{
                        const card = document.createElement('div');
                        card.className = 'slide-card';
                        card.innerHTML = `
                            <div class="slide-header">
                                <span class="slide-number"></span>
                                <span class="slide-type"></span>
                            </div>
                            <div class="slide-body">
                                <div class="slide-title"></div>
                                <div class="slide-content"></div>
                            </div>
                        `;
                        card.querySelector('.slide-number').textContent = idx + 1;
                        card.querySelector('.slide-type').textContent = slide.slide_type || 'content';
                        card.querySelector('.slide-title').textContent = slide.title || '';
                        card.querySelector('.slide-content').textContent = slide.content || '';
                        container.appendChild(card);
                    }});
                }}
            }})
            .catch(err => console.error('加载幻灯片失败:', err));
    </script>
</body>
</html>"""

# =============================================================================
# API 路由
# =============================================================================


@router.post("/pptx/generate_task", response_model=TaskResponse)
async def generate_ppt_task(
    req: PPTGenerationRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db)
):
    """
    使用任务队列生成 PPT（统一增强版）
    
    功能：
    - 支持多种输出格式 (PPTX, HTML, Markdown)
    - 支持模板系统与 Skills
    - 支持会话历史与素材文件绑定
    - PPTX 格式支持高级视觉决策与布局
    """
    user_id = token.get("sub", "anonymous")
    conversation_id = req.conversation_id

    if os.getenv("PPT_USE_CELERY", "false").lower() in {"1", "true", "yes"}:
        from app.services.ppt_dispatch_service import dispatch_ppt_to_celery

        task_id, celery_task_id = await dispatch_ppt_to_celery(
            db,
            int(user_id),
            req.model_dump(mode="json"),
        )
        task_options = req.options or {}
        if task_options.get("outline_id"):
            task_record = await db.scalar(
                select(Task).where(Task.task_id == task_id, Task.user_id == int(user_id))
            )
            if task_record:
                task_record.outline_id = str(task_options["outline_id"])
                task_record.outline_version = int(task_options["outline_version"])
                task_record.quality_mode = task_options.get("quality_mode", "standard")
                await db.commit()
        _register_ppt_owner(task_id, str(user_id))
        return TaskResponse(
            task_id=task_id,
            celery_task_id=celery_task_id,
            task_type="ppt_generation",
            status="pending",
            progress=0,
            progress_message="等待 Celery worker...",
            created_at=datetime.now().isoformat(),
        )
    
    # 解析 material_file_ids (兼容旧版 prompt 格式)
    material_file_ids = req.material_file_ids or []
    if not material_file_ids:
        try:
            match = re.search(r'<material>\[(.*?)\]</material>', req.topic)
            if match:
                material_file_ids = [int(x.strip()) for x in match.group(1).split(',')]
        except Exception: pass

    ppt_id = str(uuid.uuid4())
    
    async def run_ppt_generation(task_id: str, **kwargs):
        _register_ppt_owner(task_id, user_id)
        async def update_progress(progress: int = 0, message: str = "", status: str = None, result_data: str = None, **_kwargs):
            await task_manager.update_progress(
                task_id,
                progress,
                message,
                status=status,
                result_data=result_data,
                error_message=_kwargs.get("error_message"),
            )

        try:
            await update_progress(progress=5, message="正在准备上下文...")
            
            # 1. 构建完整 Prompt
            full_prompt = req.topic
            context_parts = []
            
            # 会话历史
            if conversation_id:
                try:
                    history_context = await compress_conversation_history(
                        db, int(user_id), int(conversation_id), max_messages=5
                    )
                    if history_context:
                        context_parts.append(f"\n[对话历史]\n{history_context}")
                except Exception as e:
                    logger.warning(f"获取会话历史失败: {e}")

            # 素材文件
            if material_file_ids:
                material_info = []
                for file_id in material_file_ids:
                    try:
                        result = await db.execute(select(File).where(File.id == file_id, File.user_id == int(user_id)))
                        file_record = result.scalar_one_or_none()
                        if file_record:
                            if not req.session_id or file_record.conversation_id == req.session_id:
                                material_info.append(f"- {file_record.filename}")
                    except Exception: pass
                if material_info:
                    context_parts.append(f"\n[参考素材]\n" + "\n".join(material_info))
            
            if context_parts:
                req.topic = f"{full_prompt}\n\n{''.join(context_parts)}"

            # 2. 生成大纲
            await update_progress(progress=20, message="正在生成 PPT 大纲...")
            approved_outline = (req.options or {}).get("approved_outline")
            outline = (
                _normalize_approved_outline(approved_outline)
                if approved_outline
                else await generate_ppt_outline(req, user_id=user_id)
            )

            # 立即保存大纲快照（用于恢复/增量生成）
            output_dir = PPT_OUTPUT_DIR
            output_dir.mkdir(exist_ok=True)
            snapshot_path = output_dir / f"{task_id}_slides.json"
            with open(snapshot_path, 'w', encoding='utf-8') as f:
                json.dump(outline.get('slides', []), f, ensure_ascii=False, indent=2)
            logger.info(f"保存大纲快照 | task_id={task_id} | slides={len(outline.get('slides', []))}")
            trace = build_ppt_trace_context(
                req.options,
                req.template,
                (req.options or {}).get("quality_mode", "standard"),
            )
            await save_ppt_stage_checkpoint(
                db,
                task_id,
                int(user_id),
                1,
                "planning",
                trace,
                {"slides": outline.get("slides", [])},
            )
            await db.commit()
            
            # 3. 根据格式生成文件
            output_dir = PPT_OUTPUT_DIR
            output_dir.mkdir(exist_ok=True)
            
            ext_map = {
                OutputFormat.PPTX: "pptx",
                OutputFormat.HTML: "html",
                OutputFormat.MARKDOWN: "md",
                OutputFormat.PDF: "pdf"
            }
            ext = ext_map.get(req.output_format, "pptx")

            filepath = output_dir / f"{task_id}.{ext}"
            slides_data = outline.get('slides', [])
            quality_mode = (req.options or {}).get("quality_mode", "standard")
            quality_slides = [
                {
                    **slide,
                    "id": slide.get("id", f"slide-{index + 1}"),
                    "elements": slide.get("elements", []),
                    "render_metadata": build_render_metadata(
                        slide,
                        token_version=str(trace["template_version"]),
                    ),
                }
                for index, slide in enumerate(slides_data)
            ]
            quality_slides, quality_report = await run_quality_pipeline(quality_slides, quality_mode)
            outline = {**outline, "slides": quality_slides}
            slides_data = quality_slides
            await save_ppt_stage_checkpoint(
                db,
                task_id,
                int(user_id),
                2,
                "rule_qa",
                trace,
                {"quality": serialize_quality_report(quality_report)},
            )
            await db.commit()

            if req.output_format == OutputFormat.PDF:
                pptx_filepath = output_dir / f"{task_id}.pptx"
                await generate_pptx_file_enhanced(pptx_filepath, outline, req, update_progress=update_progress)
                await _convert_pptx_to_pdf(pptx_filepath, filepath)
            elif req.output_format == OutputFormat.PPTX:
                # 使用高级视觉决策逻辑生成 PPTX
                await generate_pptx_file_enhanced(filepath, outline, req, update_progress=update_progress)
            elif req.output_format == OutputFormat.HTML:
                await update_progress(progress=60, message="正在生成 HTML 格式...")
                await generate_html_ppt(filepath, outline, req)
            elif req.output_format == OutputFormat.MARKDOWN:
                await update_progress(progress=60, message="正在生成 Markdown 格式...")
                await generate_markdown_ppt(filepath, outline, req)
            else:
                # 默认回退到 PPTX
                await generate_pptx_file_enhanced(filepath, outline, req, update_progress=update_progress)

            result = {
                "filename": filepath.name,
                "ppt_id": task_id,
                "download_url": f"/api/v1/pptx/download/{task_id}?format={ext}",
                "preview_url": _preview_url(task_id, req.output_format) if req.output_format in {OutputFormat.PPTX, OutputFormat.HTML, OutputFormat.MARKDOWN} else None
            }
            await persist_ppt_generation_result(
                db,
                task_id,
                int(user_id),
                filepath,
                "pptx" if ext == "pptx" else ext,
                result,
                trace,
                slides_data,
                quality_report,
            )
            await update_progress(
                progress=100,
                message="PPT 生成完成",
                status="completed",
                result_data=json.dumps(result)
            )
            return result
            
        except asyncio.CancelledError:
            # 取消时保存已有的中间状态
            await update_progress(
                status="cancelled",
                message="任务已取消，中间状态已保存",
                result_data=json.dumps({
                    "ppt_id": task_id,
                    "outline_saved": True,
                    "download_url": f"/api/v1/pptx/{task_id}/slides"
                })
            )
            logger.info(f"PPT 任务已取消 | task_id={task_id}")
        except Exception as e:
            await update_progress(
                status="failed",
                message=f"PPT 生成失败：{str(e)}",
                error_message=str(e)
            )
            logger.error(f"PPT 生成任务失败 | task_id: {task_id} | error: {str(e)}")
    
    task_response = await task_manager.create_task(
        task_type="ppt_generation",
        user_id=user_id,
        func=run_ppt_generation,
        params={}
    )

    # 将批准的大纲快照写入任务记录，保证生成任务可追溯。
    task_options = req.options or {}
    if task_options.get("outline_id"):
        task_record = await db.scalar(
            select(Task).where(Task.task_id == task_response, Task.user_id == int(user_id))
        )
        if task_record:
            task_record.outline_id = str(task_options["outline_id"])
            task_record.outline_version = int(task_options["outline_version"])
            task_record.quality_mode = task_options.get("quality_mode", "standard")
            await db.commit()

    logger.info(f"创建 PPT 生成任务 | task_id: {task_response} | user: {user_id}")
    return TaskResponse(
        task_id=task_response,
        task_type="ppt_generation",
        status="pending",
        progress=0,
        progress_message="等待中...",
        created_at=datetime.now().isoformat()
    )


@router.post("/pptx/generate", response_model=PPTGenerationResponse)
async def generate_ppt(
    req: PPTGenerationRequest,
    token: dict = Depends(verify_token)
):
    """
    生成 PPT（阻塞式同步接口）
    适用于页数较少或需要即时结果的场景
    """
    user_id = token.get("sub", "anonymous")
    ppt_id = str(uuid.uuid4())
    logger.info(f"PPT 同步生成请求 | user: {user_id} | topic: {req.topic[:50]}...")
    
    try:
        outline = await generate_ppt_outline(req, user_id=user_id)
        
        output_dir = PPT_OUTPUT_DIR
        output_dir.mkdir(exist_ok=True)
        
        ext_map = {OutputFormat.PPTX: "pptx", OutputFormat.HTML: "html", OutputFormat.MARKDOWN: "md"}
        ext = ext_map.get(req.output_format, "pptx")
        filepath = output_dir / f"{ppt_id}.{ext}"
        
        if req.output_format == OutputFormat.PPTX:
            await generate_pptx_file_enhanced(filepath, outline, req)
        elif req.output_format == OutputFormat.HTML:
            await generate_html_ppt(filepath, outline, req)
        elif req.output_format == OutputFormat.MARKDOWN:
            await generate_markdown_ppt(filepath, outline, req)
        else:
            await generate_pptx_file_enhanced(filepath, outline, req)

        _register_ppt_owner(ppt_id, user_id)
        return PPTGenerationResponse(
            id=ppt_id,
            status="completed",
            topic=req.topic[:50],
            slide_count=len(outline.get('slides', [])),
            output_format=req.output_format.value,
            created_at=datetime.now().isoformat(),
            download_url=f"/api/v1/pptx/download/{ppt_id}?format={ext}",
            preview_url=_preview_url(ppt_id, req.output_format) if req.output_format in {OutputFormat.PPTX, OutputFormat.HTML, OutputFormat.MARKDOWN} else None,
            slides=outline.get('slides', [])
        )
    except Exception as e:
        logger.error(f"PPT 同步生成失败 | error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/pptx/download/{ppt_id}")
async def download_ppt(
    ppt_id: str,
    format: str = Query(default="pptx", description="下载格式"),
    token: dict = Depends(verify_token)
):
    """下载 PPT 文件"""
    user_id = token.get("sub", "anonymous")
    _verify_ppt_owner(ppt_id, user_id)
    output_dir = PPT_OUTPUT_DIR
    
    requested_extension = "md" if format == "markdown" else format
    supported_extensions = ["pptx", "pdf", "html", "md"]
    if requested_extension in supported_extensions:
        possible_extensions = [requested_extension]
    else:
        raise HTTPException(status_code=400, detail="不支持的下载格式")
        
    filepath = None
    for ext in possible_extensions:
        test_path = output_dir / f"{ppt_id}.{ext}"
        if test_path.exists():
            filepath = test_path
            break
    
    if not filepath:
        raise HTTPException(status_code=404, detail="PPT 文件不存在")
    
    actual_format = filepath.suffix.lstrip('.')
    logger.info(f"下载 PPT | user: {user_id} | file: {filepath.name}")
    
    async def file_stream():
        with open(filepath, "rb") as f:
            while chunk := f.read(65536):
                yield chunk
    
    mime_types = {
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "html": "text/html",
        "md": "text/markdown",
        "pdf": "application/pdf"
    }
    
    return StreamingResponse(
        file_stream(),
        media_type=mime_types.get(actual_format, "application/octet-stream"),
        headers={
            "Content-Disposition": f'attachment; filename="{filepath.name}"'
        }
    )


@router.get("/pptx/preview/{ppt_id}", response_class=HTMLResponse)
async def preview_ppt(
    ppt_id: str,
    format: str = Query(default="pptx", description="预览格式"),
    token: dict = Depends(verify_token)
):
    """在线预览 PPT"""
    user_id = token.get("sub", "anonymous")
    _verify_ppt_owner(ppt_id, user_id)
    output_dir = PPT_OUTPUT_DIR
    if format not in {"pptx", "html", "markdown", "md"}:
        raise HTTPException(status_code=400, detail="不支持的预览格式")
    extension = "md" if format == "markdown" else format
    artifact_path = output_dir / f"{ppt_id}.{extension}"
    if not artifact_path.exists():
        raise HTTPException(status_code=404, detail="PPT 文件不存在")

    logger.info(f"预览 PPT | user: {user_id} | ppt_id: {ppt_id}")
    if extension == "html":
        return HTMLResponse(content=artifact_path.read_text(encoding="utf-8"))
    if extension == "md":
        markdown = html.escape(artifact_path.read_text(encoding="utf-8"))
        return HTMLResponse(content=f"<html><body><pre>{markdown}</pre></body></html>")
    return HTMLResponse(content=generate_preview_html(ppt_id))


@router.get("/pptx/{ppt_id}/slides")
async def get_ppt_slides(
    ppt_id: str,
    token: dict = Depends(verify_token)
):
    """获取 PPT 幻灯片数据 (JSON)"""
    _verify_ppt_owner(ppt_id, token.get("sub", "anonymous"))
    output_dir = PPT_OUTPUT_DIR
    json_path = output_dir / f"{ppt_id}_slides.json"
    
    if json_path.exists():
        with open(json_path, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
        return {"slides": slides_data}
    
    raise HTTPException(status_code=404, detail="幻灯片数据不存在")


@router.delete("/pptx/{task_id}/cancel")
async def cancel_ppt_task(
    task_id: str,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    """取消正在进行的 PPT 生成任务，并保存中间状态"""
    user_id = token.get("sub")

    task_info = await task_manager.get_task_info_async(task_id)
    if task_info:
        if task_info.get("user_id") != user_id:
            raise HTTPException(status_code=403, detail="无权取消此任务")

        if task_info.get("status") in ["success", "failed", "cancelled"]:
            return {"status": "already_finished", "message": f"任务已结束（{task_info['status']}）"}

        success = await task_manager.cancel_task(task_id)
        if success:
            logger.info(f"取消 PPT 任务 | task_id={task_id} | user={user_id}")
            return {"status": "cancelled", "message": "任务已取消"}
        return {"status": "not_running", "message": "任务未在执行"}

    try:
        task_record = await get_owned_task(db, task_id, int(user_id))
    except Exception as error:
        raise HTTPException(status_code=404, detail="任务不存在") from error

    if task_record.status in ["success", "failed", "cancelled"]:
        return {"status": "already_finished", "message": f"任务已结束（{task_record.status}）"}

    if task_record.celery_task_id:
        celery_app.control.revoke(task_record.celery_task_id, terminate=True, signal="SIGTERM")
        await transition_task(db, task_id, int(user_id), "cancelled", progress=task_record.progress or 0)
        await db.commit()
        logger.info(f"取消 Celery PPT 任务 | task_id={task_id} | user={user_id}")
        return {"status": "cancelled", "message": "任务已取消"}

    return {"status": "not_running", "message": "任务未在执行"}


@router.post("/pptx/{task_id}/update", response_model=TaskResponse)
async def update_ppt_task(
    task_id: str,
    body: PPTGenerationRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db)
):
    """基于已有的大纲/中间状态增量生成 PPT"""
    user_id = token.get("sub")
    _verify_ppt_owner(task_id, user_id)
    output_dir = PPT_OUTPUT_DIR

    # 查找之前的幻灯片数据
    json_path = output_dir / f"{task_id}_slides.json"
    if not json_path.exists():
        raise HTTPException(status_code=404, detail="找不到中间状态数据，无法增量生成")

    with open(json_path, 'r', encoding='utf-8') as f:
        existing_slides = json.load(f)

    async def run_incremental_ppt(task_id: str, **kwargs):
        async def update_progress(progress: int = 0, message: str = "", status: str = None, result_data: str = None, **_kwargs):
            await task_manager.update_progress(
                task_id,
                progress,
                message,
                status=status,
                result_data=result_data,
                error_message=_kwargs.get("error_message"),
            )

        try:
            await update_progress(progress=5, message="正在加载已有内容...")

            # 生成新的增量大纲
            new_req = body
            # 修改 prompt 要求增量生成
            incremental_prompt = f"""
已有 PPT 大纲内容：
{json.dumps(existing_slides, ensure_ascii=False, indent=2)}

请在已有内容基础上，根据以下新需求进行增量更新（新增幻灯片或修改内容）：

{new_req.topic}

要求：
- 保留已有幻灯片的核心内容
- 新增或修改内容以满足新需求
- 返回完整的合并后 JSON 大纲（包含所有幻灯片）
- 格式与已有大纲一致
"""
            original_prompt = new_req.topic
            new_req.topic = incremental_prompt

            await update_progress(progress=20, message="正在生成增量内容...")
            new_outline = await generate_ppt_outline(new_req, user_id=user_id)

            new_slides = new_outline.get('slides', existing_slides)
            new_title = new_outline.get('title', existing_slides[0].get('title', 'PPT') if existing_slides else 'PPT')

            output_id = task_id
            filepath = output_dir / f"{output_id}.{new_req.output_format.value}"
            slides_data = new_slides

            if new_req.output_format == OutputFormat.PPTX:
                await update_progress(progress=50, message="正在渲染 PPTX...")
                merged_outline = {"title": new_title, "slides": new_slides}
                await generate_pptx_file_enhanced(filepath, merged_outline, new_req, update_progress=update_progress)
            elif new_req.output_format == OutputFormat.HTML:
                await update_progress(progress=60, message="正在生成 HTML...")
                merged_outline = {"title": new_title, "slides": new_slides}
                await generate_html_ppt(filepath, merged_outline, new_req)
            elif new_req.output_format == OutputFormat.MARKDOWN:
                await update_progress(progress=60, message="正在生成 Markdown...")
                merged_outline = {"title": new_title, "slides": new_slides}
                await generate_markdown_ppt(filepath, merged_outline, new_req)
            else:
                await update_progress(progress=50, message="正在渲染 PPTX...")
                merged_outline = {"title": new_title, "slides": new_slides}
                await generate_pptx_file_enhanced(filepath, merged_outline, new_req, update_progress=update_progress)

            # 只有新文件生成成功后才更新快照，保留上一版可恢复内容。
            new_json_path = output_dir / f"{output_id}_slides.json"
            with open(new_json_path, 'w', encoding='utf-8') as f:
                json.dump(new_slides, f, ensure_ascii=False, indent=2)
            _register_ppt_owner(output_id, user_id)
            logger.info(f"更新大纲快照 | task_id={output_id} | slides={len(new_slides)}")

            ext_map = {OutputFormat.PPTX: "pptx", OutputFormat.HTML: "html", OutputFormat.MARKDOWN: "md"}
            ext = ext_map.get(new_req.output_format, "pptx")

            result = {
                "filename": filepath.name,
                "ppt_id": output_id,
                "download_url": f"/api/v1/pptx/download/{output_id}?format={ext}",
                "preview_url": _preview_url(output_id, new_req.output_format) if new_req.output_format in {OutputFormat.PPTX, OutputFormat.HTML, OutputFormat.MARKDOWN} else None
            }
            await update_progress(
                progress=100,
                message="PPT 更新完成",
                status="completed",
                result_data=json.dumps(result)
            )
            return result

        except Exception as e:
            await update_progress(
                status="failed",
                message=f"PPT 更新失败：{str(e)}",
                error_message=str(e)
            )
            logger.error(f"PPT 更新任务失败 | task_id: {task_id} | error: {str(e)}")

    task_response = await task_manager.create_task(
        task_type="ppt_update",
        user_id=user_id,
        func=run_incremental_ppt,
        params={}
    )

    logger.info(f"创建 PPT 更新任务 | task_id: {task_response} | base_task: {task_id} | user: {user_id}")
    return TaskResponse(
        task_id=task_response,
        task_type="ppt_update",
        status="pending",
        progress=0,
        progress_message="等待中...",
        created_at=datetime.now().isoformat()
    )


# =============================================================================
# 视觉增强 PPT 修改
# =============================================================================


@router.post("/pptx/{task_id}/modify")
async def modify_ppt_visual_endpoint(
    task_id: str,
    body: PPTModifyRequest,
    token: dict = Depends(verify_token),
):
    """
    视觉增强 PPT 修改

    基于已有 PPTX 文件，通过自然语言修改需求进行精确修改。
    支持修改字体、颜色、布局等属性。

    流程：
    1. 解析用户修改需求
    2. 视觉分析目标幻灯片
    3. 应用修改
    4. 返回修改结果和预览图
    """
    user_id = token.get("sub", "anonymous")
    output_dir = PPT_OUTPUT_DIR
    _verify_ppt_owner(task_id, user_id)

    # 查找已有的 PPTX 文件
    pptx_path = output_dir / f"{task_id}.pptx"
    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail=f"找不到 PPT 文件：{task_id}")

    try:
        from app.utils.pptx.visual_modifier import modify_ppt_visual

        # 生成新的 task_id 用于存储修改后的文件
        new_task_id = str(uuid.uuid4())
        new_pptx_path = output_dir / f"{new_task_id}.pptx"

        logger.info("PPT 视觉修改请求 | user: %s | base_task: %s | input: %s",
                    user_id, task_id, body.user_input[:100])

        # 执行视觉增强修改
        result = await modify_ppt_visual(
            pptx_path=str(pptx_path),
            user_input=body.user_input,
            output_path=str(new_pptx_path),
            api_key_token=body.api_key_token,
            user_id=user_id,
            analyze_before_modify=body.analyze_before_modify
        )

        if result["success"]:
            _register_ppt_owner(new_task_id, user_id)
            return {
                "success": True,
                "message": result["message"],
                "task_id": new_task_id,
                "download_url": f"/api/v1/pptx/download/{new_task_id}",
                "preview_url": f"/api/v1/pptx/preview/{new_task_id}",
                "intent": result.get("intent"),
                "analysis": result.get("analysis"),
                "preview_count": result.get("preview_count", 0)
            }
        else:
            return {
                "success": False,
                "message": result["message"],
                "intent": result.get("intent"),
                "analysis": result.get("analysis")
            }

    except Exception as exc:
        logger.error("PPT 视觉修改失败 | task_id: %s | error: %s", task_id, exc)
        raise HTTPException(status_code=500, detail=f"修改失败: {exc}")


@router.get("/pptx/{task_id}/analyze")
async def analyze_ppt_endpoint(
    task_id: str,
    slide_number: Optional[int] = Query(None, description="指定幻灯片编号"),
    token: dict = Depends(verify_token),
):
    """
    分析 PPT 视觉状态（不执行修改）

    返回 PPT 的布局、字体、颜色等信息，用于了解当前状态。
    """
    _verify_ppt_owner(task_id, token.get("sub", "anonymous"))
    output_dir = PPT_OUTPUT_DIR
    pptx_path = output_dir / f"{task_id}.pptx"

    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail=f"找不到 PPT 文件：{task_id}")

    try:
        from app.utils.pptx.visual_modifier import analyze_ppt_for_modification

        result = await analyze_ppt_for_modification(
            pptx_path=str(pptx_path),
            slide_number=slide_number
        )

        return result

    except Exception as exc:
        logger.error("PPT 分析失败 | task_id: %s | error: %s", task_id, exc)
        raise HTTPException(status_code=500, detail=f"分析失败: {exc}")


# =============================================================================
# 模板与历史管理
# =============================================================================


@router.get("/pptx/templates")
async def list_ppt_templates(
    category: Optional[str] = Query(None, description="按分类筛选"),
    scenario: Optional[str] = Query(
        None,
        pattern="^(business|data_report|product_pitch|academic|education|general)$",
        description="按演示场景推荐",
    ),
    topic: str = Query("", max_length=5000, description="用于场景识别的主题"),
    token: dict = Depends(verify_token)
):
    """获取可用 PPT 模板列表"""
    if scenario or topic:
        recommendation = TemplateManager().recommend_for_scenario(topic, scenario=scenario)
        return recommendation
    templates = []
    for tpl_id, tpl_config in PPT_TEMPLATES.items():
        if category and not tpl_id.startswith(category):
            continue
        templates.append({
            "id": tpl_id,
            "name": tpl_config["name"],
            "primary_color": tpl_config["primary_color"],
            "background": tpl_config["background"],
        })
    return {"templates": templates}


@router.get("/pptx/history")
async def list_ppt_history(
    page: int = Query(1, ge=1),
    page_size: int = Query(20, ge=1, le=100),
    token: dict = Depends(verify_token)
):
    """获取用户的 PPT 生成历史"""
    user_id = token.get("sub", "anonymous")
    output_dir = PPT_OUTPUT_DIR

    if not output_dir.exists():
        return {"records": [], "total": 0}

    records = []
    for json_path in sorted(output_dir.glob("*_slides.json"), reverse=True):
        ppt_id = json_path.name.replace("_slides.json", "")
        try:
            owner_path = PPT_OWNER_DIR / f"{ppt_id}.json"
            if not owner_path.exists():
                continue
            with open(owner_path, "r", encoding="utf-8") as owner_file:
                if str(json.load(owner_file).get("user_id")) != str(user_id):
                    continue
            with open(json_path, "r", encoding="utf-8") as f:
                slides_data = json.load(f)
            pptx_path = output_dir / f"{ppt_id}.pptx"
            records.append({
                "task_id": ppt_id,
                "title": slides_data[0].get("title", "未命名") if slides_data else "未命名",
                "slide_count": len(slides_data),
                "has_file": pptx_path.exists(),
                "created_at": datetime.fromtimestamp(json_path.stat().st_mtime).isoformat(),
            })
        except Exception:
            continue

    total = len(records)
    start = (page - 1) * page_size
    end = start + page_size
    return {"records": records[start:end], "total": total}


@router.delete("/pptx/history/{task_id}")
async def delete_ppt_history(
    task_id: str,
    token: dict = Depends(verify_token)
):
    """删除指定 PPT 历史记录及其文件"""
    user_id = token.get("sub", "anonymous")
    _verify_ppt_owner(task_id, user_id)
    output_dir = PPT_OUTPUT_DIR
    json_path = output_dir / f"{task_id}_slides.json"

    if not json_path.exists():
        raise HTTPException(status_code=404, detail="记录不存在")

    json_path.unlink(missing_ok=True)
    for ext in ["pptx", "html", "md"]:
        (output_dir / f"{task_id}.{ext}").unlink(missing_ok=True)

    return {"success": True, "message": f"已删除 {task_id}"}


@router.get("/pptx/history/stats")
async def get_ppt_stats(
    token: dict = Depends(verify_token)
):
    """获取 PPT 生成统计信息"""
    user_id = str(token.get("sub", "anonymous"))
    output_dir = PPT_OUTPUT_DIR
    if not output_dir.exists():
        return {"total": 0, "completed": 0, "failed": 0}

    owned_ids = []
    for owner_path in PPT_OWNER_DIR.glob("*.json"):
        try:
            with open(owner_path, "r", encoding="utf-8") as owner_file:
                if str(json.load(owner_file).get("user_id")) == user_id:
                    owned_ids.append(owner_path.stem)
        except (OSError, ValueError):
            continue
    total = sum((output_dir / f"{ppt_id}_slides.json").exists() for ppt_id in owned_ids)
    completed = sum((output_dir / f"{ppt_id}.pptx").exists() for ppt_id in owned_ids)
    return {"total": total, "completed": completed, "failed": 0}


# =============================================================================
# 文本防溢出
# =============================================================================

def prevent_text_overflow(
    text: str, 
    max_chars_per_line: int = 70, 
    max_lines: int = 6, 
    shrink_font: bool = True
) -> List[str]:
    """
    防止文本在 PPT 幻灯片中溢出。
    
    使用新的智能文本处理模块:
    - 中文按字数换行
    - 英文按单词边界换行
    - 自动字号调整
    
    Args:
        text: 原始文本
        max_chars_per_line: 每行最大字符数
        max_lines: 最大行数限制
        shrink_font: 是否建议缩小字号
        
    Returns:
        处理后的文本行列表
    """
    layout = prevent_text_overflow_v2(
        text,
        max_chars_per_line=max_chars_per_line,
        max_lines=max_lines,
    )
    
    if layout.needs_overflow_warning:
        logger.warning(layout.overflow_message)
    
    return layout.lines


# =============================================================================
# 自动搜图
# =============================================================================

# 图片搜索管理器 (延迟初始化)
_image_search_manager: Optional[ImageSearchManager] = None


def get_image_search_manager() -> ImageSearchManager:
    """获取图片搜索管理器单例"""
    global _image_search_manager
    if _image_search_manager is None:
        _image_search_manager = ImageSearchManager(
            bing_key=os.environ.get("BING_IMAGE_SEARCH_KEY"),
            unsplash_key=os.environ.get("UNSPLASH_ACCESS_KEY"),
            pexels_key=os.environ.get("PEXELS_API_KEY"),
        )
    return _image_search_manager


async def search_image_url(keyword: str) -> Optional[str]:
    """
    搜索图片 URL
    
    使用多源聚合搜索:
    1. Bing Image Search (需要 API Key)
    2. Unsplash (需要 API Key)
    3. Pexels (需要 API Key)
    4. 占位图降级
    """
    manager = get_image_search_manager()
    return await manager.search_image(keyword)


async def download_image(url: str, save_path: Path) -> bool:
    """下载图片到本地"""
    import aiohttp
    
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=10) as resp:
                if resp.status == 200:
                    data = await resp.read()
                    save_path.write_bytes(data)
                    return True
    except Exception as e:
        logger.warning(f"图片下载失败 {url}: {e}")
    
    return False


IMAGE_CACHE_DIR = Path("./static/images/cache")


async def get_image_for_slide(keywords: List[str], slide_index: int) -> Optional[str]:
    """获取幻灯片配图 (缓存优先)"""
    manager = get_image_search_manager()
    
    for kw in keywords[:2]:
        # 检查缓存
        cached = await manager.get_cached_image(kw)
        if cached:
            return str(cached)
        
        # 搜索并下载
        url = await manager.search_image(kw)
        if url:
            path = await manager.download_and_cache(kw, url)
            if path:
                return str(path)
    
    return None


# =============================================================================
# PPT Agent API 端点
# =============================================================================

class OutlineSlide(BaseModel):
    """大纲中的单页幻灯片"""
    type: str = Field(..., description="幻灯片类型")
    title: str = Field(..., description="幻灯片标题")
    bullets: List[str] = Field(default_factory=list, description="要点列表")
    image_keywords: List[str] = Field(default_factory=list, description="配图关键词")
    notes: str = Field(default="", description="备注")
    narrative_role: str = Field(default="", description="商业叙事角色")
    content_blocks: List[Dict[str, Any]] = Field(default_factory=list, description="结构化商业内容块")

class OutlineGenerationRequest(BaseModel):
    """大纲生成请求"""
    topic: str = Field(..., description="PPT 主题", max_length=500)
    description: str = Field(default="", description="详细描述", max_length=2000)
    num_slides: int = Field(default=10, ge=1, le=PPT_MAX_SLIDES, description="幻灯片数量")
    model: str = Field(default=PPT_DEFAULT_MODEL, description="AI 模型")
    api_key_token: Optional[str] = Field(None, description="用户 API Key Token")

class OutlineGenerationResponse(BaseModel):
    """大纲生成响应"""
    title: str
    slides: List[OutlineSlide]
    total_slides: int
    outline_id: Optional[str] = None
    outline_version: Optional[int] = None
    status: Optional[str] = None


@router.post("/generate-text", response_model=OutlineGenerationResponse)
async def generate_ppt_from_text(
    req: OutlineGenerationRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    """
    自然语言生成 PPT 大纲 (仅返回结构化数据，不生成文件)
    """
    user_id = token.get("sub", "anonymous")
    logger.info("PPT Agent 请求 | user: %s | topic: %s", user_id, req.topic[:50])

    try:
        draft = await persist_create_ppt_outline(
            db,
            str(user_id),
            OutlineCreateRequest(
                topic=req.topic,
                description=req.description,
                num_slides=req.num_slides,
                model=req.model,
                api_key_token=req.api_key_token,
            ),
        )

        return OutlineGenerationResponse(
            title=draft.title,
            slides=[
                OutlineSlide(
                    type=slide.slide_type,
                    title=slide.title,
                    bullets=[block.content for block in slide.content_blocks],
                    image_keywords=slide.asset_intent.keywords if slide.asset_intent else [],
                    notes=slide.speaker_notes,
                    narrative_role=slide.narrative_role,
                    content_blocks=[block.model_dump(mode="json") for block in slide.content_blocks],
                )
                for slide in draft.slides
            ],
            total_slides=len(draft.slides),
            outline_id=draft.id,
            outline_version=draft.version,
            status=draft.status,
        )

    except Exception as exc:
        logger.error("PPT Agent 大纲生成失败 | error: %s", exc)
        raise HTTPException(status_code=500, detail=f"大纲生成失败: {exc}")


@router.post("/generate-from-text", response_model=TaskResponse)
async def generate_ppt_from_text_task(
    req: OutlineGenerationRequest,
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    """
    端到端: 自然语言 -> 大纲 -> PPTX 文件 (任务队列)
    """
    user_id = token.get("sub", "anonymous")
    logger.info("PPT Agent 端到端请求 | user: %s | topic: %s", user_id, req.topic[:50])

    try:
        draft = await persist_create_ppt_outline(
            db,
            str(user_id),
            OutlineCreateRequest(
                topic=req.topic,
                description=req.description,
                num_slides=req.num_slides,
                model=req.model,
                api_key_token=req.api_key_token,
            ),
        )
        approved = await persist_approve_ppt_outline(db, str(user_id), draft.id)
        return await generate_ppt_from_approved_outline(
            approved.id,
            OutlineGenerateRequest(outline_version=approved.version),
            token,
            db,
        )

    except Exception as exc:
        logger.error("PPT Agent 端到端失败 | error: %s", exc)
        raise HTTPException(status_code=500, detail=f"生成失败: {exc}")


# =============================================================================
# 文件上传生成 PPT
# =============================================================================

# 支持的文件扩展名
_PPT_FILE_EXTENSIONS = {
    ".txt", ".md", ".pdf", ".docx", ".doc",
    ".py", ".js", ".ts", ".json", ".yaml", ".yml",
    ".csv", ".log", ".rst", ".html", ".css",
}


async def _stream_upload_to_path(file: UploadFile, destination: Path, max_size: int) -> int:
    """Persist an upload in bounded chunks and return its byte count."""
    total_size = 0
    with destination.open("wb") as output:
        while chunk := await file.read(1024 * 1024):
            total_size += len(chunk)
            if total_size > max_size:
                raise HTTPException(status_code=400, detail=f"文件过大，最大支持 {max_size // 1024 // 1024}MB")
            output.write(chunk)
    return total_size

@router.post("/pptx/generate_from_file", response_model=TaskResponse)
async def generate_ppt_from_file(
    file: UploadFile = FastAPIFile(..., description="上传文件 (PDF/Word/TXT/MD 等)"),
    template: str = Form(default="modern", description="模板风格"),
    slide_count: int = Form(default=10, ge=5, le=PPT_MAX_SLIDES, description="页数"),
    output_format: str = Form(default="pptx", description="输出格式 (pptx/html/markdown)"),
    extra_prompt: str = Form(default="", description="额外提示词"),
    api_key_token: Optional[str] = Form(default=None, description="用户 API Key Token"),
    token: dict = Depends(verify_token),
    db: AsyncSession = Depends(get_db),
):
    """
    上传文件生成 PPT

    支持格式: PDF, Word, TXT, Markdown, 代码文件等
    流程: 上传文件 → 解析内容 → AI 生成大纲 → 生成 PPT
    """
    user_id = token.get("sub", "anonymous")

    # 验证文件扩展名
    filename = file.filename or "unknown"
    suffix = Path(filename).suffix.lower()
    if suffix not in _PPT_FILE_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件格式: {suffix}。支持: {', '.join(sorted(_PPT_FILE_EXTENSIONS))}"
        )

    # 保存上传文件到临时目录
    upload_dir = Path("./uploads/ppt_uploads")
    upload_dir.mkdir(parents=True, exist_ok=True)
    file_id = str(uuid.uuid4())
    temp_path = upload_dir / f"{file_id}{suffix}"

    try:
        await _stream_upload_to_path(file, temp_path, 50 * 1024 * 1024)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件保存失败: {e}")

    # 解析文件内容
    try:
        from app.utils.aicloud.knowledge_processor import parse_document
        parsed_text = parse_document(str(temp_path))
        if not parsed_text or not parsed_text.strip():
            raise HTTPException(status_code=400, detail="文件内容为空或无法解析")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件解析失败: {e}")

    # 截断过长内容（避免 token 超限）
    max_chars = 30000
    if len(parsed_text) > max_chars:
        parsed_text = parsed_text[:max_chars] + f"\n\n[内容过长，已截断至 {max_chars} 字符]"

    # 构建 PPT 请求
    topic = f"根据以下文件内容生成 PPT：\n\n文件名: {filename}\n\n文件内容:\n{parsed_text}"
    if extra_prompt:
        topic += f"\n\n用户要求: {extra_prompt}"

    req = PPTGenerationRequest(
        topic=topic,
        template=template,
        slide_count=slide_count,
        output_format=output_format,
        api_key_token=api_key_token,
    )

    ppt_id = str(uuid.uuid4())

    async def run_file_ppt_generation(task_id: str, **kwargs):
        _register_ppt_owner(task_id, user_id)
        async def update_progress(progress: int = 0, message: str = "", status: str = None, result_data: str = None, **_kwargs):
            await task_manager.update_progress(
                task_id,
                progress,
                message,
                status=status,
                result_data=result_data,
                error_message=_kwargs.get("error_message"),
            )

        try:
            await update_progress(progress=5, message="正在解析文件内容...")

            # 生成大纲
            await update_progress(progress=20, message="正在根据文件内容生成 PPT 大纲...")
            outline = await generate_ppt_outline(req, user_id=user_id)

            # 保存大纲快照
            output_dir = PPT_OUTPUT_DIR
            output_dir.mkdir(exist_ok=True)
            snapshot_path = output_dir / f"{task_id}_slides.json"
            with open(snapshot_path, 'w', encoding='utf-8') as f:
                json.dump(outline.get('slides', []), f, ensure_ascii=False, indent=2)

            # 根据格式生成文件
            ext_map = {"pptx": "pptx", "html": "html", "markdown": "md"}
            ext = ext_map.get(output_format, "pptx")
            filepath = output_dir / f"{task_id}.{ext}"

            if output_format == "pptx":
                await generate_pptx_file_enhanced(filepath, outline, req, update_progress=update_progress)
            elif output_format == "html":
                await update_progress(progress=60, message="正在生成 HTML 格式...")
                await generate_html_ppt(filepath, outline, req)
            elif output_format == "markdown":
                await update_progress(progress=60, message="正在生成 Markdown 格式...")
                await generate_markdown_ppt(filepath, outline, req)
            else:
                await generate_pptx_file_enhanced(filepath, outline, req, update_progress=update_progress)

            result = {
                "filename": filepath.name,
                "ppt_id": task_id,
                "source_file": filename,
                "download_url": f"/api/v1/pptx/download/{task_id}?format={ext}",
                "preview_url": f"/api/v1/pptx/preview/{task_id}?format={output_format}" if output_format in {"pptx", "html", "markdown"} else None,
            }
            await update_progress(
                progress=100,
                message="PPT 生成完成",
                status="completed",
                result_data=json.dumps(result)
            )
            return result

        except asyncio.CancelledError:
            await update_progress(status="cancelled", message="任务已取消")
        except Exception as e:
            await update_progress(status="failed", message=f"生成失败: {e}", error_message=str(e))
            logger.error(f"文件 PPT 生成失败 | task_id: {task_id} | error: {e}")
        finally:
            # 清理临时文件
            try:
                if temp_path.exists():
                    temp_path.unlink()
            except Exception:
                pass

    task_response = await task_manager.create_task(
        task_type="ppt_generation",
        user_id=user_id,
        func=run_file_ppt_generation,
        params={},
    )

    logger.info(f"文件 PPT 生成任务 | task_id: {task_response} | file: {filename} | user: {user_id}")
    return TaskResponse(
        task_id=task_response,
        task_type="ppt_generation",
        status="pending",
        progress=0,
        progress_message="等待中...",
        created_at=datetime.now().isoformat()
    )


# =============================================================================
# 自定义模板上传
# =============================================================================

@router.post("/pptx/templates/upload")
async def upload_custom_template(
    file: UploadFile = FastAPIFile(..., description="上传 PPTX 模板文件"),
    name: str = Form(..., description="模板名称"),
    description: str = Form(default="", description="模板描述"),
    token: dict = Depends(verify_token),
):
    """
    上传自定义 PPT 模板

    上传 .pptx 文件作为模板，系统会解析母版配置（配色、字体、布局）
    """
    user_id = token.get("sub", "anonymous")

    # 验证文件类型
    filename = file.filename or "unknown"
    if not filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="仅支持 .pptx 格式")

    # 保存模板文件
    template_dir = Path("./configs/ppt/custom_templates")
    template_dir.mkdir(parents=True, exist_ok=True)

    template_id = f"custom_{user_id}_{uuid.uuid4().hex[:8]}"
    template_path = template_dir / f"{template_id}.pptx"

    try:
        await _stream_upload_to_path(file, template_path, 20 * 1024 * 1024)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"模板保存失败: {e}")

    # 解析模板配置
    try:
        from app.utils.pptx.custom_template import CustomTemplateParser
        parser = CustomTemplateParser()
        config = parser.parse(str(template_path))

        # 保存配置
        config_path = template_dir / f"{template_id}.json"
        import json as _json
        with open(config_path, 'w', encoding='utf-8') as f:
            _json.dump(config, f, ensure_ascii=False, indent=2)

        return {
            "template_id": template_id,
            "name": name,
            "description": description,
            "config": config,
            "message": "模板上传成功",
        }
    except Exception as e:
        logger.warning(f"模板解析失败: {e}")
        # 即使解析失败也保留文件，用户可以手动使用
        return {
            "template_id": template_id,
            "name": name,
            "description": description,
            "config": None,
            "message": f"模板上传成功，但自动解析失败: {e}",
        }


@router.get("/pptx/templates/custom")
async def list_custom_templates(
    token: dict = Depends(verify_token),
):
    """列出用户上传的自定义模板"""
    user_id = token.get("sub", "anonymous")
    template_dir = Path("./configs/ppt/custom_templates")

    if not template_dir.exists():
        return {"templates": []}

    templates = []
    for json_file in template_dir.glob("*.json"):
        try:
            import json as _json
            with open(json_file, 'r', encoding='utf-8') as f:
                config = _json.load(f)
            template_id = json_file.stem
            # 只返回当前用户的模板
            if f"custom_{user_id}_" in template_id:
                templates.append({
                    "template_id": template_id,
                    "config": config,
                })
        except Exception:
            pass

    return {"templates": templates}


# =============================================================================
# PDF 导出
# =============================================================================

@router.get("/pptx/download/{ppt_id}/pdf")
async def download_ppt_as_pdf(
    ppt_id: str,
    token: dict = Depends(verify_token),
):
    """将 PPT 转换为 PDF 并下载"""
    _verify_ppt_owner(ppt_id, token.get("sub", "anonymous"))
    # 查找 PPTX 文件
    pptx_path = PPT_OUTPUT_DIR / f"{ppt_id}.pptx"
    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail="PPT 文件不存在")

    pdf_path = PPT_OUTPUT_DIR / f"{ppt_id}.pdf"

    try:
        await _convert_pptx_to_pdf(pptx_path, pdf_path)
    except FileNotFoundError:
        raise HTTPException(
            status_code=501,
            detail="PDF 导出需要安装 LibreOffice。请在服务器上安装: apt-get install libreoffice-impress"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF 转换失败: {e}")

    if not pdf_path.exists():
        raise HTTPException(status_code=500, detail="PDF 转换失败")

    return FileResponse(
        path=str(pdf_path),
        filename=f"{ppt_id}.pdf",
        media_type="application/pdf",
    )


async def _convert_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> None:
    """Convert a generated PPTX using the server's LibreOffice installation."""
    import asyncio
    import subprocess

    def convert() -> None:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(pdf_path.parent), str(pptx_path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice 转换失败: {result.stderr}")
        generated = pptx_path.with_suffix(".pdf")
        if generated != pdf_path and generated.exists():
            generated.replace(pdf_path)
        if not pdf_path.exists():
            raise RuntimeError("LibreOffice 未生成 PDF 文件")

    await asyncio.to_thread(convert)


# =============================================================================
# WebSocket 进度推送
# =============================================================================

@router.websocket("/ws/ppt/{task_id}")
async def ppt_progress_websocket(websocket: WebSocket, task_id: str):
    """
    WebSocket 端点用于实时 PPT 生成进度推送

    客户端连接后会自动接收任务进度更新。
    消息格式:
    - {"type": "progress", "progress": 0.5, "step": "generating", "message": "正在生成第 5 页..."}
    - {"type": "completed", "progress": 1, "step": "completed", "message": "任务完成"}
    - {"type": "error", "error": "错误信息"}
    """
    access_token = websocket.query_params.get("token", "")
    valid, payload, close_code, reason = verify_token_ws(access_token)
    if not valid or not payload:
        await websocket.close(code=close_code or 1008, reason=reason or "未授权")
        return
    try:
        _verify_ppt_owner(task_id, payload.get("sub", ""))
    except HTTPException as error:
        await websocket.close(code=1008, reason=str(error.detail))
        return

    await websocket.accept()
    logger.info(f"PPT WebSocket 连接已建立 | task_id={task_id}")

    try:
        # 获取任务管理器
        from app.utils.task_manager import task_manager

        last_progress = -1
        last_status = ""
        try:
            last_event_sequence = max(0, int(websocket.query_params.get("after_sequence", "0")))
        except ValueError:
            last_event_sequence = 0
        snapshot_sent = False

        while True:
            # SQL 事件日志提供断线后的可靠重放来源。
            try:
                from app.services.task_event_service import replay_task_events

                async with async_session() as event_db:
                    events = await replay_task_events(
                        event_db,
                        task_id,
                        int(payload.get("sub", 0)),
                        after_sequence=last_event_sequence,
                        limit=100,
                    )
                for event in events:
                    await websocket.send_json({
                        "type": event.event_type,
                        "sequence": event.sequence,
                        "status": event.status,
                        "progress": event.progress,
                        "payload": event.payload_json,
                        "created_at": event.created_at.isoformat() if event.created_at else None,
                    })
                    last_event_sequence = event.sequence
                if last_event_sequence > 0 and not events and not snapshot_sent:
                    from app.services.task_checkpoint_service import get_latest_checkpoint

                    async with async_session() as checkpoint_db:
                        checkpoint = await get_latest_checkpoint(
                            checkpoint_db,
                            task_id,
                            int(payload.get("sub", 0)),
                        )
                    if checkpoint:
                        await websocket.send_json({
                            "type": "snapshot_recovery",
                            "sequence": last_event_sequence,
                            "revision": checkpoint.revision,
                            "step": checkpoint.step,
                            "state": checkpoint.state_json,
                        })
                    snapshot_sent = True
            except Exception as event_error:
                logger.debug(f"PPT SQL 事件重放暂不可用 | task_id={task_id} | error={event_error}")

            # 查询任务状态
            task_info = await task_manager.get_task_info_async(task_id)

            if not task_info:
                async with async_session() as task_db:
                    task_result = await task_db.execute(
                        select(Task).where(
                            Task.task_id == task_id,
                            Task.user_id == int(payload.get("sub", 0)),
                        )
                    )
                    task_record = task_result.scalar_one_or_none()
                if task_record:
                    task_info = {
                        "task_id": task_record.task_id,
                        "status": task_record.status,
                        "progress": task_record.progress or 0,
                        "progress_message": task_record.progress_message or "",
                        "result": task_record.result or task_record.result_json or {},
                        "error_message": task_record.error_message,
                    }

            if not task_info:
                await websocket.send_json({
                    "type": "error",
                    "error": "任务不存在"
                })
                break

            current_progress = task_info.get("progress", 0)
            current_status = task_info.get("status", "")
            current_message = task_info.get("progress_message", "")

            # 只在进度或状态变化时发送更新
            if current_progress != last_progress or current_status != last_status:
                if current_status == "success":
                    await websocket.send_json({
                        "type": "completed",
                        "progress": 1,
                        "step": "completed",
                        "message": "任务完成",
                        "result": task_info.get("result", {})
                    })
                    break
                elif current_status == "failed":
                    await websocket.send_json({
                        "type": "error",
                        "error": task_info.get("error_message", "任务失败"),
                        "progress": current_progress / 100,
                        "step": "error",
                        "message": current_message
                    })
                    break
                elif current_status == "cancelled":
                    await websocket.send_json({
                        "type": "error",
                        "error": "任务已取消",
                        "progress": current_progress / 100,
                        "step": "cancelled",
                        "message": "任务已取消"
                    })
                    break
                else:
                    await websocket.send_json({
                        "type": "progress",
                        "progress": current_progress / 100,
                        "step": current_status,
                        "message": current_message
                    })

                last_progress = current_progress
                last_status = current_status

            # 检查是否有客户端消息（如 ping）
            try:
                data = await asyncio.wait_for(websocket.receive_text(), timeout=1.0)
                if data == "ping":
                    await websocket.send_text("pong")
            except asyncio.TimeoutError:
                pass
            except WebSocketDisconnect:
                break

            # 等待一段时间再查询
            await asyncio.sleep(0.5)

    except WebSocketDisconnect:
        logger.info(f"PPT WebSocket 连接断开 | task_id={task_id}")
    except Exception as e:
        logger.error(f"PPT WebSocket 错误 | task_id={task_id} | error={e}")
        try:
            await websocket.send_json({
                "type": "error",
                "error": f"服务器错误: {str(e)}"
            })
        except:
            pass
    finally:
        try:
            await websocket.close()
        except:
            pass
