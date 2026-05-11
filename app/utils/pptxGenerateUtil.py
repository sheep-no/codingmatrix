import asyncio
import json
import re
import gc
import os
import io
import httpx
from datetime import datetime
from enum import Enum
from dataclasses import dataclass, field
from typing import List, Literal, Dict, Optional, AsyncGenerator, Any
from app.utils.AiCodeUtil import call_siliconflow


class SlideType(Enum):
    TITLE = "title"
    CONTENT = "content"
    CHAPTER = "chapter"
    BULLET = "bullet"
    IMAGE = "image"
    CHART = "chart"
    END = "end"


@dataclass
class Slide:
    """单页幻灯片数据结构"""
    type: SlideType
    title: str
    content: List[str] = field(default_factory=list)
    subtitle: str = ""
    image_keywords: List[str] = field(default_factory=list)
    layout: str = "default"
    notes: str = ""
    design_suggestions: Dict[str, Any] = field(default_factory=dict)
    image_layout: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Presentation:
    """完整演示文稿"""
    title: str
    subtitle: str
    author: str
    theme: str
    slides: List[Slide]
    created_at: str = field(default_factory=lambda: datetime.now().isoformat())
    design_scheme: Dict[str, Any] = field(default_factory=dict)


class PptGenerator:
    """
    PPT智能体生成器（美化增强版）
    工作流程：
    1. AI分析主题+基本大纲结构
    2. AI撰写每页内容+填充slide对象（带字数限制）
    3. AI生成每页的美化建议（配色、布局、字体）
    4. 配图策略：提取关键词+网页搜索/占位符
    5. python-pptx渲染->输出.pptx文件（智能图片插入+自动清理占位符）
    """
    MODELS = {
        "outline": "Qwen/Qwen2.5-7B-Instruct",
        "content": "deepseek-ai/DeepSeek-R1-0528-Qwen3-8B",
        "polish": "Qwen/Qwen2.5-7B-Instruct",
        "design": "Qwen/Qwen2.5-7B-Instruct"
    }

    # 预定义的主题配色方案
    COLOR_SCHEMES = {
        "business": {
            "primary": "1F4E79",
            "secondary": "2E75B6",
            "accent": "70AD47",
            "background": "FFFFFF",
            "text": "333333",
            "title_font": "微软雅黑",
            "body_font": "微软雅黑"
        },
        "academic": {
            "primary": "203864",
            "secondary": "4472C4",
            "accent": "ED7D31",
            "background": "F8F9FA",
            "text": "212529",
            "title_font": "Times New Roman",
            "body_font": "宋体"
        },
        "pitch": {
            "primary": "C00000",
            "secondary": "FF6B6B",
            "accent": "FFD700",
            "background": "1A1A1A",
            "text": "FFFFFF",
            "title_font": "Arial Black",
            "body_font": "Arial"
        },
        "tutorial": {
            "primary": "5B9BD5",
            "secondary": "9DC3E6",
            "accent": "FFC000",
            "background": "FFFFFF",
            "text": "404040",
            "title_font": "Segoe UI",
            "body_font": "Segoe UI"
        }
    }

    # 【新增】单页字数限制配置 - 防止内容溢出
    MAX_CHARS_PER_SLIDE = {
        SlideType.TITLE: 50,  # 封面：标题10字 + 副标题20字
        SlideType.CHAPTER: 80,  # 章节页：标题20字 + 副标题30字
        SlideType.BULLET: 200,  # 要点页：标题20字 + 5要点×30字
        SlideType.CONTENT: 300,  # 内容页：标题20字 + 4段落×60字
        SlideType.IMAGE: 150,  # 图文页：标题20字 + 2说明×50字
        SlideType.CHART: 150,  # 图表页：标题20字 + 3洞察×40字
        SlideType.END: 100  # 结束页：标题20字 + 要点
    }

    def __init__(self):
        print(f"[{datetime.now().strftime('%H:%M:%S')}] PptGenerator 初始化...")
        self.image_search_cache = {}
        self.semaphore = asyncio.Semaphore(3)
        self.temp_image_dir = os.path.join(os.path.dirname(__file__), "temp_images")
        os.makedirs(self.temp_image_dir, exist_ok=True)
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 临时图片目录: {self.temp_image_dir}")

    async def _call_ai_with_retry(self, prompt: str, model: str, temperature: float = 0.3,
                                  max_tokens: int = 2048, max_retries: int = 3) -> Dict:
        """【新增】带重试机制的AI调用"""
        for attempt in range(max_retries):
            try:
                print(f"    AI调用尝试 {attempt + 1}/{max_retries}...")
                resp = await call_siliconflow(
                    prompt=prompt,
                    model=model,
                    temperature=temperature,
                    max_tokens=max_tokens
                )
                return resp
            except Exception as e:
                print(f"    尝试 {attempt + 1} 失败: {e}")
                if attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 2
                    print(f"    等待 {wait_time} 秒后重试...")
                    await asyncio.sleep(wait_time)
                else:
                    raise

    async def create_outline(
            self,
            topic: str,
            pages: int = 10,
            style: Literal["business", "academic", "pitch", "tutorial"] = "business",
            audience: str = "general"
    ) -> List[Dict]:
        """
        Step 1: AI 生成 PPT 大纲结构（包含美化需求）
        【修复】更严格的提示词和错误处理
        """
        print(
            f"[{datetime.now().strftime('%H:%M:%S')}] 开始生成大纲: topic='{topic}', pages={pages}, style='{style}'")

        color_scheme = self.COLOR_SCHEMES.get(style, self.COLOR_SCHEMES["business"])

        # 【修复】更严格、更清晰的提示词，增加字数限制说明
        prompt = f"""你是一位专业的PPT结构设计师。请为主题"{topic}"设计一个{pages}页的演示大纲。

要求：
1. 输出必须是有效的JSON数组格式
2. 不要包含任何markdown代码块标记（如```json）
3. 不要包含任何解释性文字
4. 只输出JSON数组，从[开始，到]结束
5. 【重要】每页内容必须简洁，标题不超过15字，描述不超过30字

受众：{audience}
风格：{style}

每页必须包含以下字段：
- type: 页面类型（title/chapter/content/bullet/image/chart/end）
- title: 页面标题（简洁有力，不超过15字）
- purpose: 这页的核心目的（简短描述，不超过20字）
- image_keywords: 建议配图关键词（数组，可为空，每个词不超过6字）
- layout_suggestion: 布局建议（如"左文右图"、"全屏背景"等）
- design_notes: 设计要点说明（不超过30字）

页面结构要求：
- 第1页必须是 "type": "title"
- 最后1页必须是 "type": "end"
- 中间包含1-2个 "type": "chapter"
- 其余为 content/bullet/image/chart

示例格式：
[{{"type": "title", "title": "校园安全", "purpose": "封面展示", "image_keywords": ["校园"], "layout_suggestion": "全屏背景+居中文字", "design_notes": "大字号标题"}}]

现在直接输出JSON数组："""

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 调用 AI 生成大纲...")
        start_time = datetime.now()

        try:
            resp = await self._call_ai_with_retry(
                prompt=prompt,
                model=self.MODELS['outline'],
                temperature=0.2,  # 【修复】降低temperature提高确定性
                max_tokens=2048,
                max_retries=3
            )

            elapsed = (datetime.now() - start_time).total_seconds()
            print(f"[{datetime.now().strftime('%H:%M:%S')}] 大纲生成完成，耗时 {elapsed:.2f} 秒")

            content = resp["choices"][0]["message"]["content"]

            # 【修复】添加内容清理
            content = self._clean_ai_response(content)

            result = self._extract_json(content)

            # 验证结果
            if not isinstance(result, list):
                raise ValueError(f"AI返回的不是数组，而是 {type(result)}")

            if len(result) == 0:
                raise ValueError("AI返回了空数组")

            print(f"[{datetime.now().strftime('%H:%M:%S')}] 大纲解析成功，共 {len(result)} 页")

            for i, item in enumerate(result, 1):
                print(f"  页 {i}: [{item.get('type', 'unknown')}] {item.get('title', '无标题')}")

            return result

        except Exception as e:
            print(f"[{datetime.now().strftime('%H:%M:%S')}] 大纲生成失败: {e}")
            # 【修复】返回默认大纲作为fallback
            print("使用默认大纲作为备用方案...")
            return self._get_default_outline(topic, pages)

    def _clean_ai_response(self, text: str) -> str:
        """【新增】清理AI响应内容"""
        if not text:
            return text

        # 移除markdown代码块
        text = re.sub(r'```json\s*', '', text)
        text = re.sub(r'```\s*', '', text)

        # 移除开头的非JSON字符
        text = re.sub(r'^[^{[]+', '', text)

        # 移除结尾的非JSON字符
        text = re.sub(r'[^}\]]+$', '', text)

        # 修复常见的JSON格式错误
        text = text.replace('"""', '"')
        text = text.replace('""', '"')

        return text.strip()

    def _get_default_outline(self, topic: str, pages: int) -> List[Dict]:
        """【新增】获取默认大纲（当AI失败时使用）"""
        print(f"生成默认大纲: {topic}, {pages}页")

        outline = [
            {
                "type": "title",
                "title": topic[:15],
                "purpose": "封面展示",
                "image_keywords": [topic[:6]],
                "layout_suggestion": "居中布局",
                "design_notes": "简洁大气"
            }
        ]

        # 生成中间页
        content_types = ["content", "bullet", "image", "chart"]
        for i in range(pages - 2):
            outline.append({
                "type": content_types[i % len(content_types)],
                "title": f"第{i + 1}部分"[:15],
                "purpose": "内容展示",
                "image_keywords": [] if i % 2 == 0 else [topic[:6]],
                "layout_suggestion": "左文右图" if i % 2 == 0 else "全文字",
                "design_notes": "清晰易读"
            })

        # 结束页
        outline.append({
            "type": "end",
            "title": "谢谢观看",
            "purpose": "结束",
            "image_keywords": [],
            "layout_suggestion": "居中布局",
            "design_notes": "简洁收尾"
        })

        return outline

    async def generate_slide_content(self, outline_item: Dict, context: str, style: str = "business") -> Slide:
        """Step 2: 为单页大纲生成详细内容（带字数限制）"""
        slide_type = SlideType(outline_item.get('type', 'content'))
        print(
            f"[{datetime.now().strftime('%H:%M:%S')}] 生成内容: [{slide_type.value}] {outline_item.get('title', '无标题')[:30]}...")

        # 获取该类型页面的字数限制
        max_chars = self.MAX_CHARS_PER_SLIDE.get(slide_type, 200)

        async with self.semaphore:
            prompts = {
                SlideType.TITLE: f"""为主题"{outline_item.get('title', '')}"写PPT封面文案。
【重要限制】总字数不超过{max_chars}字，标题不超过10字，副标题不超过20字。
要求：输出纯JSON，不要markdown，不要解释。
格式：{{"title": "主标题", "subtitle": "副标题", "content": []}}""",

                SlideType.BULLET: f"""为主题"{outline_item.get('title', '')}"写5个核心要点。
【重要限制】总字数不超过{max_chars}字，每点不超过25字，标题不超过15字。
要求：输出纯JSON，不要markdown。
格式：{{"title": "页面标题", "content": ["要点1", "要点2", ...], "notes": "备注"}}""",

                SlideType.CONTENT: f"""为主题"{outline_item.get('title', '')}"写详细内容页。
【重要限制】总字数不超过{max_chars}字，每段不超过50字，标题不超过15字，最多4段。
要求：输出纯JSON，不要markdown，不要解释。
格式：{{"title": "...", "content": ["段落1", ...], "notes": "..."}}""",

                SlideType.CHAPTER: f"""写章节过渡页文案。
【重要限制】总字数不超过{max_chars}字，标题不超过15字，副标题不超过25字。
要求：输出纯JSON，不要markdown。
格式：{{"title": "章节名", "subtitle": "一句话总结", "content": []}}""",

                SlideType.END: f"""写结束页文案。
【重要限制】总字数不超过{max_chars}字，标题不超过15字，副标题不超过20字。
要求：输出纯JSON，不要markdown。
格式：{{"title": "感谢观看", "subtitle": "...", "content": ["核心要点"]}}""",

                SlideType.IMAGE: f"""为主题"{outline_item.get('title', '')}"写图文页文案。
【重要限制】总字数不超过{max_chars}字，说明每句不超过40字，标题不超过15字。
要求：输出纯JSON，不要markdown。
格式：{{"title": "...", "content": ["说明1", ...], "notes": "..."}}""",

                SlideType.CHART: f"""为主题"{outline_item.get('title', '')}"写数据图表页。
【重要限制】总字数不超过{max_chars}字，每个洞察不超过40字，标题不超过15字。
要求：输出纯JSON，不要markdown。
格式：{{"title": "...", "content": ["洞察1", ...], "notes": "..."}}"""
            }

            prompt = prompts.get(slide_type, prompts[SlideType.CONTENT])

            try:
                resp = await self._call_ai_with_retry(
                    prompt=prompt,
                    model=self.MODELS['content'],
                    temperature=0.2,
                    max_tokens=1024,
                    max_retries=2
                )

                content = resp["choices"][0]["message"]["content"]
                content = self._clean_ai_response(content)
                data = self._extract_json(content)

                # 【新增】截断超出字数限制的内容
                data = self._truncate_slide_content(data, slide_type, max_chars)

                print(
                    f"[{datetime.now().strftime('%H:%M:%S')}] 内容生成完成: {data.get('title', '无标题')[:30]} (字数: {self._count_chars(data)})")

                # 生成美化建议
                design_suggestions = await self._generate_design_suggestions(
                    slide_type=slide_type,
                    title=data.get("title", outline_item.get('title', '')),
                    content=data.get("content", []),
                    layout_hint=outline_item.get("layout_suggestion", ""),
                    style=style
                )

                image_layout = self._parse_image_layout(
                    outline_item.get("layout_suggestion", ""),
                    slide_type
                )

                await asyncio.sleep(0.05)

                return Slide(
                    type=slide_type,
                    title=data.get("title", outline_item.get('title', '无标题')),
                    content=data.get("content", []),
                    subtitle=data.get("subtitle", ""),
                    image_keywords=outline_item.get("image_keywords", []),
                    notes=data.get("notes", ""),
                    design_suggestions=design_suggestions,
                    image_layout=image_layout
                )

            except Exception as e:
                print(f"[{datetime.now().strftime('%H:%M:%S')}] 内容生成失败，使用默认内容: {e}")
                return self._get_default_slide(slide_type, outline_item)

    def _count_chars(self, data: Dict) -> int:
        """【新增】计算幻灯片内容字数"""
        total = len(data.get("title", ""))
        total += len(data.get("subtitle", ""))
        for item in data.get("content", []):
            total += len(str(item))
        total += len(data.get("notes", ""))
        return total

    def _truncate_slide_content(self, data: Dict, slide_type: SlideType, max_chars: int) -> Dict:
        """【新增】截断超出字数限制的内容"""
        current_chars = self._count_chars(data)

        if current_chars <= max_chars:
            return data

        print(f"    警告: 内容超出字数限制 ({current_chars}/{max_chars})，进行截断...")

        # 截断标题
        title = data.get("title", "")
        if len(title) > 20:
            data["title"] = title[:18] + "..."

        # 截断副标题
        subtitle = data.get("subtitle", "")
        if len(subtitle) > 30:
            data["subtitle"] = subtitle[:28] + "..."

        # 截断内容列表
        content = data.get("content", [])
        truncated_content = []
        current_total = len(data.get("title", "")) + len(data.get("subtitle", "")) + len(data.get("notes", ""))

        for item in content:
            item_len = len(str(item))
            if current_total + item_len > max_chars - 10:  # 留10字余量
                # 截断最后一个项目
                remaining = max_chars - current_total - 10
                if remaining > 10:
                    truncated_item = str(item)[:remaining] + "..."
                    truncated_content.append(truncated_item)
                break
            truncated_content.append(item)
            current_total += item_len

        data["content"] = truncated_content

        # 截断备注
        notes = data.get("notes", "")
        if len(notes) > 50:
            data["notes"] = notes[:48] + "..."

        return data

    def _get_default_slide(self, slide_type: SlideType, outline_item: Dict) -> Slide:
        """【新增】获取默认幻灯片内容"""
        max_chars = self.MAX_CHARS_PER_SLIDE.get(slide_type, 200)

        # 根据类型生成默认内容
        if slide_type == SlideType.TITLE:
            content = []
        elif slide_type == SlideType.BULLET:
            content = ["要点一", "要点二", "要点三"]
        elif slide_type == SlideType.CHAPTER:
            content = []
        elif slide_type == SlideType.END:
            content = ["感谢观看"]
        else:
            content = ["内容待补充"]

        return Slide(
            type=slide_type,
            title=outline_item.get('title', '无标题')[:15],
            content=content,
            subtitle='',
            image_keywords=outline_item.get('image_keywords', []),
            notes='',
            design_suggestions=self._get_default_design_suggestions('business'),
            image_layout=self._parse_image_layout('', slide_type)
        )

    async def _generate_design_suggestions(
            self,
            slide_type: SlideType,
            title: str,
            content: List[str],
            layout_hint: str,
            style: str
    ) -> Dict[str, Any]:
        """为单页生成详细的美化建议"""
        color_scheme = self.COLOR_SCHEMES.get(style, self.COLOR_SCHEMES["business"])

        prompt = f"""为以下PPT页面提供设计建议，输出纯JSON：

页面类型: {slide_type.value}
标题: {title}
布局: {layout_hint}

输出格式（严格JSON）：
{{
  "color_usage": {{"title": "#{color_scheme['primary']}", "background": "#{color_scheme['background']}", "accent": "#{color_scheme['accent']}"}},
  "font_scheme": {{"title_size": 32, "body_size": 18, "title_font": "{color_scheme['title_font']}", "body_font": "{color_scheme['body_font']}"}},
  "spacing": {{"line_spacing": 1.5, "paragraph_spacing": 12}},
  "visual_elements": ["简洁设计"],
  "image_position": {{"position": "right", "width_ratio": 0.4}},
  "special_effects": []
}}

只输出JSON，不要其他内容："""

        try:
            resp = await self._call_ai_with_retry(
                prompt=prompt,
                model=self.MODELS['design'],
                temperature=0.2,
                max_tokens=512,
                max_retries=2
            )
            content_text = resp["choices"][0]["message"]["content"]
            content_text = self._clean_ai_response(content_text)
            suggestions = self._extract_json(content_text)
            print(f"    美化建议生成完成")
            return suggestions
        except Exception as e:
            print(f"    美化建议生成失败，使用默认值: {e}")
            return self._get_default_design_suggestions(style)

    def _get_default_design_suggestions(self, style: str) -> Dict[str, Any]:
        """获取默认美化建议"""
        scheme = self.COLOR_SCHEMES.get(style, self.COLOR_SCHEMES["business"])
        return {
            "color_usage": {
                "title": f"#{scheme['primary']}",
                "background": f"#{scheme['background']}",
                "accent": f"#{scheme['accent']}"
            },
            "font_scheme": {
                "title_size": 32,
                "body_size": 18,
                "title_font": scheme['title_font'],
                "body_font": scheme['body_font']
            },
            "spacing": {
                "line_spacing": 1.5,
                "paragraph_spacing": 12,
                "margin": 0.5
            },
            "visual_elements": ["保持简洁"],
            "image_position": {
                "position": "right",
                "width_ratio": 0.4,
                "height_ratio": 0.6
            },
            "special_effects": []
        }

    def _parse_image_layout(self, layout_suggestion: str, slide_type: SlideType) -> Dict[str, Any]:
        """解析布局建议为具体的图片布局配置"""
        suggestion_lower = layout_suggestion.lower() if layout_suggestion else ""

        default_layout = {
            "position": "right",
            "left": 8.0,
            "top": 1.3,
            "width": 4.5,
            "height": 5.5,
            "text_wrap": True
        }

        if slide_type == SlideType.TITLE:
            return {
                "position": "background",
                "left": 0,
                "top": 0,
                "width": 13.333,
                "height": 7.5,
                "opacity": 0.3,
                "text_wrap": False
            }
        elif slide_type == SlideType.CHAPTER:
            return {
                "position": "left_accent",
                "left": 0,
                "top": 0,
                "width": 2.0,
                "height": 7.5,
                "text_wrap": False
            }
        elif "左图右文" in layout_suggestion or "left" in suggestion_lower:
            return {
                "position": "left",
                "left": 0.5,
                "top": 1.3,
                "width": 5.0,
                "height": 5.5,
                "text_wrap": True
            }
        elif "上图下文" in layout_suggestion or "top" in suggestion_lower:
            return {
                "position": "top",
                "left": 1.0,
                "top": 1.5,
                "width": 11.333,
                "height": 3.5,
                "text_wrap": True
            }
        elif "全屏" in layout_suggestion or "full" in suggestion_lower:
            return {
                "position": "background",
                "left": 0,
                "top": 0,
                "width": 13.333,
                "height": 7.5,
                "opacity": 0.4,
                "text_wrap": False
            }

        return default_layout

    async def build_presentation(
            self,
            topic: str,
            pages: int = 10,
            style: str = "business",
            audience: str = "general",
            author: str = "AI Assistant"
    ) -> Presentation:
        """完整流程：大纲 → 内容 → 美化建议 → Presentation 对象"""
        print(f"\n{'=' * 60}")
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 开始构建 PPT: '{topic}'")
        print(f"{'=' * 60}\n")

        start_time = datetime.now()

        # 1. 生成大纲（带错误处理）
        outline = await self.create_outline(
            topic=topic,
            pages=pages,
            style=style,
            audience=audience
        )

        # 2. 分批生成内容
        slides = []
        batch_size = 3
        total = len(outline)

        print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 开始生成 {total} 页内容...")

        for i in range(0, total, batch_size):
            batch = outline[i:i + batch_size]
            batch_num = i // batch_size + 1
            total_batches = (total - 1) // batch_size + 1

            print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 批次 {batch_num}/{total_batches}")

            batch_start = datetime.now()
            batch_slides = await asyncio.gather(*[
                self.generate_slide_content(item, topic, style)
                for item in batch
            ], return_exceptions=True)

            # 处理可能的异常
            for slide in batch_slides:
                if isinstance(slide, Exception):
                    print(f"    警告: 某页生成失败: {slide}")
                    continue
                slides.append(slide)

            batch_elapsed = (datetime.now() - batch_start).total_seconds()
            print(f"[{datetime.now().strftime('%H:%M:%S')}] 批次 {batch_num} 完成，耗时 {batch_elapsed:.2f} 秒")

            if i > 0 and i % 3 == 0:
                gc.collect()
                await asyncio.sleep(0.3)

        total_elapsed = (datetime.now() - start_time).total_seconds()
        print(
            f"\n[{datetime.now().strftime('%H:%M:%S')}] PPT 内容生成完成！总耗时 {total_elapsed:.2f} 秒，共 {len(slides)} 页")

        design_scheme = self._generate_global_design_scheme(style, slides)

        ppt = Presentation(
            title=topic,
            subtitle=f"{style} style presentation",
            author=author,
            theme=style,
            slides=slides,
            design_scheme=design_scheme
        )

        return ppt

    def _generate_global_design_scheme(self, style: str, slides: List[Slide]) -> Dict[str, Any]:
        """生成整体设计规范"""
        scheme = self.COLOR_SCHEMES.get(style, self.COLOR_SCHEMES["business"])
        return {
            "color_palette": scheme,
            "total_slides": len(slides),
            "image_slides": len([s for s in slides if s.image_keywords]),
            "content_density": "medium",
            "consistency_rules": [
                "所有标题使用统一字体和主色",
                "正文使用统一字号和行距",
                "图片风格保持一致"
            ]
        }

    async def search_images(self, keywords: List[str], max_results: int = 3) -> List[str]:
        """使用高品图像 API 搜索图片"""
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 搜索图片: keywords={keywords[:2]}")

        urls = []
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Referer": "https://www.gaopinimages.com/",
            "Content-Type": "application/json"
        }

        async with httpx.AsyncClient(headers=headers) as client:
            for keyword in keywords[:2]:
                try:
                    payload = {
                        "keyType": 1,
                        "sortOrder": "1",
                        "from": 1,
                        "size": max_results,
                        "qk": keyword,
                        "style": 1,
                        "materialCategory": "",
                        "color": "",
                        "imagesShape": "",
                        "licenseType": "",
                        "materialType": "",
                        "modelSex": "",
                        "peopleNum": "",
                        "portraitureRight": "",
                        "secondCategory": "",
                        "huiXuanSelected": ""
                    }

                    resp = await client.post(
                        "https://www.gaopinimages.com/crest/search/searchImageV2",
                        json=payload,
                        timeout=10.0
                    )

                    if resp.status_code == 200:
                        data = resp.json()
                        if data.get("return_code") == "000000":
                            images = data.get("return_data", {}).get("data", [])
                            for img in images[:1]:
                                url = img.get("thumbnailUrl300C")
                                if url:
                                    urls.append(url)
                                    break

                    await asyncio.sleep(0.3)

                except Exception as e:
                    print(f"[{datetime.now().strftime('%H:%M:%S')}] 搜索失败 '{keyword}': {e}")
                    continue

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 图片搜索完成，共找到 {len(urls)} 张")
        return urls

    async def download_and_convert_image(self, url: str, output_dir: str) -> Optional[str]:
        """下载 webp 并转换为 PNG"""
        try:
            from PIL import Image

            async with httpx.AsyncClient() as client:
                resp = await client.get(url, timeout=15.0)

                if resp.status_code != 200:
                    return None

                img = Image.open(io.BytesIO(resp.content))

                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')

                img_id = url.split('/')[-1].split('.')[0]
                png_path = os.path.join(output_dir, f"{img_id}.png")
                img.save(png_path, 'PNG', quality=95)

                return png_path

        except Exception as e:
            print(f"[{datetime.now().strftime('%H:%M:%S')}] 图片转换失败: {e}")
            return None

    async def prepare_images(self, presentation: Presentation) -> Dict[int, str]:
        """为所有 slide 提前准备图片"""
        print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 开始准备图片资源...")

        image_map = {}
        need_image_count = sum(1 for s in presentation.slides if s.image_keywords)
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 共 {need_image_count} 页需要图片")

        for idx, slide in enumerate(presentation.slides):
            if slide.image_keywords:
                urls = await self.search_images(slide.image_keywords, max_results=1)

                if urls:
                    local_path = await self.download_and_convert_image(urls[0], self.temp_image_dir)
                    if local_path:
                        image_map[idx] = local_path
                        print(f"[{datetime.now().strftime('%H:%M:%S')}] 页 {idx + 1} 图片准备完成")

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 图片资源准备完成: {len(image_map)}/{need_image_count}")
        return image_map

    def render_to_pptx(
            self,
            presentation: Presentation,
            output_path: str,
            image_map: Dict[int, str] = None,
            remove_placeholders: bool = True
    ) -> str:
        """使用 python-pptx 生成 PPTX 文件（增强版）"""
        print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 开始渲染 PPTX: {output_path}")
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 可用图片: {len(image_map) if image_map else 0} 张")
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 清理空白占位符: {remove_placeholders}")

        try:
            from pptx import Presentation as PptxPresentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
        except ImportError as e:
            raise ImportError("请先安装 python-pptx: pip install python-pptx")

        start_time = datetime.now()

        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        color_scheme = self.COLOR_SCHEMES.get(presentation.theme, self.COLOR_SCHEMES["business"])

        layout_map = {
            SlideType.TITLE: 6,
            SlideType.CHAPTER: 5,
            SlideType.BULLET: 1,
            SlideType.CONTENT: 1,
            SlideType.IMAGE: 5,
            SlideType.CHART: 5,
            SlideType.END: 6
        }

        for idx, slide_data in enumerate(presentation.slides):
            layout_idx = layout_map.get(slide_data.type, 6)
            slide_layout = prs.slide_layouts[layout_idx]
            pptx_slide = prs.slides.add_slide(slide_layout)

            print(
                f"[{datetime.now().strftime('%H:%M:%S')}] 渲染页 {idx + 1}/{len(presentation.slides)}: [{slide_data.type.value}] {slide_data.title[:30]}...")

            self._render_slide_content(
                pptx_slide, slide_data, idx, image_map, color_scheme
            )

            if slide_data.notes:
                try:
                    notes_slide = pptx_slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data.notes
                except:
                    pass

        # 【关键修复】在保存前删除所有空白占位符
        if remove_placeholders:
            print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 开始清理空白占位符...")
            removed_count = self._remove_all_placeholders(prs)
            print(f"[{datetime.now().strftime('%H:%M:%S')}] 清理完成，删除 {removed_count} 个空白占位符")

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 所有页面渲染完成，开始保存文件...")

        save_start = datetime.now()
        prs.save(output_path)
        save_time = (datetime.now() - save_start).total_seconds()

        file_size = os.path.getsize(output_path)
        total_time = (datetime.now() - start_time).total_seconds()

        print(f"[{datetime.now().strftime('%H:%M:%S')}] PPTX 保存成功!")
        print(f"  路径: {os.path.abspath(output_path)}")
        print(f"  大小: {file_size} bytes ({file_size / 1024:.1f} KB)")
        print(f"  保存耗时: {save_time:.2f} 秒")
        print(f"  总渲染耗时: {total_time:.2f} 秒")

        return output_path

    def _render_slide_content(self, pptx_slide, slide_data, slide_idx, image_map, color_scheme):
        """【增强版】将 slide 数据渲染到 python-pptx 的 slide 对象"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        has_real_image = image_map and slide_idx in image_map

        title_color = RGBColor.from_string(color_scheme['primary'])
        text_color = RGBColor.from_string(color_scheme['text'])

        img_layout = slide_data.image_layout
        img_pos = img_layout.get("position", "right")

        print(
            f"    页 {slide_idx + 1} 渲染: type={slide_data.type.value}, has_real_image={has_real_image}, img_pos={img_pos}")

        if slide_data.type == SlideType.TITLE:
            # 封面页：支持背景图
            if has_real_image and img_pos == "background":
                try:
                    img_path = image_map[slide_idx]
                    bg_shape = pptx_slide.shapes.add_picture(
                        img_path,
                        Inches(0), Inches(0),
                        width=Inches(13.333)
                    )
                    # 将背景移到最底层
                    spTree = pptx_slide.shapes._spTree
                    sp = bg_shape._element
                    spTree.remove(sp)
                    spTree.insert(2, sp)
                    print(f"    已添加背景图片")
                except Exception as e:
                    print(f"    背景图片添加失败: {e}")

            # 标题
            title_box = pptx_slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(11.333), Inches(2)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.alignment = PP_ALIGN.CENTER

            if slide_data.subtitle:
                sub_box = pptx_slide.shapes.add_textbox(
                    Inches(1), Inches(4.5), Inches(11.333), Inches(1)
                )
                sub_box.text_frame.text = slide_data.subtitle
                sub_para = sub_box.text_frame.paragraphs[0]
                sub_para.font.size = Pt(24)
                sub_para.font.color.rgb = text_color
                sub_para.alignment = PP_ALIGN.CENTER

        elif slide_data.type == SlideType.CHAPTER:
            # 章节页：左侧色块装饰
            accent_color = RGBColor.from_string(color_scheme['accent'])

            # 添加左侧装饰条
            left_bar = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(0.8), Inches(7.5)
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = accent_color
            left_bar.line.fill.background()

            title_box = pptx_slide.shapes.add_textbox(
                Inches(1.5), Inches(3), Inches(10.833), Inches(1.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.alignment = PP_ALIGN.LEFT

            if slide_data.subtitle:
                sub_box = pptx_slide.shapes.add_textbox(
                    Inches(1.5), Inches(4.5), Inches(10.833), Inches(1)
                )
                sub_box.text_frame.text = slide_data.subtitle
                sub_para = sub_box.text_frame.paragraphs[0]
                sub_para.font.size = Pt(20)
                sub_para.font.color.rgb = text_color
                sub_para.alignment = PP_ALIGN.LEFT

        elif slide_data.type == SlideType.BULLET:
            # 要点页
            title_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(12), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = title_color

            content_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(12), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, point in enumerate(slide_data.content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
                p.space_after = Pt(14)
                p.level = 0

        elif slide_data.type in [SlideType.CONTENT, SlideType.IMAGE, SlideType.CHART]:
            # 内容页：智能图文布局

            # 根据图片位置调整文本框
            if img_pos == "right":
                text_left = 0.5
                text_width = 7.0
                img_left = 8.0
                img_width = 4.5
            elif img_pos == "left":
                text_left = 6.0
                text_width = 7.0
                img_left = 0.5
                img_width = 5.0
            elif img_pos == "top":
                text_left = 0.5
                text_width = 12.0
                img_left = 1.0
                img_width = 11.333
            else:  # default right
                text_left = 0.5
                text_width = 7.0
                img_left = 8.0
                img_width = 4.5

            # 标题
            title_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(12), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = title_color

            # 内容文本框
            content_box = pptx_slide.shapes.add_textbox(
                Inches(text_left), Inches(1.3), Inches(text_width), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, para in enumerate(slide_data.content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = para
                p.font.size = Pt(16)
                p.font.color.rgb = text_color
                p.space_after = Pt(12)

            # 插入真实图片或占位符
            if slide_data.type == SlideType.IMAGE:
                if has_real_image:
                    img_path = image_map[slide_idx]
                    try:
                        print(f"    插入真实图片: {os.path.basename(img_path)}")

                        # 计算图片尺寸（保持比例）
                        from PIL import Image as PILImage
                        with PILImage.open(img_path) as img:
                            orig_width, orig_height = img.size
                            aspect_ratio = orig_height / orig_width

                        # 根据布局计算实际尺寸
                        max_width = Inches(img_width)
                        max_height = Inches(5.5)

                        calc_height = max_width * aspect_ratio
                        if calc_height > max_height:
                            calc_height = max_height
                            calc_width = max_height / aspect_ratio
                        else:
                            calc_width = max_width

                        # 垂直居中
                        top_offset = Inches(1.3) + (Inches(5.5) - calc_height) / 2

                        pptx_slide.shapes.add_picture(
                            img_path,
                            Inches(img_left), top_offset,
                            width=calc_width, height=calc_height
                        )
                        print(f"    图片插入成功: {calc_width.inches:.2f}x{calc_height.inches:.2f} inches")
                    except Exception as e:
                        print(f"    插入图片失败: {e}，将使用占位符")
                        self._add_image_placeholder(pptx_slide, Inches(img_left), Inches(1.3),
                                                    Inches(img_width), Inches(5.5), slide_data.image_keywords)
                else:
                    print(f"    无可用图片，创建占位符")
                    self._add_image_placeholder(pptx_slide, Inches(img_left), Inches(1.3),
                                                Inches(img_width), Inches(5.5), slide_data.image_keywords)

        elif slide_data.type == SlideType.END:
            # 结束页
            title_box = pptx_slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(11.333), Inches(1.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.alignment = PP_ALIGN.CENTER

            if slide_data.subtitle:
                sub_box = pptx_slide.shapes.add_textbox(
                    Inches(1), Inches(4.5), Inches(11.333), Inches(1)
                )
                sub_box.text_frame.text = slide_data.subtitle
                sub_para = sub_box.text_frame.paragraphs[0]
                sub_para.font.size = Pt(20)
                sub_para.font.color.rgb = text_color
                sub_para.alignment = PP_ALIGN.CENTER

            if slide_data.content:
                content_box = pptx_slide.shapes.add_textbox(
                    Inches(1), Inches(5.5), Inches(11.333), Inches(1.5)
                )
                tf = content_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"核心要点：{slide_data.content[0]}"
                p.font.size = Pt(16)
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.CENTER

    def _add_image_placeholder(self, pptx_slide, left, top, width, height, keywords):
        """【增强版】添加图片占位符"""
        from pptx.util import Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        placeholder = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(1)

        text_frame = placeholder.text_frame
        text_frame.word_wrap = True
        keywords_str = ', '.join(keywords[:3]) if keywords else "无"

        p1 = text_frame.paragraphs[0]
        p1.text = "[图片占位符]"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(128, 128, 128)
        p1.alignment = 1  # CENTER

        p2 = text_frame.add_paragraph()
        p2.text = f"搜索关键词:"
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(150, 150, 150)
        p2.alignment = 1
        p2.space_before = Pt(6)

        p3 = text_frame.add_paragraph()
        p3.text = keywords_str
        p3.font.size = Pt(9)
        p3.font.color.rgb = RGBColor(100, 100, 100)
        p3.alignment = 1

    def _remove_all_placeholders(self, prs):
        """
        【关键修复】删除所有空白占位符
        遍历所有slide，删除包含"[图片占位符]"文本的形状
        使用 sp.getparent().remove(sp) 方法
        """
        removed_count = 0

        for slide_idx, slide in enumerate(prs.slides):
            shapes_to_remove = []

            for shape in slide.shapes:
                # 检查是否是占位符（通过文本内容判断）
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    # 如果包含"[图片占位符]"或"[图片区域]"标记
                    if "[图片占位符]" in text or "[图片区域]" in text:
                        shapes_to_remove.append(shape)
                        print(f"    标记删除占位符 (页 {slide_idx + 1}): {text[:30]}...")

            # 执行删除（从后往前删避免索引问题）
            for shape in reversed(shapes_to_remove):
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    removed_count += 1
                except Exception as e:
                    print(f"    删除占位符失败: {e}")

        return removed_count

    def _extract_json(self, text: str) -> Any:
        """【增强版】健壮性 JSON 提取"""
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 解析 JSON，文本长度: {len(text)}")

        # 打印原始内容前 500 字符用于调试
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 原始内容预览: {text[:500]}...")

        # 预处理：去掉 // 注释
        def remove_comments(s):
            lines = s.split('\n')
            cleaned = []
            for line in lines:
                # 去掉 // 注释
                if '//' in line:
                    line = line[:line.index('//')]
                cleaned.append(line)
            return '\n'.join(cleaned)

        # 尝试直接解析
        try:
            result = json.loads(text)
            print(f"[{datetime.now().strftime('%H:%M:%S')}] JSON 直接解析成功")
            return result
        except json.JSONDecodeError:
            pass

        # 预处理后解析
        cleaned_text = remove_comments(text)
        try:
            result = json.loads(cleaned_text)
            print(f"[{datetime.now().strftime('%H:%M:%S')}] 清理注释后解析成功")
            return result
        except json.JSONDecodeError:
            pass

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 直接解析失败，尝试提取代码块...")

        # 尝试提取 ```json 代码块
        matches = re.findall(r'```json\s*(.*?)\s*```', text, re.DOTALL)
        for match in matches:
            json_str = remove_comments(match)
            try:
                result = json.loads(json_str)
                print(f"[{datetime.now().strftime('%H:%M:%S')}] 从 ```json 代码块解析成功")
                return result
            except json.JSONDecodeError:
                continue

        # 尝试提取 ``` 代码块（无语言标记）
        matches = re.findall(r'```\s*(\{.*?\})\s*```', text, re.DOTALL)
        for match in matches:
            json_str = remove_comments(match)
            try:
                result = json.loads(json_str)
                print(f"[{datetime.now().strftime('%H:%M:%S')}] 从 ``` 代码块解析成功")
                return result
            except json.JSONDecodeError:
                continue

        # 尝试匹配 { ... } 或 [ ... ]
        matches = re.findall(r'(\{[^{}]*\})', text, re.DOTALL)
        for match in matches:
            json_str = remove_comments(match)
            try:
                result = json.loads(json_str)
                print(f"[{datetime.now().strftime('%H:%M:%S')}] 从 {{...}} 匹配解析成功")
                return result
            except json.JSONDecodeError:
                continue

        # 尝试匹配 JSON 数组
        match = re.search(r'(\[.*\])', text, re.DOTALL)
        if match:
            json_str = remove_comments(match.group(1))
            try:
                result = json.loads(json_str)
                print(f"[{datetime.now().strftime('%H:%M:%S')}] 从 [...] 匹配解析成功")
                return result
            except json.JSONDecodeError:
                pass

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 所有解析方式失败")
        raise ValueError(f"无法解析 JSON，原始内容: {text[:500]}")


# ==================== 使用示例 ====================

async def demo():
    print(f"\n{'#' * 60}")
    print(f"# PPT 生成器演示开始（美化增强版）")
    print(f"{'#' * 60}\n")

    gen = PptGenerator()

    # 生成完整 PPT
    print(f"[{datetime.now().strftime('%H:%M:%S')}] 步骤1: 生成 PPT 内容（含美化建议）...")
    ppt = await gen.build_presentation(
        topic="校园安全",
        pages=8,
        style="academic",
        audience="老师",
        author="Dr. AI"
    )

    # 打印美化方案预览
    print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 美化方案预览:")
    for i, slide in enumerate(ppt.slides[:3], 1):
        print(f"  页 {i}: {slide.title[:20]}...")
        if slide.design_suggestions:
            print(f"    配色: {slide.design_suggestions.get('color_usage', {}).get('title', '默认')}")
            print(f"    布局: {slide.image_layout.get('position', '默认')}")

    # 准备图片
    print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 步骤2: 准备图片资源...")
    image_map = await gen.prepare_images(ppt)
    print(f"[{datetime.now().strftime('%H:%M:%S')}] 图片准备结果: {len(image_map)} 张可用")

    # 渲染为 PPTX 文件（自动清理占位符）
    print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 步骤3: 渲染 PPTX 文件（智能图片插入+清理占位符）...")
    output_dir = os.path.join(os.path.dirname(__file__), "..", "..", "output")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "ai_medical.pptx")

    result_path = gen.render_to_pptx(ppt, output_path, image_map, remove_placeholders=True)

    print(f"\n{'=' * 60}")
    print(f"[{datetime.now().strftime('%H:%M:%S')}] 演示完成!")
    print(f"最终文件: {os.path.abspath(result_path)}")
    print(f"{'=' * 60}\n")

    return ppt


# 流式生成示例（WebSocket 场景）
async def demo_stream():
    print(f"\n{'#' * 60}")
    print(f"# 流式生成演示开始（美化增强版）")
    print(f"{'#' * 60}\n")

    gen = PptGenerator()

    async for update in gen.build_presentation_stream(
            topic="2026年技术趋势报告",
            pages=12,
            style="business"
    ):
        if update["stage"] == "page_ready":
            preview = update['preview']
            design_info = f" | 设计: {', '.join(preview.get('design_summary', []))}" if preview.get(
                'design_summary') else ""
            print(f"[STREAM] 第 {update['page_number']} 页完成: {preview['title']}{design_info}")
        elif update["stage"] == "complete":
            print(f"\n[{datetime.now().strftime('%H:%M:%S')}] 流式生成完成，准备图片...")
            ppt = update["presentation"]
            image_map = await gen.prepare_images(ppt)

            output_dir = os.path.join(os.path.dirname(__file__), "..", "..", "output")
            os.makedirs(output_dir, exist_ok=True)
            out_file = os.path.join(output_dir, "tech_trends_2026.pptx")

            gen.render_to_pptx(ppt, out_file, image_map, remove_placeholders=True)
            print(f"\n流式演示完成: {os.path.abspath(out_file)}")


# 运行
if __name__ == "__main__":
    asyncio.run(demo())