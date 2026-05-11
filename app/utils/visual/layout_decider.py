"""
智能布局决策器 - 根据视觉决策生成 python-pptx 渲染指令

功能：
1. 根据视觉规划生成 python-pptx 渲染代码
2. 处理图片布局计算
3. 文字与图片的相对位置
4. 装饰元素的添加
"""
import logging
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from enum import Enum

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.utils.visual.visual_analyzer import (
    SlideVisualDecision, 
    PPTVisualPlan, 
    ImageType, 
    ImagePosition,
    VisualAnalyzer
)
from app.utils.visual.image_manager import ImageManager, ImageAsset

logger = logging.getLogger(__name__)


class LayoutType(Enum):
    """布局类型"""
    TITLE_SLIDE = "title_slide"       # 标题页
    CONTENT_WITH_IMAGE = "content_with_image"  # 带图片的内容页
    CONTENT_ONLY = "content_only"      # 纯文字内容页
    TWO_COLUMN = "two_column"          # 双栏布局
    CENTER_FOCUS = "center_focus"      # 中心聚焦


@dataclass
class LayoutElement:
    """布局元素"""
    element_type: str  # shape, textbox, image, line
    left: Inches
    top: Inches
    width: Inches
    height: Inches
    properties: Dict[str, Any] = None
    
    def __post_init__(self):
        if self.properties is None:
            self.properties = {}


@dataclass
class SlideLayoutPlan:
    """幻灯片布局计划"""
    slide_index: int
    layout_type: LayoutType
    
    # 元素列表
    elements: List[LayoutElement] = None
    
    # 背景样式
    background_color: Tuple[int, int, int] = (255, 255, 255)  # RGB
    use_gradient: bool = False
    
    # 装饰设置
    add_header_bar: bool = True
    add_footer_bar: bool = True
    add_decorations: bool = True
    decoration_style: str = "minimal"  # minimal, rich, corporate
    
    # 页码
    page_number: int = 1
    total_pages: int = 1
    
    def __post_init__(self):
        if self.elements is None:
            self.elements = []


class LayoutDecider:
    """智能布局决策器"""
    
    # 幻灯片尺寸（16:9）- 使用 float 存储
    SLIDE_WIDTH_INCHES = 13.333
    SLIDE_HEIGHT_INCHES = 7.5
    
    # 边距
    MARGIN_LEFT_INCHES = 0.5
    MARGIN_RIGHT_INCHES = 0.5
    MARGIN_TOP_INCHES = 0.5
    MARGIN_BOTTOM_INCHES = 0.3
    
    # 内容区域（float）
    CONTENT_WIDTH_INCHES = SLIDE_WIDTH_INCHES - MARGIN_LEFT_INCHES - MARGIN_RIGHT_INCHES
    CONTENT_HEIGHT_INCHES = SLIDE_HEIGHT_INCHES - MARGIN_TOP_INCHES - MARGIN_BOTTOM_INCHES - 0.5
    
    def __init__(self):
        self.visual_analyzer = VisualAnalyzer()
        self.image_manager = ImageManager()
        
        # 转换为 Inches 对象供渲染使用
        self.SLIDE_WIDTH = Inches(self.SLIDE_WIDTH_INCHES)
        self.SLIDE_HEIGHT = Inches(self.SLIDE_HEIGHT_INCHES)
        self.MARGIN_LEFT = Inches(self.MARGIN_LEFT_INCHES)
        self.MARGIN_RIGHT = Inches(self.MARGIN_RIGHT_INCHES)
        self.MARGIN_TOP = Inches(self.MARGIN_TOP_INCHES)
        self.MARGIN_BOTTOM = Inches(self.MARGIN_BOTTOM_INCHES)
        self.CONTENT_WIDTH = Inches(self.CONTENT_WIDTH_INCHES)
        self.CONTENT_HEIGHT = Inches(self.CONTENT_HEIGHT_INCHES)

    def plan_slide_layout(
        self,
        slide_decision: SlideVisualDecision,
        page_number: int,
        total_pages: int
    ) -> SlideLayoutPlan:
        """
        为单页幻灯片生成布局计划

        Args:
            slide_decision: 视觉决策
            page_number: 页码
            total_pages: 总页数

        Returns:
            SlideLayoutPlan: 布局计划
        """
        elements = []

        # 根据是否有主图片决定布局类型
        main_image = slide_decision.get_main_image()
        if main_image and main_image.image_type != ImageType.NONE:
            layout_type = LayoutType.CONTENT_WITH_IMAGE
            elements = self._plan_with_image(slide_decision, main_image)
        else:
            layout_type = LayoutType.CONTENT_ONLY
            elements = self._plan_content_only(slide_decision)

        # 添加装饰图片（如果需要）
        if slide_decision.images:
            for img in slide_decision.images:
                if img.is_decoration and img.image_type != ImageType.NONE:
                    dec_elem = self._plan_decoration_image(img)
                    if dec_elem:
                        elements.append(dec_elem)

        return SlideLayoutPlan(
            slide_index=slide_decision.slide_index,
            layout_type=layout_type,
            elements=elements,
            page_number=page_number,
            total_pages=total_pages,
            add_decorations=slide_decision.add_decoration,
            decoration_style=slide_decision.decoration_style if hasattr(slide_decision, 'decoration_style') else 'minimal'
        )

    def _plan_with_image(self, decision: SlideVisualDecision, main_image: 'ImageDecision') -> List[LayoutElement]:
        """规划带图片的布局"""
        elements = []

        # 使用 float 计算来避免 Inches 运算问题
        content_width_val = self.CONTENT_WIDTH.inches
        margin_left_val = self.MARGIN_LEFT.inches
        slide_width_val = self.SLIDE_WIDTH.inches

        # 根据主图片位置决定布局
        position = main_image.position
        img_width_ratio = main_image.width_ratio

        # 初始化默认值
        text_top = Inches(1.5)

        if position == ImagePosition.LEFT:
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(4.5)
            img_left = self.MARGIN_LEFT
            img_top = Inches(1.5)
            text_left = Inches(margin_left_val + img_width.inches + 0.3)
            text_width = Inches(content_width_val - img_width.inches - 0.3)

        elif position == ImagePosition.RIGHT:
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(4.5)
            img_left = Inches(margin_left_val + content_width_val - img_width.inches)
            img_top = Inches(1.5)
            text_left = self.MARGIN_LEFT
            text_width = Inches(img_left.inches - margin_left_val - 0.3)

        elif position == ImagePosition.CENTER:
            img_width = Inches(6)
            img_height = Inches(4)
            img_left = Inches((slide_width_val - img_width.inches) / 2)
            img_top = Inches(1.2)
            text_left = self.MARGIN_LEFT
            text_width = Inches(content_width_val)

        elif position == ImagePosition.BACKGROUND:
            img_width = Inches(slide_width_val)
            img_height = Inches(7.5)
            img_left = Inches(0)
            img_top = Inches(0)
            text_left = self.MARGIN_LEFT
            text_width = Inches(content_width_val)

        elif position == ImagePosition.TOP:
            # 图片在顶部居中，文字在下方
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(3.5)
            img_left = Inches((slide_width_val - img_width.inches) / 2)
            img_top = Inches(0.6)
            text_left = self.MARGIN_LEFT
            text_width = Inches(content_width_val)
            text_top = Inches(img_top.inches + img_height.inches + 0.3)

        elif position == ImagePosition.TOP_RIGHT:
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(3.5)
            img_left = Inches(margin_left_val + content_width_val - img_width.inches)
            img_top = Inches(0.8)
            text_left = self.MARGIN_LEFT
            text_width = Inches(content_width_val)
            text_top = Inches(img_top.inches + img_height.inches + 0.3)

        elif position == ImagePosition.TOP_LEFT:
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(3.5)
            img_left = self.MARGIN_LEFT
            img_top = Inches(0.8)
            text_left = Inches(margin_left_val + img_width.inches + 0.3)
            text_width = Inches(content_width_val - img_width.inches - 0.3)
            text_top = Inches(img_top.inches + img_height.inches + 0.3)

        elif position == ImagePosition.BOTTOM:
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(3)
            img_left = Inches((slide_width_val - img_width.inches) / 2)
            img_top = Inches(4.5)
            text_left = self.MARGIN_LEFT
            text_width = Inches(content_width_val)
            text_top = Inches(1.3)

        elif position == ImagePosition.INLINE:
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(3.5)
            img_left = Inches((slide_width_val - img_width.inches) / 2)
            img_top = Inches(1.3)
            text_left = self.MARGIN_LEFT
            text_width = Inches(content_width_val)
            text_top = Inches(img_top.inches + img_height.inches + 0.3)

        else:
            # 默认右图左文
            img_width = Inches(content_width_val * img_width_ratio)
            img_height = Inches(4.5)
            img_left = Inches(margin_left_val + content_width_val - img_width.inches)
            img_top = Inches(1.5)
            text_left = self.MARGIN_LEFT
            text_width = Inches(img_left.inches - margin_left_val - 0.3)

        # 主图片元素
        if position != ImagePosition.BACKGROUND:
            elements.append(LayoutElement(
                element_type="image",
                left=img_left,
                top=img_top,
                width=img_width,
                height=img_height,
                properties={
                    "image_description": main_image.description,
                    "image_keywords": main_image.keywords,
                    "image_type": main_image.image_type.value,
                    "opacity": main_image.opacity
                }
            ))

        # 标题（使用新的 title_style）
        title_height = Inches(1)
        title_props = {
            "text": decision.title,
            "font_family": decision.title_style.font_family,
            "font_size": decision.title_style.font_size,
            "font_color": decision.title_style.font_color,
            "font_weight": decision.title_style.font_weight,
            "alignment": "left"
        }
        elements.append(LayoutElement(
            element_type="title",
            left=self.MARGIN_LEFT,
            top=Inches(0.3),
            width=Inches(content_width_val),
            height=title_height,
            properties=title_props
        ))

        # 分隔线（如果需要）
        if decision.has_separator_line:
            elements.append(LayoutElement(
                element_type="line",
                left=self.MARGIN_LEFT,
                top=Inches(1.3),
                width=Inches(content_width_val),
                height=Inches(0.03),
                properties={"color": decision.separator_color}
            ))

        # 内容区域
        if position == ImagePosition.CENTER:
            content_top = Inches(img_top.inches + img_height.inches + 0.3)
        elif position in (ImagePosition.TOP, ImagePosition.TOP_RIGHT, ImagePosition.TOP_LEFT, ImagePosition.BOTTOM, ImagePosition.INLINE):
            content_top = text_top
        else:
            content_top = Inches(1.5)
        content_height = self.CONTENT_HEIGHT.inches - content_top.inches - self.MARGIN_BOTTOM.inches

        content_props = {
            "items": decision.content_summary,
            "font_family": decision.content_style.font_family,
            "font_size": decision.content_style.font_size,
            "font_color": decision.content_style.font_color,
            "highlight_points": decision.highlight_words,
            "highlight_color": decision.highlight_color,
            "bullet_style": decision.bullet_style.style,
            "bullet_color": decision.bullet_style.color
        }
        elements.append(LayoutElement(
            element_type="content",
            left=text_left,
            top=content_top,
            width=text_width,
            height=Inches(max(content_height, 3)),
            properties=content_props
        ))

        return elements

    def _plan_decoration_image(self, img: 'ImageDecision') -> Optional[LayoutElement]:
        """规划装饰图片"""
        slide_width_val = self.SLIDE_WIDTH.inches
        slide_height_val = self.SLIDE_HEIGHT.inches

        # 计算装饰图片尺寸
        img_width = Inches(slide_width_val * img.width_ratio)
        img_height = Inches(slide_height_val * img.width_ratio)

        # 根据位置设置坐标
        if img.position == ImagePosition.CORNER:
            # 右下角装饰（默认）
            img_left = Inches(slide_width_val - img_width.inches)
            img_top = Inches(slide_height_val - img_height.inches)
        elif img.position == ImagePosition.TOP_RIGHT:
            # 右上角装饰
            img_left = Inches(slide_width_val - img_width.inches)
            img_top = Inches(0)
        elif img.position == ImagePosition.TOP_LEFT:
            # 左上角装饰
            img_left = Inches(0)
            img_top = Inches(0)
        elif img.position == ImagePosition.BOTTOM:
            # 底部居中装饰
            img_left = Inches((slide_width_val - img_width.inches) / 2)
            img_top = Inches(slide_height_val - img_height.inches)
        elif img.position == ImagePosition.TOP:
            # 顶部居中装饰
            img_left = Inches((slide_width_val - img_width.inches) / 2)
            img_top = Inches(0)
        else:
            # 默认右下角
            img_left = Inches(slide_width_val - img_width.inches)
            img_top = Inches(slide_height_val - img_height.inches)

        return LayoutElement(
            element_type="image",
            left=img_left,
            top=img_top,
            width=img_width,
            height=img_height,
            properties={
                "image_description": img.description,
                "image_keywords": img.keywords,
                "image_type": img.image_type.value,
                "opacity": img.opacity * 0.3
            }
        )

    def _plan_content_only(self, decision: SlideVisualDecision) -> List[LayoutElement]:
        """规划纯文字内容布局"""
        elements = []

        content_width_val = self.CONTENT_WIDTH.inches

        # 标题（使用新的 title_style）
        title_props = {
            "text": decision.title,
            "font_family": decision.title_style.font_family,
            "font_size": decision.title_style.font_size,
            "font_color": decision.title_style.font_color,
            "font_weight": decision.title_style.font_weight,
            "alignment": "left"
        }
        elements.append(LayoutElement(
            element_type="title",
            left=self.MARGIN_LEFT,
            top=Inches(0.3),
            width=Inches(content_width_val),
            height=Inches(1),
            properties=title_props
        ))
        
        # 分隔线
        if decision.has_separator_line:
            elements.append(LayoutElement(
                element_type="line",
                left=self.MARGIN_LEFT,
                top=Inches(1.3),
                width=Inches(content_width_val),
                height=Inches(0.03),
                properties={"color": decision.separator_color}
            ))

        # 内容（使用新的 content_style）
        content_props = {
            "items": decision.content_summary,
            "font_family": decision.content_style.font_family,
            "font_size": decision.content_style.font_size,
            "font_color": decision.content_style.font_color,
            "highlight_points": decision.highlight_words,
            "highlight_color": decision.highlight_color,
            "bullet_style": decision.bullet_style.style,
            "bullet_color": decision.bullet_style.color
        }
        elements.append(LayoutElement(
            element_type="content",
            left=self.MARGIN_LEFT,
            top=Inches(1.5),
            width=Inches(content_width_val),
            height=Inches(5.5),
            properties=content_props
        ))

        return elements
    
    def render_slide(
        self,
        prs: Presentation,
        layout_plan: SlideLayoutPlan,
        image_asset: Optional[ImageAsset] = None,
        style=None
    ) -> None:
        """
        根据布局计划渲染幻灯片
        
        Args:
            prs: Presentation 对象
            layout_plan: 布局计划
            image_asset: 图片资源（如果有）
            style: 样式对象
        """
        from app.api.v1.aiGeneratorPptx import PPTStyle
        
        if style is None:
            style = PPTStyle()
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # 设置背景
        self._add_background(slide, prs, style, layout_plan)
        
        # 添加装饰元素
        if layout_plan.add_decorations:
            self._add_decorations(slide, prs, style, layout_plan)
        
        # 渲染各个元素
        for element in layout_plan.elements:
            if element.element_type == "title":
                self._render_title(slide, element, style)
            elif element.element_type == "content":
                self._render_content(slide, element, style)
            elif element.element_type == "image" and image_asset:
                self._render_image(slide, element, image_asset, style)
            elif element.element_type == "line":
                self._render_line(slide, element, style)
        
        # 添加页码
        self._add_page_number(slide, prs, layout_plan.page_number, layout_plan.total_pages, style)
    
    def _add_background(
        self,
        slide,
        prs: Presentation,
        style,
        layout_plan: SlideLayoutPlan
    ):
        """添加背景"""
        bg_shape = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(*layout_plan.background_color)
        bg_shape.line.fill.background()
    
    def _add_decorations(
        self,
        slide,
        prs: Presentation,
        style,
        layout_plan: SlideLayoutPlan
    ):
        """添加装饰元素"""
        # 顶部装饰条
        top_bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = style.PRIMARY_COLOR
        top_bar.line.fill.background()
        
        # 左侧装饰条
        left_bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(0.1), prs.slide_height
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = style.PRIMARY_LIGHT
        left_bar.line.fill.background()
        
        # 底部装饰条
        bottom_bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(7.35),
            prs.slide_width, Inches(0.15)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = style.PRIMARY_DARK
        bottom_bar.line.fill.background()
    
    def _render_title(
        self,
        slide,
        element: LayoutElement,
        style
    ):
        """渲染标题（支持多种字体和颜色）"""
        title_box = slide.shapes.add_textbox(
            element.left,
            element.top,
            element.width,
            element.height
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = element.properties.get("text", "")

        # 获取样式属性，支持新格式
        props = element.properties
        font_family = props.get("font_family", style.FONT_MAIN)
        font_size = props.get("font_size", Pt(36))
        font_color = props.get("font_color", style.PRIMARY_COLOR)
        font_weight = props.get("font_weight", "bold")

        # 如果是 Pt 对象，直接使用
        if isinstance(font_size, Pt):
            p.font.size = font_size
        else:
            p.font.size = Pt(int(font_size))

        p.font.name = font_family
        p.font.bold = font_weight in ("bold", "medium")
        p.font.color.rgb = self._parse_color(font_color)

        # 对齐方式
        alignment = props.get("alignment", "left")
        if alignment == "center":
            p.alignment = PP_ALIGN.CENTER
        elif alignment == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

    def _render_content(
        self,
        slide,
        element: LayoutElement,
        style
    ):
        """渲染内容（支持多种字体、颜色、大小）"""
        content_box = slide.shapes.add_textbox(
            element.left,
            element.top,
            element.width,
            element.height
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        items = element.properties.get("items", "")
        highlight_points = element.properties.get("highlight_points", [])
        highlight_color = element.properties.get("highlight_color", "FF6600")

        # 获取样式属性
        props = element.properties
        font_family = props.get("font_family", style.FONT_MAIN)
        font_size = props.get("font_size", 20)
        font_color = props.get("font_color", style.TEXT_DARK)
        bullet_style = props.get("bullet_style", "circle")
        bullet_color = props.get("bullet_color", style.ACCENT_COLOR)

        if isinstance(items, str):
            items = [items]

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            item_text = str(item).strip()

            # 处理高亮词
            for point in highlight_points:
                if point in item_text:
                    # 高亮用特殊标记，实际渲染时分开处理
                    item_text = item_text.replace(point, f"『{point}』")

            # bullet 样式
            bullet_symbol = self._get_bullet_symbol(bullet_style, i)
            p.text = f"{bullet_symbol} {item_text}"

            p.font.name = font_family
            if isinstance(font_size, Pt):
                p.font.size = font_size
            else:
                p.font.size = Pt(int(font_size))
            p.font.color.rgb = self._parse_color(font_color)
            p.space_after = Pt(10)

    def _get_bullet_symbol(self, style: str, index: int) -> str:
        """获取 bullet 符号"""
        if style == "circle":
            return "●"
        elif style == "square":
            return "■"
        elif style == "arrow":
            return "→"
        elif style == "number":
            return f"{index + 1}."
        elif style == "icon":
            return "◆"
        else:
            return "●"

    def _parse_color(self, color_str: str) -> RGBColor:
        """解析颜色字符串为 RGBColor"""
        try:
            if isinstance(color_str, RGBColor):
                return color_str
            # 去掉 # 号
            color_str = color_str.lstrip("#")
            # 解析为 RGB
            if len(color_str) == 6:
                r = int(color_str[0:2], 16)
                g = int(color_str[2:4], 16)
                b = int(color_str[4:6], 16)
                return RGBColor(r, g, b)
            elif len(color_str) == 3:
                r = int(color_str[0] * 2, 16)
                g = int(color_str[1] * 2, 16)
                b = int(color_str[2] * 2, 16)
                return RGBColor(r, g, b)
        except Exception:
            pass
        # 默认返回深灰色
        return RGBColor(0x33, 0x33, 0x33)
    
    def _render_image(
        self,
        slide,
        element: LayoutElement,
        image_asset: ImageAsset,
        style
    ):
        """渲染图片"""
        if not image_asset or not image_asset.local_path:
            return
        
        try:
            # 添加图片
            slide.shapes.add_picture(
                image_asset.local_path,
                element.left,
                element.top,
                element.width,
                element.height
            )
        except Exception as e:
            logger.error(f"图片渲染失败: {str(e)}")
    
    def _render_line(
        self,
        slide,
        element: LayoutElement,
        style
    ):
        """渲染分隔线"""
        line = slide.shapes.add_shape(
            1,
            element.left,
            element.top,
            element.width,
            element.height
        )
        line.fill.solid()
        color = element.properties.get("color", (74, 144, 217))
        line.fill.fore_color.rgb = RGBColor(*color)
        line.line.fill.background()
    
    def _add_page_number(
        self,
        slide,
        prs: Presentation,
        page_num: int,
        total_pages: int,
        style
    ):
        """添加页码"""
        # 页码框
        num_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(1),
            prs.slide_height - Inches(0.5),
            Inches(0.8),
            Inches(0.3)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{page_num} / {total_pages}"
        p.font.name = style.FONT_MAIN
        p.font.size = Pt(12)
        p.font.color.rgb = style.TEXT_GRAY
        p.alignment = PP_ALIGN.RIGHT


# 全局实例
layout_decider = LayoutDecider()