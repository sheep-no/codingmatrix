"""
PPTX 预览转换器

将 PPTX 文件转换为 HTML 预览格式，支持：
- 提取幻灯片文本内容
- 生成缩略图导航
- 全屏预览模式
- 响应式布局
"""

import logging
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Dict, Any

from lxml import etree
from pptx import Presentation

logger = logging.getLogger(__name__)

# PPTML 命名空间
NS_MAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


@dataclass
class SlidePreview:
    """单页幻灯片预览数据"""
    index: int
    title: str
    content: List[str] = field(default_factory=list)
    notes: str = ""
    has_image: bool = False
    slide_type: str = "content"  # title, chapter, content, bullet, image, end


class PreviewConverterError(Exception):
    """预览转换异常"""
    pass


class PreviewConverter:
    """
    PPTX 预览转换器
    
    将 PPTX 文件转换为可交互的 HTML 预览页面。
    """
    
    def __init__(self, output_dir: Optional[Path] = None):
        """
        初始化预览转换器
        
        Args:
            output_dir: 输出目录，默认使用临时目录
        """
        self._output_dir = output_dir or Path("./tmp/ppt-preview")
        self._output_dir.mkdir(parents=True, exist_ok=True)
    
    async def convert(
        self,
        pptx_path: str,
        output_path: Optional[str] = None,
        include_notes: bool = False,
    ) -> str:
        """
        将 PPTX 转换为 HTML 预览
        
        Args:
            pptx_path: PPTX 文件路径
            output_path: 输出 HTML 路径，不传则自动生成
            include_notes: 是否包含演讲者备注
            
        Returns:
            生成的 HTML 文件路径
        """
        pptx = Path(pptx_path)
        
        if not pptx.exists():
            raise PreviewConverterError(f"PPTX 文件不存在：{pptx_path}")
        
        # 生成输出路径
        if output_path is None:
            output_path = str(self._output_dir / f"{pptx.stem}-preview.html")
        
        # 提取幻灯片数据
        slides_data = self._extract_slides(str(pptx), include_notes)
        
        # 生成 HTML
        html_content = self._generate_html(slides_data)
        
        # 写入文件
        Path(output_path).write_text(html_content, encoding='utf-8')
        
        logger.info(f"PPTX 预览生成成功 | input={pptx_path} | output={output_path}")
        
        return output_path
    
    def _extract_slides(self, pptx_path: str, include_notes: bool) -> List[SlidePreview]:
        """
        提取幻灯片内容
        
        Args:
            pptx_path: PPTX 文件路径
            include_notes: 是否包含备注
            
        Returns:
            幻灯片预览列表
        """
        prs = Presentation(pptx_path)
        slides_data = []
        
        for idx, slide in enumerate(prs.slides):
            slide_data = self._parse_slide(slide, idx + 1, include_notes)
            slides_data.append(slide_data)
        
        logger.info(f"提取幻灯片 | count={len(slides_data)}")
        
        return slides_data
    
    def _parse_slide(self, slide, slide_index: int, include_notes: bool) -> SlidePreview:
        """
        解析单页幻灯片
        
        Args:
            slide: python-pptx Slide 对象
            slide_index: 幻灯片索引（从 1 开始）
            include_notes: 是否包含备注
            
        Returns:
            幻灯片预览数据
        """
        title = ""
        content = []
        has_image = False
        slide_type = "content"
        
        # 提取标题
        title_shape = getattr(slide.shapes, 'title', None)
        if title_shape:
            title = title_shape.text_frame.text.strip()
        
        # 提取内容
        title_shape_id = id(title_shape) if title_shape else None
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                
                # 跳过标题
                if id(shape) == title_shape_id:
                    continue
                
                if text:
                    content.append(text)
            
            # 检查是否为图片类型（shape_type 13 = 图片）
            if getattr(shape, 'shape_type', None) == 13:
                has_image = True
        
        # 判断幻灯片类型
        slide_type = self._classify_slide(title, content, has_image)
        
        # 提取备注
        notes = ""
        if include_notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes = notes_slide.notes_text_frame.text.strip()
        
        return SlidePreview(
            index=slide_index,
            title=title,
            content=content,
            notes=notes,
            has_image=has_image,
            slide_type=slide_type,
        )
    
    def _classify_slide(
        self,
        title: str,
        content: List[str],
        has_image: bool,
    ) -> str:
        """
        分类幻灯片类型
        
        Args:
            title: 标题
            content: 内容列表
            has_image: 是否包含图片
            
        Returns:
            幻灯片类型
        """
        if not title and not content:
            return "blank"
        
        if not content and not has_image:
            if title:
                if len(title) > 15:
                    return "chapter"
                return "title"
            return "blank"
        
        if has_image:
            return "image"
        
        if len(content) > 5:
            return "bullet"
        
        return "content"
    
    def _generate_html(self, slides: List[SlidePreview]) -> str:
        """
        生成 HTML 预览页面
        
        Args:
            slides: 幻灯片预览数据列表
            
        Returns:
            HTML 内容
        """
        # 生成缩略图导航
        thumbnails_html = self._generate_thumbnails(slides)
        
        # 生成幻灯片内容
        slides_html = self._generate_slides(slides)
        
        # 完整 HTML 模板
        html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PPT 预览</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
            background: #1a1a1a;
            color: #fff;
            height: 100vh;
            overflow: hidden;
        }}
        
        .container {{
            display: flex;
            height: 100vh;
        }}
        
        /* 侧边栏缩略图 */
        .sidebar {{
            width: 200px;
            background: #2d2d2d;
            overflow-y: auto;
            border-right: 1px solid #444;
        }}
        
        .thumbnail {{
            width: 160px;
            height: 100px;
            margin: 10px;
            background: #fff;
            border: 2px solid transparent;
            border-radius: 4px;
            cursor: pointer;
            overflow: hidden;
            position: relative;
            transition: all 0.2s;
        }}
        
        .thumbnail:hover {{
            border-color: #409eff;
        }}
        
        .thumbnail.active {{
            border-color: #409eff;
            box-shadow: 0 0 10px rgba(64, 158, 255, 0.5);
        }}
        
        .thumbnail-number {{
            position: absolute;
            top: 4px;
            left: 6px;
            background: rgba(0,0,0,0.7);
            color: #fff;
            font-size: 12px;
            padding: 2px 6px;
            border-radius: 2px;
        }}
        
        .thumbnail-content {{
            padding: 20px 8px 8px;
            font-size: 8px;
            color: #333;
        }}
        
        .thumbnail-title {{
            font-weight: bold;
            font-size: 10px;
            margin-bottom: 4px;
        }}
        
        /* 主内容区 */
        .main-content {{
            flex: 1;
            display: flex;
            flex-direction: column;
        }}
        
        /* 工具栏 */
        .toolbar {{
            height: 50px;
            background: #2d2d2d;
            display: flex;
            align-items: center;
            padding: 0 20px;
            border-bottom: 1px solid #444;
        }}
        
        .toolbar h2 {{
            flex: 1;
            font-size: 16px;
            font-weight: normal;
        }}
        
        .toolbar button {{
            background: #409eff;
            color: #fff;
            border: none;
            padding: 8px 16px;
            border-radius: 4px;
            cursor: pointer;
            margin-left: 10px;
            transition: background 0.2s;
        }}
        
        .toolbar button:hover {{
            background: #66b1ff;
        }}
        
        /* 幻灯片显示区 */
        .slide-viewer {{
            flex: 1;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 30px;
            background: #1a1a1a;
        }}
        
        .slide-display {{
            width: 100%;
            max-width: 960px;
            aspect-ratio: 16/9;
            background: #fff;
            border-radius: 8px;
            box-shadow: 0 10px 40px rgba(0,0,0,0.5);
            padding: 60px;
            overflow: hidden;
            position: relative;
        }}
        
        .slide-title {{
            font-size: 32px;
            color: #333;
            margin-bottom: 30px;
            font-weight: bold;
        }}
        
        .slide-note {{
            /* 备注信息样式 */
        }}
        
        .slide-bullet {{
            font-size: 20px;
            color: #666;
            margin-left: 30px;
            margin-bottom: 15px;
            line-height: 1.5;
        }}
        
        /* 全屏模式 */
        .fullscreen .sidebar,
        .fullscreen .toolbar {{
            display: none;
        }}
        
        .fullscreen .slide-viewer {{
            padding: 0;
        }}
        
        .fullscreen .slide-display {{
            width: 100vw;
            height: 100vh;
            max-width: 100%;
            border-radius: 0;
            box-shadow: none;
        }}
        
        /* 导航按钮 */
        .nav-buttons {{
            position: absolute;
            bottom: 20px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 10px;
        }}
        
        .nav-btn {{
            width: 40px;
            height: 40px;
            border-radius: 50%;
            background: rgba(0,0,0,0.5);
            color: #fff;
            border: none;
            cursor: pointer;
            font-size: 18px;
            display: flex;
            align-items: center;
            justify-content: center;
        }}

        .nav-btn:hover {{
            background: rgba(0,0,0,0.8);
        }}

        .nav-btn:disabled {{
            opacity: 0.3;
            cursor: not-allowed;
        }}
    </style>
</head>
<body>
    <div class="container">
        <!-- 缩略图侧边栏 -->
        <div class="sidebar">
            {thumbnails_html}
        </div>
        
        <!-- 主内容区 -->
        <div class="main-content">
            <!-- 工具栏 -->
            <div class="toolbar">
                <h2 id="current-title">幻灯片 1 / {len(slides)}</h2>
                <button id="fullscreen-btn" onclick="toggleFullscreen()" class="fullscreen-button">全屏</button>
            </div>
            
            <!-- 幻灯片显示区 -->
            <div class="slide-viewer">
                <div class="slide-display" id="slide-display">
                    {slides_html}
                </div>
                
                <!-- 导航按钮 -->
                <div class="nav-buttons">
                    <button class="nav-btn" id="prev-btn" onclick="navigateSlide(-1)" disabled>◀</button>
                    <button class="nav-btn" id="next-btn" onclick="navigateSlide(1)">▶</button>
                </div>
            </div>
        </div>
    </div>
    
    <script>
        // 幻灯片数据
        const slides = {self._generate_slide_json(slides)};
        
        let currentSlideIndex = 0;
        
        function navigateSlide(direction) {{
            const newIndex = currentSlideIndex + direction;
            
            if (newIndex >= 0 && newIndex < slides.length) {{
                currentSlideIndex = newIndex;
                displaySlide(currentSlideIndex);
            }}
        }}
        
        function displaySlide(index) {{
            currentSlideIndex = index;
            const slide = slides[index];
            
            // 更新标题
            document.getElementById('current-title').textContent = 
                `幻灯片 ${{index + 1}} / ${{slides.length}}: ${{slide.title}}`;
            
            // 更新显示内容
            updateSlideContent(slide);
            
            // 更新导航按钮状态
            document.getElementById('prev-btn').disabled = (index === 0);
            document.getElementById('next-btn').disabled = (index === slides.length - 1);
            
            // 更新缩略图高亮
            document.querySelectorAll('.thumbnail').forEach((thumb, i) => {{
                thumb.classList.toggle('active', i === index);
            }});
        }}
        
        function updateSlideContent(slide) {{
            const display = document.getElementById('slide-display');
            
            let html = '';
            
            if (slide.title) {{
                html += `<h1 class="slide-title">${{slide.title}}</h1>`;
            }}
            
            if (slide.content && slide.content.length > 0) {{
                for (const item of slide.content) {{
                    html += `<p class="slide-bullet">${{item}}</p>`;
                }}
            }}
            
            if (slide.notes && slide.notes.length > 0) {{
                html += `<p class="slide-note"><strong>备注：</strong>${{slide.notes}}</p>`;
            }}
            
            display.innerHTML = html;
        }}
        
        function toggleFullscreen() {{
            if (!document.fullscreenElement) {{
                document.documentElement.requestFullscreen();
            }} else {{
                document.exitFullscreen();
            }}
        }}
        
        // 键盘导航
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowLeft') {{
                navigateSlide(-1);
            }} else if (e.key === 'ArrowRight') {{
                navigateSlide(1);
            }} else if (e.key === 'Escape') {{
                if (document.fullscreenElement) {{
                    document.exitFullscreen();
                }}
            }} else if (e.key === 'f' || e.key === 'F') {{
                toggleFullscreen();
            }}
        }});
        
        // 初始化
        displaySlide(0);
    </script>
</body>
</html>"""
        
        return html
    
    def _generate_thumbnails(self, slides: List[SlidePreview]) -> str:
        """生成缩略图 HTML"""
        thumbnails = []
        
        for slide in slides:
            title_preview = slide.title[:30] + ("..." if len(slide.title) > 30 else "")
            content_preview = " ".join(slide.content[:2])[:60] + ("..." if len(" ".join(slide.content)) > 60 else "")
            
            thumbnail = f"""
            <div class="thumbnail" id="thumb-{slide.index}" onclick="displaySlide({slide.index - 1})">
                <span class="thumbnail-number">{slide.index}</span>
                <div class="thumbnail-content">
                    <div class="thumbnail-title">{title_preview}</div>
                    <div>{content_preview}</div>
                </div>
            </div>
            """
            thumbnails.append(thumbnail)
        
        return "\n".join(thumbnails)
    
    def _generate_slides(self, slides: List[SlidePreview]) -> str:
        """生成首张幻灯片的初始显示内容（后续由 JS 动态更新）"""
        if not slides:
            return "<p>无幻灯片</p>"
        
        slide = slides[0]
        content_html = ""
        
        if slide.title:
            content_html += f"<h1 class=\"slide-title\">{slide.title}</h1>"
        
        if slide.content:
            for item in slide.content:
                content_html += f"<p class=\"slide-bullet\">{item}</p>"
        
        return content_html
    
    def _generate_slide_json(self, slides: List[SlidePreview]) -> str:
        """生成幻灯片数据的 JSON 表示"""
        import json
        
        data = []
        for slide in slides:
            data.append({
                "index": slide.index,
                "title": slide.title,
                "content": slide.content,
                "notes": slide.notes,
                "has_image": slide.has_image,
                "slide_type": slide.slide_type,
            })
        
        return json.dumps(data, ensure_ascii=False, indent=2)


# 全局单例
preview_converter = PreviewConverter()
