"""
PPT 统一渲染引擎

整合两套生成逻辑的核心渲染功能：
- 使用 PPTAgent 生成大纲（来自 aiGeneratorPptx.py）
- 整合 pptxGenerateUtil.py 的内容生成和字数限制逻辑
- 使用视觉决策引擎进行配图和布局
- 使用 python-pptx 进行最终渲染
"""

import asyncio
import json
import logging
import re
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from app.agent.ppt_agent import PPTAgent, PresentationOutline, SlideOutline
from app.services.ppt.config_loader import PPTConfig
from app.services.ppt.template_registry import TemplateRegistry, template_registry, TemplateConfig
from app.utils.pptx.animation_engine import AnimationEngine, AnimationPresets

logger = logging.getLogger(__name__)


@dataclass
class SlideContent:
    """单页幻灯片内容（统一数据结构）"""
    type: str  # title, chapter, content, bullet, image, chart, end
    title: str
    bullets: List[str] = field(default_factory=list)
    subtitle: str = ""
    notes: str = ""
    image_keywords: List[str] = field(default_factory=list)
    image_url: Optional[str] = None
    layout_type: str = "default"
    design_suggestions: Dict[str, Any] = field(default_factory=dict)
    max_chars: int = 300


@dataclass
class RenderedSlide:
    """渲染后的幻灯片信息"""
    slide_index: int
    title: str
    type: str
    has_image: bool = False
    has_animation: bool = False


class ContentGeneratorError(Exception):
    """内容生成异常"""
    pass


class RendererError(Exception):
    """渲染异常"""
    pass


class PPTXRenderer:
    """
    PPT 统一渲染引擎
    
    整合大纲生成、内容填充、视觉决策和 PPTX 渲染的完整流程。
    """
    
    # 单页字数限制配置（整合自 pptxGenerateUtil.py）
    MAX_CHARS_PER_SLIDE = {
        "title": 50,
        "chapter": 80,
        "bullet": 200,
        "content": 300,
        "image": 150,
        "chart": 150,
        "end": 100,
    }
    
    def __init__(
        self,
        config: Optional[PPTConfig] = None,
        template_registry: Optional[TemplateRegistry] = None,
    ):
        """
        初始化渲染引擎
        
        Args:
            config: PPT 全局配置
            template_registry: 模板注册表
        """
        self._config = config
        self._template_registry = template_registry or template_registry
        self._ppt_agent = PPTAgent()
        
        # 临时目录
        self._temp_dir = Path("./tmp/ppt/images")
        self._temp_dir.mkdir(parents=True, exist_ok=True)
    
    async def generate_complete_pptx(
        self,
        topic: str,
        template_id: str = "modern",
        slide_count: int = 10,
        output_path: Optional[Path] = None,
        api_key_token: Optional[str] = None,
    ) -> Path:
        """
        完整 PPTX 生成流程
        
        Args:
            topic: PPT 主题
            template_id: 模板 ID
            slide_count: 页数
            output_path: 输出路径，不传则自动生成
            api_key_token: API Key Token
            
        Returns:
            生成的 PPTX 文件路径
        """
        # 获取模板配置
        template = self._template_registry.get_template(template_id)
        
        # 步骤 1: 生成大纲
        logger.info(f"开始生成 PPTX | topic={topic[:50]}... | template={template_id}")
        outline = await self._ppt_agent.generate_outline(
            topic=topic,
            num_slides=slide_count,
            api_key_token=api_key_token,
        )
        
        # 步骤 2: 生成详细内容
        slides_content = await self._generate_slide_contents(
            outline=outline,
            template=template,
        )
        
        # 步骤 3: 创建 PPTX
        if output_path is None:
            output_path = Path("./pptx_output") / f"{topic[:30]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            output_path.parent.mkdir(parents=True, exist_ok=True)
        
        await self._render_pptx(
            slides_content=slides_content,
            template=template,
            output_path=output_path,
        )
        
        logger.info(f"PPTX 生成完成 | path={output_path} | slides={len(slides_content)}")
        
        return output_path
    
    async def _generate_slide_contents(
        self,
        outline: PresentationOutline,
        template: TemplateConfig,
    ) -> List[SlideContent]:
        """
        为大纲中的每页生成详细内容
        
        Args:
            outline: PPT 大纲
            template: 模板配置
            
        Returns:
            幻灯片内容列表
        """
        slides_content = []
        max_concurrent = self._config.images.max_concurrent_search if self._config else 5
        semaphore = asyncio.Semaphore(max_concurrent)
        
        async def process_slide(slide_outline: SlideOutline) -> SlideContent:
            async with semaphore:
                # 获取该类型页面的字数限制
                max_chars = self.MAX_CHARS_PER_SLIDE.get(slide_outline.type, 200)
                
                # 生成内容（简单版本，直接使用大纲内容）
                # TODO: 可以接入 AI 生成详细内容
                bullets = slide_outline.bullets if slide_outline.bullets else []
                
                # 截断超出字数限制的内容
                truncated_bullets = self._truncate_content(
                    title=slide_outline.title,
                    bullets=bullets,
                    max_chars=max_chars,
                )
                
                return SlideContent(
                    type=slide_outline.type,
                    title=slide_outline.title,
                    bullets=truncated_bullets,
                    subtitle="",
                    notes=slide_outline.notes,
                    image_keywords=slide_outline.image_keywords,
                    layout_type=self._decide_layout(slide_outline.type),
                    max_chars=max_chars,
                )
        
        # 并发处理所有幻灯片
        tasks = [process_slide(slide) for slide in outline.slides]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # 处理结果
        for i, result in enumerate(results):
            if isinstance(result, Exception):
                logger.error(f"幻灯片 {i+1} 生成失败: {result}")
                # 创建默认内容
                slides_content.append(self._create_default_slide(outline.slides[i]))
            else:
                slides_content.append(result)
        
        return slides_content
    
    def _truncate_content(
        self,
        title: str,
        bullets: List[str],
        max_chars: int,
    ) -> List[str]:
        """
        截断超出字数限制的内容
        
        Args:
            title: 标题
            bullets: 要点列表
            max_chars: 最大字数
            
        Returns:
            截断后的要点列表
        """
        # 计算当前字数
        current_chars = len(title)
        
        truncated = []
        for bullet in bullets:
            bullet_len = len(bullet)
            if current_chars + bullet_len > max_chars - 10:  # 留 10 字余量
                # 截断最后一个要点
                remaining = max_chars - current_chars - 10
                if remaining > 10:
                    truncated.append(bullet[:remaining] + "...")
                break
            truncated.append(bullet)
            current_chars += bullet_len
        
        return truncated
    
    def _create_default_slide(self, slide_outline: SlideOutline) -> SlideContent:
        """创建默认幻灯片内容（当 AI 生成失败时使用）"""
        slide_type = slide_outline.type
        max_chars = self.MAX_CHARS_PER_SLIDE.get(slide_type, 200)
        
        # 根据类型生成默认内容
        if slide_type == "title":
            bullets = []
        elif slide_type == "bullet":
            bullets = ["要点一", "要点二", "要点三"]
        elif slide_type == "chapter":
            bullets = []
        elif slide_type == "end":
            bullets = ["感谢观看"]
        else:
            bullets = ["内容待补充"]
        
        return SlideContent(
            type=slide_type,
            title=slide_outline.title[:15],
            bullets=bullets,
            subtitle="",
            notes="",
            image_keywords=slide_outline.image_keywords,
            layout_type=self._decide_layout(slide_type),
            max_chars=max_chars,
        )
    
    def _decide_layout(self, slide_type: str) -> str:
        """根据幻灯片类型决定布局"""
        layout_map = {
            "title": "center_title",
            "chapter": "chapter_divider",
            "content": "content_with_image",
            "bullet": "bullet_list",
            "image": "full_image",
            "chart": "data_chart",
            "end": "center_title",
        }
        return layout_map.get(slide_type, "content_with_image")
    
    async def _render_pptx(
        self,
        slides_content: List[SlideContent],
        template: TemplateConfig,
        output_path: Path,
        enable_animation: bool = True,
    ) -> None:
        """
        渲染 PPTX 文件
        
        Args:
            slides_content: 幻灯片内容列表
            template: 模板配置
            output_path: 输出路径
            enable_animation: 是否启用动画
        """
        prs = Presentation()
        
        # 配置模板样式
        style = self._create_style_from_template(template)
        
        total_slides = len(slides_content)
        
        for idx, slide_content in enumerate(slides_content):
            # 创建空白幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 根据类型渲染不同布局
            await self._render_slide_by_type(
                slide=slide,
                prs=prs,
                slide_content=slide_content,
                style=style,
                slide_number=idx + 1,
                total_slides=total_slides,
            )
        
        # 应用动画（如果启用）
        if enable_animation:
            try:
                await self._apply_animations(prs, template)
                logger.info("动画应用成功")
            except Exception as e:
                logger.warning(f"动画应用失败，降级为无动画版本: {e}")
        
        # 保存文件
        prs.save(str(output_path))
    
    async def _render_slide_by_type(
        self,
        slide,
        prs,
        slide_content: SlideContent,
        style: Dict[str, Any],
        slide_number: int,
        total_slides: int,
    ) -> None:
        """根据幻灯片类型调用不同的渲染方法"""
        renderers = {
            "title": self._render_title_slide,
            "chapter": self._render_chapter_slide,
            "content": self._render_content_slide,
            "bullet": self._render_bullet_slide,
            "image": self._render_image_slide,
            "chart": self._render_chart_slide,
            "end": self._render_end_slide,
        }
        
        renderer = renderers.get(slide_content.type, self._render_content_slide)
        
        await renderer(
            slide=slide,
            prs=prs,
            slide_content=slide_content,
            style=style,
            slide_number=slide_number,
            total_slides=total_slides,
        )
    
    async def _render_title_slide(
        self, slide, prs, slide_content: SlideContent, style: Dict, slide_number: int, total_slides: int
    ) -> None:
        """渲染封面页"""
        # 背景
        self._add_background(slide, prs, style)
        
        # 主标题
        self._add_centered_title(
            slide, prs,
            title=slide_content.title,
            subtitle=slide_content.subtitle,
            style=style,
        )
    
    async def _render_chapter_slide(
        self, slide, prs, slide_content: SlideContent, style: Dict, slide_number: int, total_slides: int
    ) -> None:
        """渲染章节页"""
        self._add_background(slide, prs, style)
        
        # 章节装饰条
        self._add_section_bar(slide, prs, style)
        
        # 章节标题
        self._add_centered_title(
            slide, prs,
            title=slide_content.title,
            subtitle="",
            style=style,
        )
    
    async def _render_content_slide(
        self, slide, prs, slide_content: SlideContent, style: Dict, slide_number: int, total_slides: int
    ) -> None:
        """渲染内容页"""
        self._add_background(slide, prs, style)
        self._add_page_header(slide, prs, slide_content.title, style)
        
        # 内容区域
        if slide_content.bullets:
            self._add_bullet_list(
                slide, prs,
                bullets=slide_content.bullets,
                style=style,
                start_x=Inches(0.5),
                start_y=Inches(2),
                width=Inches(9),
                height=Inches(5),
            )
    
    async def _render_bullet_slide(
        self, slide, prs, slide_content: SlideContent, style: Dict, slide_number: int, total_slides: int
    ) -> None:
        """渲染要点页"""
        self._add_background(slide, prs, style)
        self._add_page_header(slide, prs, slide_content.title, style)
        
        # 要点列表
        self._add_bullet_list(
            slide, prs,
            bullets=slide_content.bullets,
            style=style,
            start_x=Inches(0.8),
            start_y=Inches(2),
            width=Inches(8.5),
            height=Inches(5),
        )
    
    async def _render_image_slide(
        self, slide, prs, slide_content: SlideContent, style: Dict, slide_number: int, total_slides: int
    ) -> None:
        """渲染图文页"""
        self._add_background(slide, prs, style)
        self._add_page_header(slide, prs, slide_content.title, style)
        
        # 如果有图片 URL，添加图片占位符
        if slide_content.image_url:
            # TODO: 实际下载并插入图片
            self._add_image_placeholder(slide, prs, slide_content.image_url, style)
        
        # 说明文字
        if slide_content.bullets:
            self._add_bullet_list(
                slide, prs,
                bullets=slide_content.bullets[:2],  # 最多 2 条说明
                style=style,
                start_x=Inches(0.5),
                start_y=Inches(2.5),
                width=Inches(4),
                height=Inches(4.5),
            )
    
    async def _render_chart_slide(
        self, slide, prs, slide_content: SlideContent, style: Dict, slide_number: int, total_slides: int
    ) -> None:
        """渲染图表页"""
        self._add_background(slide, prs, style)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_content.title
        p.font.size = Pt(style.get("title_font_size", 28))
        p.font.bold = True
        p.font.color.rgb = style.get("primary_color", RGBColor(0, 0, 0))
        
        # 图表占位符
        self._add_chart_placeholder(slide, prs, style)
    
    async def _render_end_slide(
        self, slide, prs, slide_content: SlideContent, style: Dict, slide_number: int, total_slides: int
    ) -> None:
        """渲染结束页"""
        self._add_background(slide, prs, style)
        
        # 居中感谢文字
        self._add_centered_title(
            slide, prs,
            title=slide_content.title,
            subtitle="",
            style=style,
        )
    
    def _create_style_from_template(self, template: TemplateConfig) -> Dict[str, Any]:
        """从模板配置创建样式字典"""
        def hex_to_rgb(hex_color: str) -> RGBColor:
            hex_color = hex_color.lstrip('#')
            return RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        
        colors = template.colors
        fonts = template.fonts
        
        primary = hex_to_rgb(colors.get("primary", "#2563eb"))
        secondary = hex_to_rgb(colors.get("secondary", "#64748b"))
        background = hex_to_rgb(colors.get("background", "#ffffff"))
        text = hex_to_rgb(colors.get("text", "#1e293b"))
        
        return {
            "primary_color": primary,
            "secondary_color": secondary,
            "background_color": background,
            "text_color": text,
            "title_font": fonts.get("title", "Arial").split(',')[0].strip(),
            "body_font": fonts.get("body", "Arial").split(',')[0].strip(),
            "title_font_size": fonts.get("title_size", 32),
            "body_font_size": fonts.get("body_size", 18),
        }
    
    def _add_background(self, slide, prs, style: Dict) -> None:
        """添加背景"""
        bg_shape = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = style.get("background_color", RGBColor(0xFF, 0xFF, 0xFF))
        bg_shape.line.fill.background()
    
    def _add_centered_title(
        self, slide, prs, title: str, subtitle: str = "", style: Dict = None
    ) -> None:
        """添加居中标题"""
        style = style or {}
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(8), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(style.get("title_font_size", 36))
        p.font.bold = True
        p.font.color.rgb = style.get("primary_color", RGBColor(0, 0, 0))
        p.font.name = style.get("title_font", "Arial")
        p.alignment = PP_ALIGN.CENTER
    
    def _add_page_header(self, slide, prs, title: str, style: Dict) -> None:
        """添加页面标题和装饰"""
        # 顶部装饰条
        header_bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.1)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = style.get("primary_color", RGBColor(0, 0, 0))
        header_bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(style.get("title_font_size", 28))
        p.font.bold = True
        p.font.color.rgb = style.get("primary_color", RGBColor(0, 0, 0))
    
    def _add_section_bar(self, slide, prs, style: Dict) -> None:
        """添加章节装饰条"""
        section_bar = slide.shapes.add_shape(
            1,
            Inches(2), Inches(3.5),
            Inches(6), Inches(0.15)
        )
        section_bar.fill.solid()
        section_bar.fill.fore_color.rgb = style.get("primary_color", RGBColor(0, 0, 0))
        section_bar.line.fill.background()
    
    def _add_bullet_list(
        self, slide, prs, bullets: List[str], style: Dict,
        start_x, start_y, width, height
    ) -> None:
        """添加要点列表"""
        content_box = slide.shapes.add_textbox(
            start_x, start_y, width, height
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        body_font_size = style.get("body_font_size", 18)
        text_color = style.get("text_color", RGBColor(0x33, 0x33, 0x33))
        body_font = style.get("body_font", "Arial")
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {bullet}"
            p.font.size = Pt(body_font_size)
            p.font.color.rgb = text_color
            p.font.name = body_font
            p.space_after = Pt(12)
    
    def _add_image_placeholder(self, slide, prs, image_url: str, style: Dict) -> None:
        """添加图片占位符"""
        # 图片区域
        img_box = slide.shapes.add_shape(
            1,
            Inches(5), Inches(2),
            Inches(4.5), Inches(4.5)
        )
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        img_box.line.color.rgb = style.get("secondary_color", RGBColor(0x99, 0x99, 0x99))
        
        # 占位符文字
        tf = img_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "图片占位"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(50)
    
    def _add_chart_placeholder(self, slide, prs, style: Dict) -> None:
        """添加图表占位符"""
        chart_box = slide.shapes.add_shape(
            1,
            Inches(1), Inches(1.5),
            Inches(8), Inches(5)
        )
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        chart_box.line.color.rgb = style.get("secondary_color", RGBColor(0x99, 0x99, 0x99))
        
        tf = chart_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "图表占位区"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(80)
    
    async def _apply_animations(self, prs, template: TemplateConfig) -> None:
        """
        应用动画到演示文稿
        
        Args:
            prs: Presentation 对象
            template: 模板配置
        """
        engine = AnimationEngine()
        
        # 获取模板的动画配置
        animations = template.animations
        transition_effect = animations.get("default", "fade")
        supported_transitions = animations.get("supported", ["fade"])
        
        # 映射过渡效果
        from app.utils.pptx.animation_engine import TransitionEffect, EntranceEffect
        
        transition_map = {
            "fade": TransitionEffect.FADE,
            "push": TransitionEffect.PUSH,
            "wipe": TransitionEffect.WIPE,
            "split": TransitionEffect.SPLIT,
            "cover": TransitionEffect.COVER,
        }
        
        # 应用页面切换动画
        effect = transition_map.get(transition_effect, TransitionEffect.FADE)
        duration = 1.0  # 默认 1 秒
        
        engine.set_default_transition(prs, effect, duration)
        
        # 应用元素进入动画
        for slide in prs.slides:
            # 标题动画
            engine.apply_title_animation(slide, delay=0.3)
            
            # 内容动画
            shapes_to_animate = [
                shape for shape in slide.shapes
                if shape.has_text_frame
                and shape.text_frame.text.strip()
                and shape != slide.shapes.title
            ]
            
            if shapes_to_animate:
                engine.apply_content_animation(slide, shapes_to_animate, stagger=0.2)
    
    def _clean_ai_response(self, text: str) -> str:
        """清理 AI 响应内容"""
        if not text:
            return text
        
        # 移除 markdown 代码块
        text = re.sub(r'```json\s*', '', text)
        text = re.sub(r'```\s*', '', text)
        
        # 移除开头的非 JSON 字符
        text = re.sub(r'^[^{[]+', '', text)
        
        # 移除结尾的非 JSON 字符
        text = re.sub(r'[^}\]]+$', '', text)
        
        return text.strip()
    
    def _extract_json(self, text: str) -> Any:
        """从文本中提取 JSON"""
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            # 尝试修复常见 JSON 错误
            text = text.replace('"""', '"').replace('""', '"')
            return json.loads(text)


# 全局单例
pptx_renderer = PPTXRenderer()
