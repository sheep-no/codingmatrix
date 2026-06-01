"""
PPT 统一编排器

负责协调 PPT 生成的完整流程：
1. 接收生成请求
2. 调用 PPTAgent 生成大纲
3. 视觉决策分析
4. 并发搜图配图
5. 智能排版
6. 渲染 PPTX
7. 应用动画效果
8. 存储文件
"""

import asyncio
import logging
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List

from app.agent.ppt_agent import PPTAgent, PresentationOutline
from app.utils.task_manager import task_manager, TaskStatus
from app.utils.visual import (
    visual_analyzer,
    image_manager,
    layout_decider,
)
from app.utils.pptx.text_processor import prevent_text_overflow
from app.utils.pptx.image_search import ImageSearchManager

logger = logging.getLogger(__name__)

# PPT 文件存储目录
PPT_OUTPUT_DIR = Path("./pptx_output")
PPT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


class PPTOrchestrationError(Exception):
    """PPT 编排异常"""
    pass


class TaskCancelledError(PPTOrchestrationError):
    """任务取消异常"""
    pass


class PPTOrchestrator:
    """
    PPT 统一编排器
    
    协调 PPT 生成的完整流程，提供统一的 API 接口。
    """
    
    def __init__(self):
        self._ppt_agent = PPTAgent()
        self._image_search = ImageSearchManager()
        self._running_tasks: Dict[str, asyncio.Task] = {}
        self._cancel_events: Dict[str, asyncio.Event] = {}
    
    async def generate(
        self,
        user_id: int,
        topic: str,
        template: str = "modern",
        slide_count: int = 10,
        output_format: str = "pptx",
        language: str = "zh-CN",
        quality: str = "high",
        api_key_token: Optional[str] = None,
        options: Optional[Dict[str, bool]] = None,
    ) -> str:
        """
        创建 PPT 生成任务
        
        Args:
            user_id: 用户 ID
            topic: PPT 主题
            template: 模板风格
            slide_count: 页数
            output_format: 输出格式 (pptx/pdf/html/markdown)
            language: 语言
            quality: 内容质量 (high/balanced/fast)
            api_key_token: API Key Token
            options: 高级选项
            
        Returns:
            task_id: 任务 ID
        """
        task_id = str(uuid.uuid4())
        cancel_event = asyncio.Event()
        
        # 创建任务记录
        await task_manager.create_task(
            task_type="ppt_generation",
            user_id=user_id,
            func=self._execute_generation,
            params={
                "task_id": task_id,
                "topic": topic,
                "template": template,
                "slide_count": slide_count,
                "output_format": output_format,
                "language": language,
                "quality": quality,
                "api_key_token": api_key_token,
                "options": options or {},
                "cancel_event": cancel_event,
            },
        )
        
        # 保存取消事件引用
        self._cancel_events[task_id] = cancel_event
        
        logger.info(f"创建 PPT 生成任务 | task_id={task_id} | topic={topic[:50]}...")
        return task_id
    
    async def get_progress(self, task_id: str) -> Optional[Dict[str, Any]]:
        """
        查询任务进度
        
        Args:
            task_id: 任务 ID
            
        Returns:
            任务信息字典，包括状态、进度、当前步骤等
        """
        task_info = await task_manager.get_task_info_async(task_id)
        
        if task_info is None:
            return None
        
        return {
            "task_id": task_id,
            "status": task_info.get("status"),
            "progress": task_info.get("progress", 0.0),
            "current_step": task_info.get("message", ""),
            "created_at": task_info.get("created_at"),
            "result": task_info.get("result"),
            "error": task_info.get("error"),
        }
    
    async def cancel(self, task_id: str) -> bool:
        """
        取消任务
        
        Args:
            task_id: 任务 ID
            
        Returns:
            是否成功取消
        """
        # 设置取消事件
        cancel_event = self._cancel_events.get(task_id)
        if cancel_event:
            cancel_event.set()
        
        # 取消 asyncio 任务
        running_task = self._running_tasks.get(task_id)
        if running_task and not running_task.done():
            running_task.cancel()
        
        # 更新任务状态
        success = await task_manager.cancel_task(task_id)
        
        if success:
            logger.info(f"取消 PPT 生成任务 | task_id={task_id}")
        
        return success
    
    async def update(
        self,
        task_id: str,
        updates: Dict[str, Any],
    ) -> Optional[Dict[str, Any]]:
        """
        增量更新任务
        
        Args:
            task_id: 任务 ID
            updates: 更新内容（如补充需求、修改模板等）
            
        Returns:
            更新后的任务信息
        """
        # TODO: 实现增量更新逻辑
        # 当前版本不支持运行中任务的增量更新
        logger.warning(f"增量更新未实现 | task_id={task_id}")
        return await self.get_progress(task_id)
    
    async def _execute_generation(
        self,
        task_id: str,
        topic: str,
        template: str,
        slide_count: int,
        output_format: str,
        language: str,
        quality: str,
        api_key_token: Optional[str],
        options: Dict[str, bool],
        cancel_event: asyncio.Event,
    ):
        """
        执行 PPT 生成流程
        
        这是实际执行生成的内部方法，由 TaskManager 调用。
        """
        try:
            # 检查取消
            self._check_cancelled(task_id, cancel_event)
            
            # 步骤 1: 生成大纲 (0-20%)
            await task_manager.update_progress(task_id, 0.05, "正在生成大纲...")
            outline = await self._generate_outline(
                task_id=task_id,
                topic=topic,
                slide_count=slide_count,
                api_key_token=api_key_token,
                cancel_event=cancel_event,
            )
            
            self._check_cancelled(task_id, cancel_event)
            
            # 步骤 2: 视觉决策分析 (20-30%)
            await task_manager.update_progress(task_id, 0.20, "正在分析视觉需求...")
            visual_plan = await self._analyze_visual_needs(
                task_id=task_id,
                outline=outline,
                template=template,
                cancel_event=cancel_event,
            )
            
            self._check_cancelled(task_id, cancel_event)
            
            # 步骤 3: 并发搜图 (30-50%)
            await task_manager.update_progress(task_id, 0.30, "正在搜索配图...")
            images = await self._fetch_images_concurrent(
                task_id=task_id,
                visual_plan=visual_plan,
                cancel_event=cancel_event,
            )
            
            self._check_cancelled(task_id, cancel_event)
            
            # 步骤 4: 智能排版 (50-60%)
            await task_manager.update_progress(task_id, 0.50, "正在排版布局...")
            layout_plans = await self._create_layout_plans(
                task_id=task_id,
                outline=outline,
                visual_plan=visual_plan,
                images=images,
                template=template,
                cancel_event=cancel_event,
            )
            
            self._check_cancelled(task_id, cancel_event)
            
            # 步骤 5: 渲染 PPTX (60-90%)
            await task_manager.update_progress(task_id, 0.60, "正在生成 PPTX...")
            pptx_path = await self._render_pptx(
                task_id=task_id,
                outline=outline,
                layout_plans=layout_plans,
                template=template,
                cancel_event=cancel_event,
            )
            
            self._check_cancelled(task_id, cancel_event)
            
            # 步骤 6: 应用动画 (90-95%)
            if options.get("enable_animations", False):
                await task_manager.update_progress(task_id, 0.90, "正在应用动画效果...")
                pptx_path = await self._apply_animations(
                    task_id=task_id,
                    pptx_path=pptx_path,
                    cancel_event=cancel_event,
                )
            
            # 步骤 7: 保存文件 (95-100%)
            await task_manager.update_progress(task_id, 0.95, "正在保存文件...")
            file_id = await self._save_file(
                task_id=task_id,
                pptx_path=pptx_path,
                output_format=output_format,
            )
            
            # 完成
            await task_manager.update_progress(task_id, 1.0, "生成完成!")
            await task_manager.mark_success(task_id, {
                "file_id": file_id,
                "slide_count": len(outline.slides),
                "download_url": f"/api/v1/pptx/download/{file_id}",
                "preview_url": f"/api/v1/pptx/preview/{file_id}",
            })
            
            logger.info(f"PPT 生成完成 | task_id={task_id} | slides={len(outline.slides)}")
            
        except TaskCancelledError:
            await task_manager.mark_cancelled(task_id)
            logger.info(f"PPT 生成任务已取消 | task_id={task_id}")
            
        except Exception as e:
            logger.error(f"PPT 生成失败 | task_id={task_id} | error={e}", exc_info=True)
            await task_manager.mark_failed(task_id, str(e))
            
        finally:
            # 清理取消事件
            self._cancel_events.pop(task_id, None)
            self._running_tasks.pop(task_id, None)
    
    def _check_cancelled(self, task_id: str, cancel_event: asyncio.Event):
        """检查任务是否已取消"""
        if cancel_event.is_set():
            raise TaskCancelledError(f"任务 {task_id} 已取消")
    
    async def _generate_outline(
        self,
        task_id: str,
        topic: str,
        slide_count: int,
        api_key_token: Optional[str],
        cancel_event: asyncio.Event,
    ) -> PresentationOutline:
        """生成 PPT 大纲"""
        self._check_cancelled(task_id, cancel_event)
        
        outline = await self._ppt_agent.generate_outline(
            topic=topic,
            num_slides=slide_count,
            api_key_token=api_key_token,
        )
        
        await task_manager.update_progress(task_id, 0.15, f"大纲生成完成，共 {len(outline.slides)} 页")
        return outline
    
    async def _analyze_visual_needs(
        self,
        task_id: str,
        outline: PresentationOutline,
        template: str,
        cancel_event: asyncio.Event,
    ) -> Dict[str, Any]:
        """视觉决策分析"""
        self._check_cancelled(task_id, cancel_event)
        
        # 将大纲转换为视觉分析器可用的格式
        slides_data = [slide.to_dict() for slide in outline.slides]
        
        # 调用视觉分析器
        visual_plan = await visual_analyzer.analyze_ppt_content(
            topic=outline.title,
            slides=slides_data,
            template=template,
        )
        
        await task_manager.update_progress(task_id, 0.25, "视觉分析完成")
        return visual_plan
    
    async def _fetch_images_concurrent(
        self,
        task_id: str,
        visual_plan: Dict[str, Any],
        cancel_event: asyncio.Event,
        max_concurrent: int = 5,
    ) -> List[Optional[str]]:
        """并发搜索图片"""
        self._check_cancelled(task_id, cancel_event)
        
        # 提取需要搜索的关键词
        queries = []
        for slide_decision in visual_plan.get("slide_decisions", []):
            if slide_decision.get("needs_image"):
                queries.append(slide_decision.get("image_keywords", []))
        
        if not queries:
            return []
        
        # 使用 Semaphore 控制并发数
        semaphore = asyncio.Semaphore(max_concurrent)
        results = [None] * len(queries)
        
        async def fetch_with_semaphore(idx, keywords):
            async with semaphore:
                self._check_cancelled(task_id, cancel_event)
                try:
                    # 搜索图片
                    image_url = await self._image_search.search_image(
                        keywords=keywords,
                        width=1920,
                        height=1080,
                    )
                    results[idx] = image_url
                    
                    # 更新进度
                    progress = 0.30 + (0.20 * (idx + 1) / len(queries))
                    await task_manager.update_progress(
                        task_id,
                        progress,
                        f"正在搜索配图 ({idx + 1}/{len(queries)})...",
                    )
                except Exception as e:
                    logger.warning(f"图片搜索失败 | idx={idx} | error={e}")
                    results[idx] = None
        
        # 并发执行
        tasks = [
            fetch_with_semaphore(idx, keywords)
            for idx, keywords in enumerate(queries)
        ]
        await asyncio.gather(*tasks, return_exceptions=True)
        
        await task_manager.update_progress(task_id, 0.50, "配图搜索完成")
        return results
    
    async def _create_layout_plans(
        self,
        task_id: str,
        outline: PresentationOutline,
        visual_plan: Dict[str, Any],
        images: List[Optional[str]],
        template: str,
        cancel_event: asyncio.Event,
    ) -> List[Dict[str, Any]]:
        """创建排版计划"""
        self._check_cancelled(task_id, cancel_event)
        
        layout_plans = []
        slides_data = [slide.to_dict() for slide in outline.slides]
        
        for idx, slide_data in enumerate(slides_data):
            # 获取对应的图片 URL
            image_url = images.pop(0) if images else None
            
            # 调用布局决策器
            layout_plan = await layout_decider.decide_layout(
                slide_type=slide_data.get("type", "content"),
                title=slide_data.get("title", ""),
                content=slide_data.get("bullets", []),
                image_url=image_url,
                template=template,
            )
            
            layout_plans.append(layout_plan)
            
            # 更新进度
            progress = 0.50 + (0.10 * (idx + 1) / len(slides_data))
            await task_manager.update_progress(
                task_id,
                progress,
                f"正在排版 ({idx + 1}/{len(slides_data)})...",
            )
        
        return layout_plans
    
    async def _render_pptx(
        self,
        task_id: str,
        outline: PresentationOutline,
        layout_plans: List[Dict[str, Any]],
        template: str,
        cancel_event: asyncio.Event,
    ) -> Path:
        """渲染 PPTX 文件"""
        self._check_cancelled(task_id, cancel_event)
        
        # TODO: 这里需要调用实际的 PPTX 渲染逻辑
        # 当前使用占位实现，后续将集成完整的渲染流程
        
        output_path = PPT_OUTPUT_DIR / f"{task_id}.pptx"
        
        # 使用 python-pptx 创建基础 PPTX
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        
        # 获取模板样式
        from app.api.v1.aiGeneratorPptx import PPT_TEMPLATES, PPTStyle
        style = PPTStyle(template_name=template)
        
        total_slides = len(outline.slides)
        
        for idx, (slide_data, layout_plan) in enumerate(zip(outline.slides, layout_plans)):
            self._check_cancelled(task_id, cancel_event)
            
            # 创建幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                Inches(9), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = style.PRIMARY_COLOR
            p.alignment = PP_ALIGN.LEFT
            
            # 添加内容
            if slide_data.bullets:
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2),
                    Inches(9), Inches(5)
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                
                for i, bullet in enumerate(slide_data.bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = style.TEXT_DARK
                    p.space_after = Pt(12)
            
            # 更新进度
            progress = 0.60 + (0.30 * (idx + 1) / total_slides)
            await task_manager.update_progress(
                task_id,
                progress,
                f"正在渲染第 {idx + 1}/{total_slides} 页...",
            )
        
        # 保存
        prs.save(str(output_path))
        
        return output_path
    
    async def _apply_animations(
        self,
        task_id: str,
        pptx_path: Path,
        cancel_event: asyncio.Event,
    ) -> Path:
        """应用动画效果"""
        self._check_cancelled(task_id, cancel_event)
        
        # TODO: 集成 AnimationEngine
        # 当前版本跳过动画应用
        logger.info(f"动画效果应用（暂未实现）| task_id={task_id}")
        
        return pptx_path
    
    async def _save_file(
        self,
        task_id: str,
        pptx_path: Path,
        output_format: str,
    ) -> str:
        """保存文件并返回文件 ID"""
        # 生成文件 ID
        file_id = str(uuid.uuid4())
        
        # TODO: 集成 FileStorageManager
        # 当前版本直接返回本地路径
        
        if output_format == "pdf":
            # TODO: 实现 PDF 转换
            logger.warning(f"PDF 输出暂未实现，返回 PPTX | task_id={task_id}")
        
        return str(pptx_path)


# 全局单例
ppt_orchestrator = PPTOrchestrator()
